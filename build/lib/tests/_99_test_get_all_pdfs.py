import os
from PyPDF2 import PdfMerger
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
import re
#===============================================================================
outputs_dir = "/pscratch/sd/a/adammwea/workspace/RBS_network_models/optimizing/CDKL5/CDKL5_DIV21/outputs"
pdf_dir = "/pscratch/sd/a/adammwea/workspace/RBS_network_models/optimizing/CDKL5/CDKL5_DIV21/slides"
#===============================================================================
keyword_exclude = []
keyword_include = [
    'cand_',
    'gen_',
    '_slide',
]
def combine_pdfs_to_pptx(
    outputs_dir, 
    pdf_dir, 
    combined_pdf_name="all_pdfs.pdf", 
    pptx_name="all_pdfs.pptx",
    keyword_exclude=None,
    keyword_include=None,
    ):
    # Ensure the PDF directory exists
    os.makedirs(pdf_dir, exist_ok=True)

    # Find all PDF files in the outputs directory
    pdf_files = []
    for root, _, files in os.walk(outputs_dir):
        for file in files:
            if file.endswith(".pdf"):
                if keyword_exclude is not None and any([keyword in file for keyword in keyword_exclude]):
                    continue
                if keyword_include is not None and not all([keyword in file for keyword in keyword_include]):
                    continue
                pdf_files.append(os.path.join(root, file))

    if not pdf_files:
        print("No PDF files found in the outputs directory.")
        return

    # Combine all PDFs into one
    combined_pdf_path = os.path.join(pdf_dir, combined_pdf_name)
    #merger = PdfMerger()
    #for pdf in sorted(pdf_files):
    mergers = {}
    for pdf in pdf_files:
        
        try:
            #get run name, gen name, and cand name
            run_name = os.path.basename(os.path.dirname(os.path.dirname(pdf)))
            gen_name = os.path.basename(os.path.dirname(pdf))
            cand_name = re.search(r'cand_?(\d{1,4})', os.path.basename(pdf))[0]
            
            # parse paths by run and gen to make more manageable files
            if run_name not in mergers:
                mergers[run_name] = {}
            if gen_name not in mergers[run_name]:
                mergers[run_name][gen_name] = PdfMerger()
            try: mergers[run_name][gen_name].append(pdf)
            except Exception as e:
                print(pdf)
                print(f"Error: {e}")
                continue
        except Exception as e:
            print(pdf)
            print(f"Error: {e}")
            continue
            
    for run_name, run in mergers.items():
        for gen_name, merger in run.items():
            combined_pdf_path = os.path.join(pdf_dir, f"{run_name}_{gen_name}.pdf")
            merger.write(combined_pdf_path)
            merger.close()
            print(f"Combined PDF saved to {combined_pdf_path}")
            
            # Convert the combined PDF to a PowerPoint presentation
            presentation = Presentation()
            try: pages = convert_from_path(combined_pdf_path) #NOTE: this is the slowest part of the process
            except Exception as e:
                print(combined_pdf_path)
                print(f"Error: {e}")
                continue

            for page in pages:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
                image_path = os.path.join(pdf_dir, "temp_page.png")
                page.save(image_path, "PNG")

                # Add the image to the slide
                left = top = Inches(0)
                slide_width = presentation.slide_width
                slide_height = presentation.slide_height
                slide.shapes.add_picture(image_path, left, top, width=slide_width, height=slide_height)

                # Clean up the temporary image
                os.remove(image_path)

            pptx_name = f"{run_name}_{gen_name}.pptx"
            pptx_path = os.path.join(pdf_dir, pptx_name)
            presentation.save(pptx_path)
            print(f"PowerPoint presentation saved to {pptx_path}")
combine_pdfs_to_pptx(
    outputs_dir, 
    pdf_dir,
    keyword_exclude=keyword_exclude,
    keyword_include=keyword_include,
    )
