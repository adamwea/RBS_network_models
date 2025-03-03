import os
import spikeinterface.sorters as ss
#import workspace.RBS_neuronal_network_models.optimizing.CDKL5_DIV21.scripts_dep._1_sims.mea_processing_library as mea_lib
#import MEAProcessingLibrary as mea_lib
from MEA_Analysis.MEAProcessingLibrary import mea_processing_library as mea_lib
import numpy as np
import sys
#import subprocess
import importlib.util
import matplotlib.pyplot as plt
import pickle
from datetime import datetime
from concurrent.futures import ProcessPoolExecutor, as_completed
import json
from pdf2image import convert_from_path
import re
from PIL import Image
#import shutil
from netpyne import sim, specs
from time import time
#===================================================================================================
'''global variables'''
global CONVOLUTION_PARAMS, PROGRESS_SLIDES_PATH, REFERENCE_DATA_NPY, SIMULATION_RUN_PATH, DEBUG_MODE
#===================================================================================================

'''functions'''
# from now on (23Dec2024) I'm just going to put newer functions at the top of the file.
def reprocess_simulation(
    sim_data_path, 
    REFERENCE_DATA_NPY, 
    #CONVOLUTION_PARAMS, 
    #cfg_script_path=None, 
    #param_script_path=None, 
    #target_script_path=None, 
    #sim_cfg_path=None, 
    #netParams_data_path=None, 
    #duration_seconds=1, 
    #permuted_cfg_path=None, 
    #save_data=False, 
    #overwrite_cfgs=False
    DEBUG_MODE=False
    ):
    
    #
    #sim_data_path = SIMULATION_RUN_PATH
    
    # surpress all prints
    #suppress_print() # surpress all prints

    # old commented code...
        # #assert simulation run path is .json file with '_data' in the name
        # assert '_data' in sim_data_path and '.json' in sim_data_path
        # privileged_print(f"\nSimulation run path: {sim_data_path}")

        # #rerun simulation of interest
        # privileged_print("Rerunning simulation of interest...")
        # start = time()
        # def rerun_simulation_of_interest(sim_data_path, duration_seconds, permuted_cfg_path=None, save_data=True, overwrite_cfgs=False):
        #     # load simulation files 
        #     # def load_simulation(sim_data_path, permuted_cfg_path=None):
        #     #     #sim.clearAll()
        #     #     print(f'simLabel: {sim.cfg.simLabel}')
        #     #     output = sim.load(sim_data_path, output=True)  #output is none when there are cells in the sim_data I guess?
        #     #     if permuted_cfg_path is not None: 
        #     #         print(f'simLabel: {sim.cfg.simLabel}')
        #     #         sim.loadSimCfg(permuted_cfg_path)
        #     #         # #sim.net.createCells()
        #     #         # #sim.clearConns()
        #     #         # sim.clearObj("conns")
        #     #         # sim.net.connectCells()
                    
        #     #         # #TODO need to write a function that recreates cells in place. use .util logic/functions.
        #     #         # # hold shape and size constant.
                    
        #     #         # #TODO: make a slide about this for Roy.
                    
                    
        #     #         # sim.net.setupRecording()
        #     #         print(f'simLabel: {sim.cfg.simLabel}') 
        #     #     filename = sim.cfg.filename
        #     # #sim_data_path = SIMULATION_RUN_PATH
        #     # load_simulation(sim_data_path, permuted_cfg_path=permuted_cfg_path)
            
        #     # load simulation files
        #     def load_simulation(sim_data_path):
        #         try:
        #             sim.clearAll() #this will fail if sim has no net yet
        #         except: pass
        #         sim.load(sim_data_path)
        #         if permuted_cfg_path is not None:
        #             sim.loadSimCfg(permuted_cfg_path)
        #     load_simulation(sim_data_path)
            
        #     # modify simulation as needed
        #     instructions = {
        #         'duration_seconds': duration_seconds,      
        #         }        
        #     def modify_simulation(instructions):        
        #         # rebuild simulation as needed            
        #         # print simulation label            
        #         print(f'simLabel: {sim.cfg.simLabel}')
                
        #         ## things that shouldn't affect the simulation
                
        #         # #activate core neuron, which is optimized for large networks and hpcs
        #         # sim.cfg.coreneuron = True            
        #         # t_vec = sim.h.Vector()             # coreneuron needs to be told to record time explicitly
        #         # t_vec.record(sim.h._ref_t)
        #         sim.cfg.dump_coreneuron_model = False # not really sure what this does... it dumps the model instead of running it?
                
        #         # # deactivate cvode and use fast imem for now
        #         sim.cfg.cvode_active = False    # pretty sure this must be false for coreneuron to work - which is fine because coreneuron i dont think I need this?
        #         #                                 # copied from documentation:
        #         #                                 # Multi order variable time step integration method which may be used in place of the default staggered fixed time step method. 
        #         #                                 # The performance benefits can be substantial (factor of more than 10) for problems in which all states vary slowly for long periods of time between fast spikes.
        #         # sim.cfg.use_fast_imem = False   # this doesnt work with coreneuron, but it might help speed up the simulation. 
        #         #                                 # 
                
        #         # ## cache efficient is a good idea in general
        #         # sim.cfg.cache_efficient = True
                
        #         # ## outputs and stuff
        #         # sim.cfg.validateNetParams = True #this only plays a role when re-building the net, which we may or may not do.
        #         sim.cfg.printSynsAfterRule = True
        #         sim.cfg.saveTiming = True # this creates a time vector in the sim data...which I wish the documentation would have made more obvious
        #         # sim.cfg.printRunTime = True
        #         # sim.cfg.savePickle = False # saving pkls just seems like a waste of time lately
        #         # sim.cfg.saveCSV = False            
        #         # sim.cfg.timestampFilename = False # this might be useful for version control, but lets keep this false for now
                
        #         # ## things that might affect the simulation
        #         sim.cfg.allowSelfConns = False
        #         sim.cfg.oneSynPerNetcon = True    
        #         sim.cfg.allowConnsWithWeight0 = True
                
        #         ## turns out that modifying the cfg object alone doesn't actually change the simulation.
        #         ## need to modify intantiated objects in the sim object using modifyCells, modifyConns, modifySynMechs, etc.
        #         ## they expect params similar to those in the netParams object.
                
        #         #sim.net.connectCells()
                
        #         netParams = sim.net.params
        #         # params = {
        #         #     'conds': sim.cfg.__dict__.copy(),
        #         # }
        #         #sim.net.
        #         #sim.net.modifyCells(params)
                
        #         # clear and rebuild the net
        #         # #params = sim.cfg.__dict__.copy()
        #         # #params = sim.net.params.__dict__.copy()
        #         # cellParams = sim.net.params.cellParams
        #         # for pop, params in cellParams.items():
        #         #     #params = cell_pop
        #         #     sim.net.modifyCells(params)
                    
        #         # cellConds ={
        #         #     # if desired, I can filter synmechs modified by any tag i guess. 
        #         #     # e.g. {'pop': 'PYR', 'ynorm': [0.1, 0.6]} targets connections of cells from the ‘PYR’ population with normalized depth within 0.1 and 0.6.
        #         # }
                
        #         # synMechParams = sim.net.params.synMechParams
        #         # for pop, params in synMechParams.items():
        #         #     #params = cell_pop
        #         #     params['cellConds'] = {'pop': 'E'}
        #         #     params['']                
        #         #     sim.net.modifySynMechs(params)
                    
        #         #     params['cellConds'] = {'pop': 'I'}
        #         #     sim.net.modifySynMechs(params)
        #         #     #sim.net.modifySynMechs()
                    
        #         # for pop, params in cellParams.items():
        #         #     #params = cell_pop
        #         #     sim.net.modifyConns(params)
                    
        #         # # sim.net.createPops()
        #         # # sim.net.createCells()
        #         # # sim.net.connectCells()

        #         #modify duration of the simulation. Intuitively, modulating this last makes sense.
        #         sim.cfg.duration = duration_seconds*1e3
                
        #         return t_vec if 't_vec' in locals() else None

        #     t_vec = modify_simulation(instructions)
            
            
        #     #run simulation
        #     sim.simulate()
            
        #     # post process as needed
        #     def post_simulation_processing(tvec):
        #         sim.gatherData()
        #         # if sim.cfg.coreneuron: #save time vector to sim data - which for some reason is lost by coreneuron
        #         #     time_vector = np.array(t_vec)
        #         #     sim.allSimData.t = time_vector    
        #     post_simulation_processing(t_vec)

        #     #save data
        #     saved_files = None
        #     if save_data:
        #         activate_print()
        #         saved_files = sim.saveData()
        #         print("Data saved.")
        #         suppress_print()
        #     return saved_files
        # saved_files = rerun_simulation_of_interest(sim_data_path, duration_seconds, permuted_cfg_path=permuted_cfg_path, save_data=save_data, overwrite_cfgs=overwrite_cfgs)
        # from pprint import pprint
        # activate_print()
        # pprint(saved_files)
        
        # # update sim_data_path if necessary
        # def update_sim_data_path(saved_files):
        #     # if any of the saved files are _data.json, update sim_data_path
        #     if any('_data.json' in f for f in saved_files):
        #         new_path = [f for f in saved_files if '_data.json' in f][0]
        #     return new_path
        # if saved_files is not None: sim_data_path = update_sim_data_path(saved_files)   
        # print("Simulation rerun complete.")
        # print(f"Time taken: {time()-start} seconds")
        # suppress_print()

    # re-fit simulation of interest
    privileged_print("Refitting simulation of interest...")
    start = time()
    #from ._1_sims.calculate_fitness_vCurrent import fitnessFunc
    #from DIV21.utils.fitness_func import fitnessFunc
    from RBS_network_models.fitnessFunc import fitnessFunc
    def refit_simulation_of_interest(
        sim_data_path, 
        #target_script_path, 
        #saved_files=None
        ):   
        
        #sim.clearAll()
        #sim.load(sim_data_path)
        try:
            simData = sim.allSimData
        except:
            print("sim.allSimData failed. Trying sim.load(sim_data_path)...")
            sim.load(sim_data_path)
            simData = sim.allSimData
            # print e and i rates for debugging
            #print(simData['popRates'])
            privileged_print(f'path: {sim_data_path}')
            privileged_print(f"Excitatory firing rate: {simData['popRates']['E']}")
            privileged_print(f"Inhibitory firing rate: {simData['popRates']['I']}")
        #convolution_params = import_module_from_path(CONVOLUTION_PARAMS)
        #script_dir = os.path.dirname(os.path.realpath(__file__))
        #target_script_path = os.path.join(script_dir, target_script_path)
        #fitnessFuncArgs = import_module_from_path(target_script_path).fitnessFuncArgs
        #from DIV21.src.conv_params import conv_params
        from RBS_network_models.CDKL5.DIV21.src.conv_params import conv_params
        #from RBS_network_models.developing.CDKL5.DIV21.src.conv_params import conv_params
        
        # def get_fitness_file_path(sim_data_path, saved_files):
        #     # if saved_files is not None:
        #     #     #find the file with data and json in the name
        #     #     data_file = [f for f in saved_files if '_data' in f and 'json' in f][0]
        #     #     #replace _data with _fitness
        #     #     fitness_file = data_file.replace('_data', '_fitness')
        #     # else:
        #     fitness_file = sim_data_path.replace('_data', '_fitness')
        #     return fitness_file
        # fitness_save_path = get_fitness_file_path(sim_data_path, saved_files)        
        
        #from DIV21.src.fitness_targets import fitnessFuncArgs
        #from RBS_network_models.developing.CDKL5.DIV21.src.fitness_targets import fitnessFuncArgs
        from RBS_network_models.CDKL5.DIV21.src.fitness_targets import fitnessFuncArgs
        fitness_save_path = sim_data_path.replace('_data', '_fitness')
        kwargs = {
            'simConfig': sim.cfg,
            #'conv_params': convolution_params.conv_params,
            'conv_params': conv_params,
            'popData': sim.net.allPops,
            'cellData': sim.net.allCells, #not actually used in the fitness calculation, but whatever
            'targets': fitnessFuncArgs['targets'],
            'maxFitness': fitnessFuncArgs['maxFitness'],
            'features': fitnessFuncArgs['features'],
            'fitness_save_path': fitness_save_path,
            'break_deals': False #no need to do deal_breakers when doing refitting or sensitivity analysis
        }
        average_fitness = fitnessFunc(simData, mode='simulated data', **kwargs)
        return average_fitness
    average_fitness = refit_simulation_of_interest(
        sim_data_path, 
        #target_script_path, 
        #saved_files=saved_files
        ) 
    
    # TODO: There seems to be some inconsistency in FR calculated by network metrics vs NETPYNE - need to investigate this.
    privileged_print(f"Average fitness: {average_fitness}")
    privileged_print("Refit complete.")
    privileged_print(f"Time taken: {time()-start} seconds")

    # re-plot simulation of interest - re-generate summary plot and all associated plots
    privileged_print("Replotting simulation of interest...")
    start = time()
    #from CDKL5_DIV21.src.conv_params import conv_params
    def replot_simulation_of_interest(sim_data_path, 
                                      #convolution_params_path, 
                                      #reference_data_npy, 
                                      average_fitness,
                                      trim_start=0,
                                      ):
        start_network_metrics = time()
        #from CDKL5_DIV21.src.conv_params import conv_params
        def get_network_metrics():
            #convolution_params = import_module_from_path(CONVOLUTION_PARAMS)
            #from DIV21.src.conv_params import conv_params
            from RBS_network_models.CDKL5.DIV21.src.conv_params import conv_params
            #from RBS_network_models.developing.CDKL5.DIV21.src.conv_params import conv_params
            #get network metrics
            kwargs = {
                'simData': sim.allSimData,
                'simConfig': sim.cfg,
                #'conv_params': convolution_params.conv_params,
                'conv_params': conv_params,
                'popData': sim.net.allPops,
                'cellData': sim.net.allCells, #not actually used in the fitness calculation, but whatever
            }
            #from DIV21.utils.fitness_helper import calculate_network_metrics
            #from RBS_network_models.developing.utils.fitness_helper import calculate_network_metrics
            from RBS_network_models.utils.fitness_helper import calculate_network_metrics
            error, kwargs = calculate_network_metrics(kwargs)
            #save network metrics
            network_metrics_path = sim_data_path.replace('_data', '_network_metrics')
            if '.pkl' in network_metrics_path:
                network_metrics_path = network_metrics_path.replace('.pkl', '.npy')
            elif '.json' in network_metrics_path:
                network_metrics_path = network_metrics_path.replace('.json', '.npy')
            np.save(network_metrics_path, kwargs)
            # with open(network_metrics_path, 'w') as f:
            #     json.dump(kwargs, f, indent=4)
            return error, kwargs
        error, kwargs = get_network_metrics()
        if error: return error
        privileged_print("\tNetwork metrics calculated - kwargs dict created.")
        privileged_print(f'\tTime taken: {time()-start_network_metrics} seconds')

        #plot raster
        start_raster = time()
        def plot_simulated_raster_wrapper(ax = None, subplot=False, trim_start = None, **kwargs):
            network_metrics = kwargs['network_metrics']
            spiking_data_by_unit = network_metrics['spiking_data']['spiking_data_by_unit']
            popData = sim.net.allPops
            E_gids = popData['E']['cellGids']
            I_gids = popData['I']['cellGids']
            
            if ax is None:
                fig, ax_raster = plt.subplots(1, 1, figsize=(16, 4.5))
            else:
                ax_raster = ax
                subplot = True
                
            ax_raster = plot_raster(ax_raster, spiking_data_by_unit, E_gids=E_gids, I_gids=I_gids, data_type='simulated')
            
            # if trim_start, trim first x seconds from the start of the simulation
            if trim_start is not None and trim_start > 0 and trim_start < ax_raster.get_xlim()[1]:
                ax_raster.set_xlim(trim_start, ax_raster.get_xlim()[1])
            elif trim_start is not None and trim_start > 0 and trim_start > ax_raster.get_xlim()[1]:
                modified_trim = ax_raster.get_xlim()[1]*0.1
                ax_raster.set_xlim(modified_trim, ax_raster.get_xlim()[1])
                print('boop')                
            
            #break if subplot
            if subplot:
                #plt.close()
                return ax_raster

            if DEBUG_MODE:
                #save local for debugging
                #dev_dir = '/pscratch/sd/a/adammwea/workspace/RBS_neuronal_network_models/optimization_projects/CDKL5_DIV21/scripts'
                dev_dir = os.path.dirname(os.path.realpath(__file__))
                plt.savefig(os.path.join(dev_dir, '_raster_plot.png'), dpi=300)
                #save local for debugging
                
            #save wherever data is saved
            #sim_data_path = SIMULATION_RUN_PATH
            raster_plot_path = sim_data_path.replace('_data', '_raster_plot')
            #remove file type and replace with png
            if '.json' in raster_plot_path:
                raster_plot_path = raster_plot_path.replace('.json', '.png')
            elif '.pkl' in raster_plot_path:
                raster_plot_path = raster_plot_path.replace('.pkl', '.png')
            #raster_plot_path = raster_plot_path.replace('.json', '.png')
            plt.savefig(raster_plot_path, dpi=300)
            print(f"Raster plot saved to {raster_plot_path}")
            
            #save as pdf
            raster_plot_path = sim_data_path.replace('_data', '_raster_plot')
            #raster_plot_path = raster_plot_path.replace('.json', '.pdf')
            if '.json' in raster_plot_path:
                raster_plot_path = raster_plot_path.replace('.json', '.pdf')
            elif '.pkl' in raster_plot_path:
                raster_plot_path = raster_plot_path.replace('.pkl', '.pdf')
            plt.savefig(raster_plot_path)
            print(f"Raster plot saved to {raster_plot_path}")
            plt.close()
            
            raster_plots_paths = [
                raster_plot_path,
                raster_plot_path.replace('.pdf', '.png'),
            ]
            
            return raster_plots_paths
        raster_plot_paths = plot_simulated_raster_wrapper(trim_start = trim_start, **kwargs)
        privileged_print("\tIndividual raster plots saved.")
        privileged_print(f'\tTime taken: {time()-start_raster} seconds')

        #plot bursting summary
        start_bursting = time()
        def plot_simulated_bursting_wrapper(ax = None, subplot=False, trim_start = None, **kwargs):
            if ax is None:
                fig, new_ax = plt.subplots(1, 1)
                fig.set_size_inches(16, 4.5)
            else:
                new_ax = ax
                subplot = True #if ax is passed in, then we are plotting on a subplot
            
            #
            conv_params = kwargs['conv_params']
            SpikeTimes = kwargs['network_metrics']['spiking_data']['spiking_times_by_unit']
            #from DIV21.utils.fitness_helper import plot_network_activity_aw
            #from RBS_network_models.developing.utils.analysis_helper import plot_network_activity_aw
            from MEA_Analysis.NetworkAnalysis.awNetworkAnalysis.network_analysis import plot_network_activity_aw
            #bursting_ax, _ = plot_network_activity_aw(new_ax, SpikeTimes, **conv_params) #TODO need to make sure this function agrees with mandar. would be best if we shared a function here.
            bursting_ax = kwargs['network_metrics']['bursting_data']['bursting_summary_data']['ax']
            mega_ax = kwargs['network_metrics']['mega_bursting_data']['bursting_summary_data']['ax']
            
            # # HACK
            # mega_conv_params = kwargs['conv_params'].copy()
            # mega_conv_params['binSize'] *= 5
            # mega_conv_params['gaussianSigma'] *= 15
            #mega_ax, _ = plot_network_activity_aw(new_ax, SpikeTimes, **mega_conv_params) #TODO need to make sure this function agrees with mandar. would be best if we shared a function here.
            
            from MEA_Analysis.NetworkAnalysis.awNetworkAnalysis.network_analysis import plot_network_bursting_experimental
            new_ax = plot_network_bursting_experimental(new_ax, bursting_ax, mega_ax=mega_ax)            
            
            # if trim_start, trim first x seconds from the start of the simulation
            if trim_start is not None and trim_start > 0 and trim_start < new_ax.get_xlim()[1]:
                new_ax.set_xlim(trim_start, new_ax.get_xlim()[1])
            elif trim_start is not None and trim_start > 0 and trim_start > new_ax.get_xlim()[1]:
                modified_trim = new_ax.get_xlim()[1]*0.1
                new_ax.set_xlim(modified_trim, new_ax.get_xlim()[1])    
            
            new_ax.set_title('Bursting summary')
            new_ax.set_xlabel('Time (s)')
            new_ax.set_ylabel('Fire rate (Hz)')
            
            #break if subplot
            if subplot:
                #plt.close()
                return new_ax
            
            plt.tight_layout()
            
            if DEBUG_MODE:
                # Save local for debugging
                dev_dir = os.path.dirname(os.path.realpath(__file__))
                plt.savefig(os.path.join(dev_dir, '_bursting_plot.png'), dpi=300)
                # Save local for debugging
            
            # save wherever data is saved
            #sim_data_path = SIMULATION_RUN_PATH
            bursting_plot_path = sim_data_path.replace('_data', '_bursting_plot')
            #remove file type and replace with png
            #bursting_plot_path = bursting_plot_path.replace('.json', '.png')
            if '.json' in bursting_plot_path:
                bursting_plot_path = bursting_plot_path.replace('.json', '.png')
            elif '.pkl' in bursting_plot_path:
                bursting_plot_path = bursting_plot_path.replace('.pkl', '.png')
            plt.savefig(bursting_plot_path, dpi=300)
            print(f"Bursting plot saved to {bursting_plot_path}")
            
            #save as pdf
            bursting_plot_path = sim_data_path.replace('_data', '_bursting_plot')
            #bursting_plot_path = bursting_plot_path.replace('.json', '.pdf')
            if '.json' in bursting_plot_path:
                bursting_plot_path = bursting_plot_path.replace('.json', '.pdf')
            elif '.pkl' in bursting_plot_path:
                bursting_plot_path = bursting_plot_path.replace('.pkl', '.pdf')
            plt.savefig(bursting_plot_path)
            print(f"Bursting plot saved to {bursting_plot_path}")
            plt.close()
            
            bursting_plot_paths = [
                bursting_plot_path,
                bursting_plot_path.replace('.pdf', '.png'),
            ]
            
            return bursting_plot_paths    
        bursting_plot_paths = plot_simulated_bursting_wrapper(trim_start = trim_start, **kwargs)
        privileged_print("\tIndividual bursting plots saved.")
        privileged_print(f'\tTime taken: {time()-start_bursting} seconds')

        # combine plots into a single summary plot
        start_summary = time()
        def plot_simulation_summary(trim_start = None, **kwargs):
            fig, ax = plt.subplots(2, 1, figsize=(16, 9))
            raster_plot_ax, bursting_plot_ax = ax
            subplot = True
            raster_plot_ax = plot_simulated_raster_wrapper(ax=raster_plot_ax, subplot=subplot, trim_start = trim_start, **kwargs)
            bursting_plot_ax = plot_simulated_bursting_wrapper(ax=bursting_plot_ax, subplot=subplot, trim_start = trim_start, **kwargs)
            
            #make both plots share the same x-axis
            raster_plot_ax.get_shared_x_axes().join(raster_plot_ax, bursting_plot_ax)    
            plt.tight_layout()
            
            if DEBUG_MODE:
                # Save local for debugging
                dev_dir = os.path.dirname(os.path.realpath(__file__))
                plt.savefig(os.path.join(dev_dir, '_summary_plot.png'), dpi=300)
                # Save local for debugging
                
            # save wherever data is saved
            #sim_data_path = SIMULATION_RUN_PATH
            summary_plot_path = sim_data_path.replace('_data', '_summary_plot')
            #remove file type and replace with png
            #summary_plot_path = summary_plot_path.replace('.json', '.png')
            if '.json' in summary_plot_path:
                summary_plot_path = summary_plot_path.replace('.json', '.png')
            elif '.pkl' in summary_plot_path:
                summary_plot_path = summary_plot_path.replace('.pkl', '.png')
            plt.savefig(summary_plot_path, dpi=300)
            print(f"Summary plot saved to {summary_plot_path}")
            
            #save as pdf
            summary_plot_path = sim_data_path.replace('_data', '_summary_plot')
            #summary_plot_path = summary_plot_path.replace('.json', '.pdf')
            if '.json' in summary_plot_path:
                summary_plot_path = summary_plot_path.replace('.json', '.pdf')
            elif '.pkl' in summary_plot_path:
                summary_plot_path = summary_plot_path.replace('.pkl', '.pdf')
            plt.savefig(summary_plot_path)  
            print(f"Summary plot saved to {summary_plot_path}")
            plt.close()
            
            summary_plot_paths = [
                summary_plot_path,
                summary_plot_path.replace('.pdf', '.png'),
            ]    
            return summary_plot_paths
        plot_simulation_summary(trim_start=trim_start, **kwargs)
        privileged_print("\tSimulation summary plot saved.")
        privileged_print(f'\tTime taken: {time()-start_summary} seconds')

        # comparison plot against reference data snippet
        start_comparison = time()
        activate_print()
        def plot_comparision_plot(ax_list=None, subplot=False, trim_start = None, **kwargs):
            #activate_print() #for debugging
            #fig, ax = plt.subplots(4, 1, figsize=(16, 9))
            if ax_list is None:
                fig, ax_list = plt.subplots(4, 1, figsize=(16, 9))
                sim_raster_ax, sim_bursting_ax, ref_raster_ax, ref_bursting_ax = ax_list
            else:
                #assert that there be 4 axes in the ax list
                assert len(ax_list) == 4, "There must be 4 axes in the ax_list."
                sim_raster_ax, sim_bursting_ax, ref_raster_ax, ref_bursting_ax = ax_list
                subplot = True
            
            #plot simulated raster
            sim_raster_ax = plot_simulated_raster_wrapper(ax=sim_raster_ax, subplot=True, trim_start=trim_start, **kwargs)
            
            #plot simulated bursting
            sim_bursting_ax = plot_simulated_bursting_wrapper(ax=sim_bursting_ax, subplot=True, trim_start=trim_start, **kwargs)
            
            # plot reference raster
            def plot_reference_raster_wrapper(ax=None, subplot=False, sim_data_length=None, **kwargs):
                #load npy ref data
                print(ax)
                print(subplot)
                #print(kwargs)
                
                #load npy ref data
                ref_data = np.load(
                    REFERENCE_DATA_NPY, 
                    allow_pickle=True
                    ).item()
                network_metrics = ref_data
                spiking_data_by_unit = network_metrics['spiking_data']['spiking_data_by_unit'].copy()
                unit_ids = [unit_id for unit_id in spiking_data_by_unit.keys()]
                
                # Trim reference data to match the length of simulated data
                if sim_data_length is not None:
                    for unit_id in unit_ids:
                        #spiking_data_by_unit[unit_id] = spiking_data_by_unit[unit_id][:sim_data_length]
                        #max_spike_time_index = np.argmax(spiking_data_by_unit[unit_id])
                        spike_times = spiking_data_by_unit[unit_id]['spike_times']
                        spike_times = spike_times[spike_times < sim_data_length]
                        spiking_data_by_unit[unit_id]['spike_times'] = spike_times
                
                if ax is None:
                    fig, ax_raster = plt.subplots(1, 1, figsize=(16, 4.5))
                else:
                    ax_raster = ax
                    subplot = True
                    
                ax_raster = plot_raster(ax_raster, spiking_data_by_unit, unit_ids=unit_ids, data_type='experimental')
                
                if subplot:
                    #plt.close()
                    print("subplot")
                    return ax_raster
                
                if DEBUG_MODE:
                    #save local for debugging
                    #dev_dir = '/pscratch/sd/a/adammwea/workspace/RBS_neuronal_network_models/optimization_projects/CDKL5_DIV21/scripts'
                    dev_dir = os.path.dirname(os.path.realpath(__file__))
                    plt.savefig(os.path.join(dev_dir, '_ref_raster_plot.png'), dpi=300)
                    #save as pdf
                    plt.savefig(os.path.join(dev_dir, '_ref_raster_plot.pdf'))
                    #save local for debugging
                
                
                #TODO: semi unfinished compared to similar function for simulated raster
            #sim_data_length = len(kwargs['network_metrics']['spiking_data']['spiking_data_by_unit'][0])
            sim_data_length = kwargs['simConfig']['duration']/1000 #in seconds
            ref_raster_ax = plot_reference_raster_wrapper(ax=ref_raster_ax, subplot=True, sim_data_length = sim_data_length, **kwargs)
            
            #plot reference bursting
            def plot_reference_bursting_wrapper(ax=None, subplot=False, sim_data_length=None, trim_start = None, **kwargs):
                #ref_data = np.load(REFERENCE_DATA_NPY, allow_pickle=True).item()
                ref_data = np.load(REFERENCE_DATA_NPY, allow_pickle=True).item()
                network_metrics = ref_data
                conv_params = kwargs['conv_params']
                SpikeTimes = network_metrics['spiking_data']['spiking_times_by_unit']
                if ax is None:
                    fig, new_ax = plt.subplots(1, 1)
                    fig.set_size_inches(16, 4.5)
                else:
                    new_ax = ax
                    subplot = True
                    
                if sim_data_length is not None:
                    for unit_id in SpikeTimes.keys():
                        spike_times = SpikeTimes[unit_id]
                        spike_times = spike_times[spike_times < sim_data_length]
                        SpikeTimes[unit_id] = spike_times
                
                #from DIV21.utils.fitness_helper import plot_network_activity_aw
                #from RBS_network_models.developing.utils.analysis_helper import plot_network_activity_aw
                from MEA_Analysis.NetworkAnalysis.awNetworkAnalysis.network_analysis import plot_network_activity_aw
                new_ax, _ = plot_network_activity_aw(new_ax, SpikeTimes, **conv_params) #TODO need to make sure this function agrees with mandar. would be best if we shared a function here.
                
                new_ax.set_title('Bursting summary')
                new_ax.set_xlabel('Time (s)')
                new_ax.set_ylabel('Fire rate (Hz)')
                
                #break if subplot
                if subplot:
                    #plt.close()
                    return new_ax
                
                plt.tight_layout()
                
                if DEBUG_MODE:
                    # Save local for debugging
                    dev_dir = os.path.dirname(os.path.realpath(__file__))
                    plt.savefig(os.path.join(dev_dir, '_ref_bursting_plot.png'), dpi=300)
                    # plot as pdf
                    plt.savefig(os.path.join(dev_dir, '_ref_bursting_plot.pdf'))
                    # Save local for debugging
                
                # TODO: semi unfinished compared to similar function for simulated bursting
            sim_data_length = kwargs['simConfig']['duration']/1000 #in seconds
            ref_bursting_ax = plot_reference_bursting_wrapper(ax=ref_bursting_ax, subplot=True, sim_data_length=sim_data_length, trim_start = trim_start, **kwargs)
            
            # # remove the first second of x-axis for all plots
            # def set_xlim_to_first_value_over_one(ax):
            #     x_data = ax.lines[0].get_xdata()
            #     first_value_over_one = next((x for x in x_data if x > 1), None)
            #     if first_value_over_one is not None:
            #         ax.set_xlim(first_value_over_one, ax.get_xlim()[1])
            
            # set_xlim_to_first_value_over_one(sim_raster_ax)
            # set_xlim_to_first_value_over_one(ref_raster_ax)
            # set_xlim_to_first_value_over_one(sim_bursting_ax)
            # set_xlim_to_first_value_over_one(ref_bursting_ax)
            
            # #print axes limits
            # print(sim_raster_ax.get_xlim())
            # print(ref_raster_ax.get_xlim())
            # print(sim_bursting_ax.get_xlim())
            # print(ref_bursting_ax.get_xlim())
            #ref_bursting_ax.append(ref_bursting_ax.get_xlim()[0] + 1)
            
            # # Ensure y-axis is the same for bursting plots
            # sim_bursting_ylim = sim_bursting_ax.get_ylim()
            # ref_bursting_ylim = ref_bursting_ax.get_ylim()
            
            # ensure xaxis of refernce plots matches simulated raster
            ref_raster_ax.set_xlim(sim_raster_ax.get_xlim())
            ref_bursting_ax.set_xlim(sim_raster_ax.get_xlim())
            sim_bursting_ax.set_xlim(ref_bursting_ax.get_xlim())
            
            # Ensure y-axis is the same for bursting plots
            def adjust_ylim_based_on_xlim(ax):
                x_data = ax.lines[0].get_xdata()
                y_data = ax.lines[0].get_ydata()
                xlim = ax.get_xlim()
                filtered_y_data = [y for x, y in zip(x_data, y_data) if xlim[0] <= x <= xlim[1]]
                if filtered_y_data:
                    ax.set_ylim(min(filtered_y_data) * 0.8, max(filtered_y_data) * 1.2)

            adjust_ylim_based_on_xlim(sim_bursting_ax)
            adjust_ylim_based_on_xlim(ref_bursting_ax)  
            
            
            # Calculate the combined y-axis limits
            sim_bursting_ylim = sim_bursting_ax.get_ylim()
            ref_bursting_ylim = ref_bursting_ax.get_ylim()
            # Calculate the combined y-axis limits
            combined_ylim = [
                min(sim_bursting_ylim[0], ref_bursting_ylim[0]) * 0.8,
                max(sim_bursting_ylim[1], ref_bursting_ylim[1]) * 1.2
            ]
            
            # Set the same y-axis limits for both axes
            sim_bursting_ax.set_ylim(combined_ylim)
            ref_bursting_ax.set_ylim(combined_ylim)
            
            # Make all four plots share the same x-axis as simulated raster
            sim_raster_ax.get_shared_x_axes().join(sim_raster_ax, sim_bursting_ax)
            sim_raster_ax.get_shared_x_axes().join(sim_raster_ax, ref_raster_ax)
            sim_raster_ax.get_shared_x_axes().join(sim_raster_ax, ref_bursting_ax)  

            #
            if subplot:
                #plt.tight_layout()
                #plt.close()
                return [sim_raster_ax, sim_bursting_ax, ref_raster_ax, ref_bursting_ax] 
            
            plt.tight_layout()
            
            if DEBUG_MODE:
                # Save local for debugging
                dev_dir = os.path.dirname(os.path.realpath(__file__))
                fig.savefig(os.path.join(dev_dir, '_comparison_plot.png'), dpi=300)
                # Save local for debugging
            
            # save wherever data is saved
            #sim_data_path = SIMULATION_RUN_PATH
            comparison_plot_path = sim_data_path.replace('_data', '_comparison_plot')
            #remove file type and replace with png
            #comparison_plot_path = comparison_plot_path.replace('.json', '.png')
            if '.json' in comparison_plot_path:
                comparison_plot_path = comparison_plot_path.replace('.json', '.png')
            elif '.pkl' in comparison_plot_path:
                comparison_plot_path = comparison_plot_path.replace('.pkl', '.png')
            fig.savefig(comparison_plot_path, dpi=300)
            print(f"Comparison plot saved to {comparison_plot_path}")
            
            #save as pdf
            comparison_plot_path = sim_data_path.replace('_data', '_comparison_plot')
            #comparison_plot_path = comparison_plot_path.replace('.json', '.pdf')
            if '.json' in comparison_plot_path:
                comparison_plot_path = comparison_plot_path.replace('.json', '.pdf')
            elif '.pkl' in comparison_plot_path:
                comparison_plot_path = comparison_plot_path.replace('.pkl', '.pdf')
            fig.savefig(comparison_plot_path)
            print(f"Comparison plot saved to {comparison_plot_path}")
            plt.close()
        plot_comparision_plot(trim_start=trim_start, **kwargs)
        privileged_print("\tComparison plot saved.")
        privileged_print(f'\tTime taken: {time()-start_comparison} seconds')

        # build comparison summary slide
        start_summary_slide = time()
        def build_comparision_summary_slide(sim_data_path, trim_start = None,):
            fig, ax_list = plt.subplots(4, 1, figsize=(16, 9))
            
            #plot comparison plot
            plot_comparision_plot(ax_list=ax_list, subplot=True, trim_start=trim_start, **kwargs)
            
            #remove x-axis labels from all plots
            for ax in ax_list:
                ax.set_xlabel('')
                
            #remove titles from all plots
            for ax in ax_list:
                ax.set_title('')
                
            #remove x-axis ticks from all plots except the bottom one
            for ax in ax_list[:-1]:
                ax.set_xticks([])
                
            # add 'simulated' and 'reference' labels above each raster plot
            ax_list[0].text(0.5, 1.1, 'Simulated', ha='center', va='center', transform=ax_list[0].transAxes, fontsize=12)
            ax_list[2].text(0.5, 1.1, 'Reference', ha='center', va='center', transform=ax_list[2].transAxes, fontsize=12)
            
            # add time(s) label to the bottom plot
            ax_list[-1].set_xlabel('Time (s)')
                    
            #create space to the right of plots for text
            fig.subplots_adjust(right=0.6)
            
            #create some space at the bottom for one line of text
            fig.subplots_adjust(bottom=0.1)
            
            #add text
            spiking_summary = kwargs['network_metrics']['spiking_data']['spiking_summary_data']
            bursting_summary = kwargs['network_metrics']['bursting_data']['bursting_summary_data']
            simulated_spiking_data = kwargs['network_metrics']['simulated_data']
            #simulated_bursting_data = kwargs['network_metrics']['simulated_data']
            
            #TODO: if meanBurstrate is not in the spiking summary, then the following will fail, add nan in that case for now
            if 'mean_Burst_Rate' not in bursting_summary:
                bursting_summary['mean_Burst_Rate'] = np.nan                
            
            text = (
                # f"Summary of comparison between simulated and reference data:\n"
                # f"Simulated raster plot is shown in the top left, reference raster plot is shown in the bottom left.\n"
                # f"Simulated bursting plot is shown in the top right, reference bursting plot is shown in the bottom right.\n"
                f"Simulated spiking summary:\n"
                f"  - Mean firing rate: {spiking_summary['MeanFireRate']} Hz\n"
                f"  - Coefficient of variation: {spiking_summary['CoVFireRate']}\n"
                f"  - Mean ISI: {spiking_summary['MeanISI']} s\n"
                f"  - Coefficient of variation of ISI: {spiking_summary['CoV_ISI']}\n"
                f"  - Mean firing rate E: {simulated_spiking_data['MeanFireRate_E']} Hz\n"
                f"  - Mean CoV FR E: {simulated_spiking_data['CoVFireRate_E']}\n"
                f"  - Mean firing rate I: {simulated_spiking_data['MeanFireRate_I']} Hz\n"
                f"  - Mean CoV FR I: {simulated_spiking_data['CoVFireRate_I']}\n"
                f"  - Mean ISI E: {simulated_spiking_data['MeanISI_E']} s\n"
                f"  - CoV ISI E: {simulated_spiking_data['CoV_ISI_E']}\n"
                f"  - Mean ISI I: {simulated_spiking_data['MeanISI_I']} s\n"
                f"  - CoV ISI I: {simulated_spiking_data['CoV_ISI_I']}\n"
                
                f"Simulated bursting summary:\n"        
                #mean burst rate
                #number of units
                f"  - Number of units: {bursting_summary['NumUnits']}\n"
                #baseline
                f"  - Baseline: {bursting_summary['baseline']} Hz\n"
                #fanofactor
                f"  - Fanofactor: {bursting_summary['fano_factor']}\n"
                #num bursts
                f"  - Number of bursts: {bursting_summary['Number_Bursts']}\n"
                #mean burst rate
                f"  - Mean burst rate: {bursting_summary['mean_Burst_Rate']} bursts/second\n"     
                #mean burst peak
                f"  - Mean burst peak: {bursting_summary['mean_Burst_Peak']} Hz\n"
                #burst peak CoV
                f"  - CoV burst peak: {bursting_summary['cov_Burst_Peak']}\n"
                #mean IBI
                f"  - Mean IBI: {bursting_summary['mean_IBI']} s\n"
                #CoV IBI
                f"  - CoV IBI: {bursting_summary['cov_IBI']}\n"
                #mean within burst isi
                f"  - Mean within burst ISI: {bursting_summary['MeanWithinBurstISI']} s\n"
                #CoV within burst isi
                f"  - CoV within burst ISI: {bursting_summary['CoVWithinBurstISI']}\n"
                #mean outside burst isi
                f"  - Mean outside burst ISI: {bursting_summary['MeanOutsideBurstISI']} s\n"
                #CoV outside burst isi
                f"  - CoV outside burst ISI: {bursting_summary['CoVOutsideBurstISI']}\n"
                #mean whole network isi
                f"  - Mean network ISI: {bursting_summary['MeanNetworkISI']} s\n"
                #CoV whole network isi
                f"  - CoV network ISI: {bursting_summary['CoVNetworkISI']}\n"
                )
            
            #append reference metrics to text
            ref_data = np.load(REFERENCE_DATA_NPY, allow_pickle=True).item()
            ref_spiking_summary = ref_data['spiking_data']['spiking_summary_data']
            ref_bursting_summary = ref_data['bursting_data']['bursting_summary_data']    
            text += (
                f"\n"
                f"\n"
                f"Reference spiking summary:\n"
                f"  - Mean firing rate: {ref_spiking_summary['MeanFireRate']} Hz\n"
                f"  - Coefficient of variation: {ref_spiking_summary['CoVFireRate']}\n"
                f"  - Mean ISI: {ref_spiking_summary['MeanISI']} s\n"
                f"  - Coefficient of variation of ISI: {ref_spiking_summary['CoV_ISI']}\n"
                
                f"Reference bursting summary:\n"
                #number of units
                f"  - Number of units: {ref_bursting_summary['NumUnits']}\n"
                #baseline
                f"  - Baseline: {ref_bursting_summary['baseline']} Hz\n"
                #fanofactor
                f"  - Fanofactor: {ref_bursting_summary['fano_factor']}\n"
                #num bursts
                f"  - Number of bursts: {ref_bursting_summary['Number_Bursts']}\n"
                #mean burst rate
                #f"  - Mean burst rate: {ref_bursting_summary['mean_Burst_Rate']} Hz\n" #TODO: this is not in the data        
                #mean burst peak
                f"  - Mean burst peak: {ref_bursting_summary['mean_Burst_Peak']} Hz\n"
                #burst peak CoV
                f"  - CoV burst peak: {ref_bursting_summary['cov_Burst_Peak']}\n"
                #mean IBI
                f"  - Mean IBI: {ref_bursting_summary['mean_IBI']} s\n"
                #CoV IBI
                f"  - CoV IBI: {ref_bursting_summary['cov_IBI']}\n"
                #mean within burst isi
                f"  - Mean within burst ISI: {ref_bursting_summary['MeanWithinBurstISI']} s\n"
                #CoV within burst isi
                f"  - CoV within burst ISI: {ref_bursting_summary['CoVWithinBurstISI']}\n"
                #mean outside burst isi
                f"  - Mean outside burst ISI: {ref_bursting_summary['MeanOutsideBurstISI']} s\n"
                #CoV outside burst isi
                f"  - CoV outside burst ISI: {ref_bursting_summary['CoVOutsideBurstISI']}\n"
                #mean whole network isi
                f"  - Mean network ISI: {ref_bursting_summary['MeanNetworkISI']} s\n"
                #CoV whole network isi
                f"  - CoV network ISI: {ref_bursting_summary['CoVNetworkISI']}\n"
                )
            
            #add average fitness to text
            text += (
                f"\n"
                f"\n"
                f"\nAverage fitness: {average_fitness}"
                ) 
            
            #add text to the right of the plots
            fig.text(0.65, 0.5, text, ha='left', va='center', fontsize=9)
            
            #add data path at the bottom of the slide
            #fig.text(0.5, 0.05, f"Data path: {SIMULATION_RUN_PATH}", ha='center', va='center', fontsize=10)
            #go a little lower and add data path
            fig.text(0.5, 0.02, f"Data path: {sim_data_path}", ha='center', va='center', fontsize=9) 
            
            if DEBUG_MODE:
                #save local for debugging
                dev_dir = os.path.dirname(os.path.realpath(__file__))
                fig.savefig(os.path.join(dev_dir, '_comparison_summary_slide.png'), dpi=300)
                #save as pdf
            
            # save wherever data is saved
            sim_data_path = sim_data_path
            comparison_summary_slide_path = sim_data_path.replace('_data', '_comparison_summary_slide')
            #remove file type and replace with png
            #comparison_summary_slide_path = comparison_summary_slide_path.replace('.json', '.png')
            if '.json' in comparison_summary_slide_path:
                comparison_summary_slide_path = comparison_summary_slide_path.replace('.json', '.png')
            elif '.pkl' in comparison_summary_slide_path:
                comparison_summary_slide_path = comparison_summary_slide_path.replace('.pkl', '.png')
            fig.savefig(comparison_summary_slide_path, dpi=300)
            privileged_print(f"Comparison summary slide saved to {comparison_summary_slide_path}")
            
            #save as pdf
            comparison_summary_slide_path = sim_data_path.replace('_data', '_comparison_summary_slide')
            #comparison_summary_slide_path = comparison_summary_slide_path.replace('.json', '.pdf')
            if '.json' in comparison_summary_slide_path:
                comparison_summary_slide_path = comparison_summary_slide_path.replace('.json', '.pdf')
            elif '.pkl' in comparison_summary_slide_path:
                comparison_summary_slide_path = comparison_summary_slide_path.replace('.pkl', '.pdf')
            fig.savefig(comparison_summary_slide_path)
            privileged_print(f"Comparison summary slide saved to {comparison_summary_slide_path}")
            plt.close()
        build_comparision_summary_slide(sim_data_path, trim_start = trim_start)
        privileged_print("\tComparison summary slide saved.")
        privileged_print(f'\tTime taken: {time()-start_summary_slide} seconds')
    replot_simulation_of_interest(
        sim_data_path, 
        #CONVOLUTION_PARAMS,
        #conv_params, 
        #REFERENCE_DATA_NPY, 
        average_fitness,
        trim_start=5,
        )
    privileged_print("Replotting complete.")
    privileged_print(f"Time taken: {time()-start} seconds")

def run_permutation(
    sim_data_path, 
    cfg, 
    cfg_param, 
    cfg_val,
    reference_data_path = None,
    plot = False,
    debug = False,
    debug_permuted_sims=True, #if true, will only generate sims, wont run them 
    # reference_data_path, 
    # conv_params_path, 
    # target_script_path,
    # duration_seconds, 
    # save_data=True, 
    # overwrite_cfgs=False,
    *args
    ):
    
    if plot: assert reference_data_path is not None, "Reference data path must be provided for plotting."
    
    try: 
        #from netpyne import sim
        # perm_simConfig = (cfg, cfg_param, cfg_val)
        # cfg, cfg_param, cfg_val = perm_simConfig
        simLabel = cfg['simLabel']
        print(f'Running permutation {simLabel}...')
        
        if cfg_param is not None and cfg_val is not None: #if none, then it's the original simConfig
            
            def prepare_permuted_sim(sim_data_path, cfg, cfg_param, cfg_val):
                # load netparams and permute
                sim.load(sim_data_path, simConfig=cfg)
                simConfig = specs.SimConfig(simConfigDict=cfg)
                netParams = sim.loadNetParams(sim_data_path, setLoaded=False)
                
                #Typical case:
                strategy = 'by_value'
                #SPECIAL CASES: gnabar, gkbar, L, diam, Ra
                handle_by_name = ['gnabar', 'gkbar', 'L', 'diam', 'Ra']
                if any([name in cfg_param for name in handle_by_name]): 
                    if '_' in cfg_param:
                        elements = cfg_param.split('_')
                        for element in elements:
                            if any([name==element for name in handle_by_name]):
                                strategy = 'by_name'
                                break
                    elif any([name==cfg_param for name in handle_by_name]):
                        strategy = 'by_name'

                def map_cfg_to_netparams(simConfig, netParams, strategy='by_value'):
                    """
                    Map attributes in simConfig to their corresponding locations in netParams based on values.
                    
                    Parameters:
                        simConfig (dict): The configuration dictionary (cfg).
                        netParams (object): The network parameters object.
                    
                    Returns:
                        dict: A mapping from simConfig parameters to their paths in netParams.
                    """
                    def find_value_in_netparams(value, netParams, current_path=""):
                        """
                        Recursively search for the value in netParams and return a list of matching paths.
                        
                        Parameters:
                            value (any): The value to search for.
                            netParams (object): The network parameters object.
                            current_path (str): The current path in the recursive search.
                        
                        Returns:
                            list: A list of paths to the matching value.
                        """
                        stack = [(netParams, current_path)]  # Stack for backtracking, contains (current_object, current_path)
                        matching_paths = []  # To store all matching paths

                        while stack:
                            current_obj, current_path = stack.pop()
                            
                            # if 'connParams' in current_path:  # Debugging: specific context output
                            #     print('found connParams')
                            #     if 'I->E' in current_path:
                            #         print('found I->E')

                            if isinstance(current_obj, dict):
                                for key, val in current_obj.items():
                                    new_path = f"{current_path}.{key}" if current_path else key
                                    if val == value:
                                        matching_paths.append(new_path)
                                    elif isinstance(val, str):  # Handle HOC string matches
                                        if str(value) in val:
                                            matching_paths.append(new_path)
                                    elif isinstance(val, (dict, list)):
                                        stack.append((val, new_path))  # Push deeper layer onto stack

                            elif isinstance(current_obj, list):
                                for i, item in enumerate(current_obj):
                                    new_path = f"{current_path}[{i}]"
                                    if item == value:
                                        matching_paths.append(new_path)
                                    elif isinstance(item, str):  # Handle HOC string matches
                                        if str(value) in item:
                                            matching_paths.append(new_path)
                                    elif isinstance(item, (dict, list)):
                                        stack.append((item, new_path))  # Push list item onto stack

                        return matching_paths  # Return all matching paths
                    
                    def find_name_in_netparams(name, netParams, current_path=""):
                        """
                        Recursively search for the name in netParams and return a list of matching paths.
                        
                        Parameters:
                            name (str): The name to search for.
                            netParams (object): The network parameters object.
                            current_path (str): The current path in the recursive search.
                        
                        Returns:
                            list: A list of paths to the matching name.
                        """
                        stack = [(netParams, current_path)]
                        
                        if '_' in name:
                            elements = name.split('_')
                            try: assert 'E' in elements or 'I' in elements
                            except: elements = None
                        else:
                            elements = None
                        
                        matching_paths = []
                        while stack:
                            current_obj, current_path = stack.pop()
                            
                            # if 'cellParams' in current_path:  # Debugging: specific context output
                            #     print('found cellParams')
                                
                            # if 'gnabar' in current_path:  # Debugging: specific context output
                            #     print('found gnabar')
                                
                            # elements=None
                            # #if _ in 
                            
                            if elements is not None:
                                if isinstance(current_obj, dict):
                                    for key, val in current_obj.items():
                                        new_path = f"{current_path}.{key}" if current_path else key
                                        if all([element in new_path for element in elements]):
                                            matching_paths.append(new_path)
                                        elif isinstance(val, (dict, list)):
                                            stack.append((val, new_path))
                                elif isinstance(current_obj, list):
                                    for i, item in enumerate(current_obj):
                                        new_path = f"{current_path}[{i}]"
                                        if all([element in new_path for element in elements]):
                                            matching_paths.append(new_path)
                                        elif isinstance(item, (dict, list)):
                                            stack.append((item, new_path))
                                
                            elif isinstance(current_obj, dict):
                                for key, val in current_obj.items():
                                    new_path = f"{current_path}.{key}" if current_path else key
                                    #if key == name:
                                    if key == name:
                                        matching_paths.append(new_path)
                                    elif isinstance(val, (dict, list)):
                                        stack.append((val, new_path))
                            elif isinstance(current_obj, list):
                                for i, item in enumerate(current_obj):
                                    new_path = f"{current_path}[{i}]"
                                    if item == name:
                                        matching_paths.append(new_path)
                                    elif isinstance(item, (dict, list)):
                                        stack.append((item, new_path))
                        return matching_paths
                        
                    # Generate the mapping
                    mapping = {}
                    for param, value in simConfig.items():
                        if strategy == 'by_name':
                            #paths = find_value_in_netparams(param, netParams)
                            paths = find_name_in_netparams(param, netParams)
                        elif strategy == 'by_value':
                            #paths = find_name_in_netparams(value, netParams)
                            paths = find_value_in_netparams(value, netParams)
                        else:
                            raise ValueError(f"Invalid strategy: {strategy}")
                        mapping[param] = paths if paths else None  # Assign None if no path is found

                    return mapping
                cfg_to_netparams_mapping = map_cfg_to_netparams(
                    {cfg_param: cfg_val}, 
                    netParams.__dict__.copy(),
                    strategy=strategy
                    )
                mapped_paths = cfg_to_netparams_mapping[cfg_param]
                
                if mapped_paths is None:
                    print(f"WARNING: mapped paths is None.")
                    print(f"No paths found for {cfg_param} = {cfg_val}")
                    return

                # update permuted params
                def getNestedParam(netParams, paramLabel):
                    if '.' in paramLabel: 
                        paramLabel = paramLabel.split('.')
                    if isinstance(paramLabel, list ) or isinstance(paramLabel, tuple):
                        container = netParams
                        for ip in range(len(paramLabel) - 1):
                            if hasattr(container, paramLabel[ip]):
                                container = getattr(container, paramLabel[ip])
                            else:
                                container = container[paramLabel[ip]]
                        return container[paramLabel[-1]]
                for mapped_path in mapped_paths:    
                    current_val = getNestedParam(netParams, mapped_path)
                    #assert cfg_val == current_val, f"Expected {cfg_val} but got {current_val}"
                    if isinstance(current_val, str): #handle hoc strings
                        assert str(cfg_val) in current_val, f"Expected {cfg_val} to be in {current_val}"
                        updated_func = current_val.replace(str(cfg_val), str(cfg[cfg_param]))                    
                        #netParams.setNestedParam(mapped_path, updated_func)
                        before_val = getNestedParam(netParams, mapped_path)
                        netParams.setNestedParam(mapped_path, updated_func)
                        after_val = getNestedParam(netParams, mapped_path)
                        assert before_val != after_val, f"Failed to update {mapped_path} from {before_val} to {after_val}"
                        print(f"Updated {mapped_path} from {before_val} to {after_val}")
                    elif strategy == 'by_name': #special case
                        assert cfg_val == current_val, f"Expected {cfg_val} but got {current_val}"
                        original_val = cfg_val
                        permuted_val = cfg[cfg_param]
                        modifier = permuted_val / original_val  # NOTE: this should end up equal to one of the 
                                                                #       level multipliers
                        #proportion = cfg[cfg_param] / current_val
                        
                        #adjust mean proportinally 
                        #netParams.setNestedParam(mapped_path, current_val * modifier)
                        
                        before_val = getNestedParam(netParams, mapped_path)
                        netParams.setNestedParam(mapped_path, current_val * modifier)
                        after_val = getNestedParam(netParams, mapped_path)
                        assert before_val != after_val, f"Failed to update {mapped_path} from {before_val} to {after_val}"
                        print(f"Updated {mapped_path} from {before_val} to {after_val}")
                    else:
                        assert cfg_val == current_val, f"Expected {cfg_val} but got {current_val}"
                        before_val = getNestedParam(netParams, mapped_path)
                        netParams.setNestedParam(mapped_path, cfg[cfg_param])
                        after_val = getNestedParam(netParams, mapped_path)
                        assert before_val != after_val, f"Failed to update {mapped_path} from {before_val} to {after_val}"
                        print(f"Updated {mapped_path} from {before_val} to {after_val}")  

                # remove previous data
                sim.clearAll()
            
                #remove mapping from netParams #TODO: figure out how to actually take advantage of this
                # if hasattr(netParams, 'mapping'):
                #     del netParams.mapping
                netParams.mapping = {}
            
                # run simulation
                # Create network and run simulation
                sim.initialize(                     # create network object and set cfg and net params
                        simConfig = simConfig,          # pass simulation config and network params as arguments
                        netParams = netParams)
                sim.net.createPops()                # instantiate network populations
                sim.net.createCells()               # instantiate network cells based on defined populations
                sim.net.connectCells()              # create connections between cells based on params
                sim.net.addStims()                  # add stimulation (usually there are none)
                sim.setupRecording()                # setup variables to record for each cell (spikes, V traces, etc)
            
            #enable print
            activate_print()
            prepare_permuted_sim(sim_data_path, cfg, cfg_param, cfg_val)
            
            # if not debug_permuted_sims:
            #     #sim.load(sim_data_path, simConfig=cfg)
            sim.runSim()                        # run parallel Neuron simulation
            sim.gatherData()                    # gather spiking data and cell info from each node
            suppress_print()
        elif not debug: # i.e. if original simConfig is being re-run for official 
                        # sensitivity analysis - confirming that the basic load and re-run method does
                        # in fact generate identical results.
            # #sim.clearAll()
            # if not debug_permuted_sims:
            sim.load(sim_data_path, simConfig=cfg)
            sim.runSim()                        # run parallel Neuron simulation
            sim.gatherData()                    # gather spiking data and cell info from each node
        elif debug: 
            print(f'Runing in debug mode. Only loading simConfig from {sim_data_path}')
            print(f'Simulation will not be re-run. it will, however, be saved to the expected location.')
            sim.load(sim_data_path, simConfig=cfg)
        else:
            print('Not sure how you got here... Something went wrong.')
                
        permuted_data_paths = sim.saveData()                      # save params, cell info and sim output to file (pickle,mat,txt,etc)
        assert len(permuted_data_paths) == 1, "Expected only one data path, the .pkl file. Got more."
        perm_sim_data_path = permuted_data_paths[0]
        # print(f'Permutation {simLabel} successfully ran!')
        # return sim
    except Exception as e:
        print(f'Error running permutation {simLabel}: {e}')
        return
        
    if plot:
        try:        
            reprocess_simulation(
                perm_sim_data_path, 
                reference_data_path, 
                #conv_params_path,
                #target_script_path=target_script_path, 
                #duration_seconds=duration_seconds,
                #save_data=save_data, 
                #overwrite_cfgs=overwrite_cfgs
                DEBUG_MODE=debug,
                #DEBUG_MODE=False
                #trim_start = 5,
            )
        except Exception as e:
            print(f'Error processing permutation {simLabel}: {e}')
            return
        
    print(f'Permutation {simLabel} successfully re-ran!')
    
def run_sensitivity_analysis(
    sim_data_path, 
    output_path,
    #evol_params_path, 
    #sensitivity_analysis_output_path,
    reference_data_path = None,
    plot = False,
    lower_bound=0.2,
    upper_bound=1.8,
    levels=2,
    duration_seconds=1,
    option='serial', #NOTE: options are 'serial' or 'parallel'
    num_workers=None, #NOTE: specify number of workers for parallel option, if None, will allocate as many as possible and distribute threads evenly
    debug=False,
    ):

    # generate permutations
    def generate_permutations(
        sim_data_path, 
        #evol_params_path, 
        saveFolder,
        lower_bound=0.2, 
        upper_bound=1.8, 
        levels=2,
        duration_seconds=1,
        verbose = False
        ):
        #from copy import deepcopy
        #from netpyne import sim    
        
        #cfg_permutations
        cfg_permutations = []
        
        #apparently need to modify simcfg before loading
        simConfig = sim.loadSimCfg(sim_data_path, setLoaded=False)
        
        #modify shared runtime options
        duration_seconds = duration_seconds
        simConfig.duration = 1e3 * duration_seconds  # likewise, I think it's necessary to modify netParams, not net.params or net
        simConfig.verbose = False
        #simConfig.verbose = True # NOTE: during connection formation, this will be VERY verbose
        simConfig.validateNetParams = True
        #simConfig.coreneuron = True
        simConfig.saveFolder=saveFolder
        simConfig.saveJson = False
        simConfig.savePickle = True
        #simConfig.coreneuron = True
        simConfig.cvode_active = False # make sure variable time step is off...not sure why it was ever on.
        simConfig.simLabel = '_'+simConfig.simLabel # NOTE this is only applied to the original simConfig 
                                                    # - only because it isnt overwritten later when generating permutations.
        
        # turn recordings off
        # remove recordCells from simConfig
        if hasattr(simConfig, 'recordCells'):
            delattr(simConfig, 'recordCells')
        
        # append original simConfig to permutations so that it is also run, plot, and saved with the others.
        # structure is maintained as a tuple to match the structure of the other permutations
        cfg_permutations.append((
            simConfig.__dict__.copy(), #cfg
            None, #permuted param
            None, #original value
            ))
        
        # load evol_params
        # evol_params_module = import_module_from_path(evol_params_path)
        # evol_params = evol_params_module.params
        #from DIV21.src.evol_params import params
        from RBS_network_models.CDKL5.DIV21.src.evol_params import params
        evol_params = params
        
        # generate permutations
        for evol_param, evol_val in evol_params.items():
            # for cfg_param, cfg_val in simConfig.items():
            #     if evol_param == cfg_param:
            if hasattr(simConfig, evol_param):
                cfg_param = evol_param
                cfg_val = getattr(simConfig, evol_param)
                #if evol_val is a list, then it's a range from min to max allowed for the parameter
                if isinstance(evol_val, list):
                    
                    #NOTE: it has occured to me that modifying certain params this way just isnt very practical or useful
                    # for example, modifying std of distribution for a param, would require getting all the values of
                    # the param, and somehow remapping them to the new std for each cell. 
                    # I dont think this would be very useful, and would be pretty complicated to implement.
                    # by contrast, if the mean of the distribution was modified, it would be much simpler to just 
                    # shift all the values by the same proportion.
                    
                    #the following if statement will skip over these kinds of params
                    excepted_param_keys = [
                        'std', # see rationale above
                        #'probLengthConst', # NOTE: this is included into a string that passed and evaluated in hoc,
                                            # 'probability': 'exp(-dist_3D / {})*{}'.format(cfg.probLengthConst, spec['prob']),
                                            # it's easier to modify probability to get a sense of how it affects the network
                                            
                                            # NOTE: nvm I figured it out. I can just modify the string directly.
                        ]
                    if any([key in cfg_param for key in excepted_param_keys]):
                        if verbose: print(f'Skipping permutations for {cfg_param}...')
                        continue
                    
                    if verbose: print(f'Generating permutations for {cfg_param}...')
                    
                    
                    # 2025-01-10 10:09:03 aw - original code. 2 levels of permutations, hard coded.
                    # going to implement code that will allow for any number of levels of permutations 
                    # between the lower and upper bounds.
                    
                    # #create two permutations of cfg in this param, 0.2 of the original, and 1.8 of the original
                    # cfg_permutation_1 = simConfig.__dict__.copy()
                    # cfg_permutation_1[cfg_param] = cfg_val * lower_bound
                    # cfg_permutation_1['simLabel'] = f'{cfg_param}_reduced'
                    # cfg_permutations.append((
                    #     cfg_permutation_1, #cfg
                    #     cfg_param, #permuted param
                    #     cfg_val, #original value                    
                    #     ))
                    
                    # cfg_permutation_2 = simConfig.__dict__.copy()
                    # cfg_permutation_2[cfg_param] = cfg_val * upper_bound
                    # cfg_permutation_2['simLabel'] = f'{cfg_param}_increased'
                    # cfg_permutations.append((
                    #     cfg_permutation_2,
                    #     cfg_param,
                    #     cfg_val,                    
                    #     ))
                    
                    
                    # 2025-01-10 10:10:08 aw new code. any number of levels of permutations between lower and upper bounds.
                    def append_permutation_levels(cfg_permutations, 
                                                  simConfig, 
                                                  cfg_param, 
                                                  cfg_val, 
                                                  lower_bound, 
                                                  upper_bound, 
                                                  levels):
                        """
                        Append permutations to the list of cfg permutations for a given parameter.
                        """
                        def permute_param(cfg_permutation, cfg_param, upper_bound, lower_bound, level, levels):                         
                            """
                            Handle special cases where the parameter should not be permuted.
                            """
                            # if verbose: print(f'Skipping permutations for {cfg_param}...')
                            # return cfg_permutations
                            # special cases
                            #print(f'Generating levels for {cfg_param}...')
                            if 'LengthConst' in cfg_param:
                                # got to typical case, this isnt actually a probability based param
                                # TODO: rename this to just length constant later.
                                upper_value = cfg_val * upper_bound
                                lower_value = cfg_val * lower_bound                       
                            elif 'prob' in cfg_param:
                                #modify upper and lower bounds such that probability based params 
                                #dont go below 0 or above 1
                                upper_value = cfg_val * upper_bound
                                lower_value = cfg_val * lower_bound
                                if upper_value > 1:
                                    upper_value = 1
                                if lower_value < 0:
                                    lower_value = 0
                                
                                # #calculate new upper and lower bounds to be used in the permutations
                                # upper_bound = 1 / cfg_val
                                # lower_bound = 0 / cfg_val
                            else:
                                #typical case
                                upper_value = cfg_val * upper_bound
                                lower_value = cfg_val * lower_bound
                                
                            
                            # return permuted value
                            permuted_vals = np.linspace(lower_value, upper_value, levels)
                            permuted_val = permuted_vals[level]
                            cfg_permutation[cfg_param] = permuted_val
                            return cfg_permutation, permuted_vals
                        #print(f'Generating levels for {cfg_param}...')
                        for i in range(levels):
                            original_cfg = simConfig.__dict__.copy()
                            cfg_permutation = simConfig.__dict__.copy()
                            #temp_upper, temp_lower = upper_bound, lower_bound # save original values for future iterations
                            cfg_permutation, permuted_vals = permute_param(cfg_permutation, cfg_param, upper_bound, lower_bound, i, levels)
                            # cfg_permutation[cfg_param] = cfg_val * (lower_bound + (i * (upper_bound - lower_bound) / (levels - 1)))
                            # upper_bound, lower_bound = temp_upper, temp_lower # set back to original values for next iteration
                            cfg_permutation['simLabel'] = f'{cfg_param}_{i}'
                            cfg_permutations.append((
                                cfg_permutation,
                                cfg_param,
                                cfg_val,
                                ))
                            
                            # aw 2025-01-17 08:38:27 - adding validation controls to ensure permutations are correct
                            # seems like there's no issue in this step - however there may be issues in the net preprocessing
                            assert cfg_permutation[cfg_param] == permuted_vals[i], f'Failed to permute {cfg_param} to {permuted_vals[i]}'
                            assert original_cfg != cfg_permutation, f'Failed to permute {cfg_param} to {permuted_vals[i]}'
                        #print(f'Generated permuted vals: {permuted_vals}')
                        return cfg_permutations
                    
                    cfg_permutations = append_permutation_levels(cfg_permutations,
                                                simConfig,
                                                cfg_param,
                                                cfg_val,
                                                lower_bound,
                                                upper_bound,
                                                levels)
                    
                    if verbose: print(f'Permutations generated for {cfg_param}!')
        #debug
        #only keep cfgs where the simLabel contains 'gk' or 'gna'
        # #TODO: figure out the issue with these params
        # cfg_permutations = [cfg_permutation for cfg_permutation in cfg_permutations if 'gk' in cfg_permutation[0]['simLabel'] or 'gna' in cfg_permutation[0]['simLabel']]
        
        print(f'Generated {len(cfg_permutations)} cfg permutations.')
        return cfg_permutations
    cfg_permutations = generate_permutations(sim_data_path, 
                                             output_path,
                                             #evol_params_path, 
                                             #sensitivity_analysis_output_path, 
                                             lower_bound=lower_bound, 
                                             upper_bound=upper_bound,
                                             levels=levels,
                                             duration_seconds=duration_seconds
                                             )

    # run permutation, test individual perms as needed
    if option == 'serial':
        for perm_simConfig in cfg_permutations:
            try:
                # #for debug only
                # if perm_simConfig[1] is None: continue
                # #if 'tau1_inh' not in perm_simConfig[1]: continue
                # if 'prob' not in perm_simConfig[1] or 'LengthConst' in perm_simConfig[1]: continue
                # #for debug only
                
                run_permutation(
                    sim_data_path,
                    reference_data_path = reference_data_path,
                    plot = plot,
                    debug = debug, 
                    *perm_simConfig
                    )
            except Exception as e:
                print(f'Error running permutation {perm_simConfig["simLabel"]}: {e}')
                print('Continuing to next permutation...')
            #break
    elif option == 'parallel':
        # run all permutations, parallelized        
        # TODO: validate this function
        # NOTE: this one seems to work, the ones above, do not.
                # running this one for now to see if it works.
        def run_all_permutations3(sim_data_path, cfg_permutations, plot=None, reference_data_path=None, num_workers=None):
            """
            Run all configuration permutations in parallel, limited by logical CPU availability.

            Args:
                sim_data_path (str): Path to simulation data.
                cfg_permutations (list): List of configuration permutations to run.
                plot: Plot-related information for `run_permutation` (optional).
                reference_data_path (str): Reference data path for `run_permutation` (optional).
            """
            import os
            from concurrent.futures import ProcessPoolExecutor, as_completed

            available_cpus = os.cpu_count()
            
            # set number of workers (i.e. processes, simultaneously running simulations)
            if num_workers is None:
                num_workers = min(len(cfg_permutations), available_cpus)
            else:
                num_workers = min(num_workers, available_cpus)
            print(f'Using {num_workers} workers out of {available_cpus} available CPUs.')

            # Evenly distribute threads among processes
            threads_per_worker = max(1, available_cpus // num_workers)
            os.environ["OMP_NUM_THREADS"] = str(threads_per_worker)
            os.environ["MKL_NUM_THREADS"] = str(threads_per_worker)
            os.environ["OPENBLAS_NUM_THREADS"] = str(threads_per_worker)
            os.environ["NUMEXPR_NUM_THREADS"] = str(threads_per_worker)

            # Prepare tasks
            tasks = []
            for cfg_tuple in cfg_permutations:
                if isinstance(cfg_tuple, tuple) and len(cfg_tuple) == 3:
                    cfg, cfg_param, cfg_val = cfg_tuple
                    tasks.append((sim_data_path, cfg, cfg_param, cfg_val, reference_data_path, plot,))
                else:
                    raise ValueError(f"Unexpected structure in cfg_permutations: {cfg_tuple}")

            # Run tasks in parallel
            with ProcessPoolExecutor(max_workers=num_workers) as executor:
                futures = {executor.submit(run_permutation, *task): task for task in tasks}

                for future in as_completed(futures):
                    task = futures[future]
                    try:
                        future.result()  # This will raise any exceptions from the worker
                    except Exception as e:
                        cfg = task[1]  # Access the configuration from the task
                        sim_label = cfg.get("simLabel", "unknown") if isinstance(cfg, dict) else "unknown"
                        print(f"Unhandled exception in permutation {sim_label}: {e}")
        run_all_permutations3(sim_data_path, cfg_permutations, plot=plot, reference_data_path=reference_data_path, num_workers=num_workers)

'''test net params'''
def add_module_to_sys_path(module_path):
    parent_dir = os.path.dirname(module_path)
    if parent_dir not in sys.path:
        sys.path.append(parent_dir)
    else:
        print(f'{parent_dir} already in sys.path')

'''sort'''
def spike_sort_experimental_data(recording_paths=None, stream_select=None, sorting_output_path=None, **kwargs):
    #assert required arguments
    assert recording_paths is not None, f"Error: recording_paths is None"
    assert sorting_output_path is not None, f"Error: sorting_output_path is None"
    
    #print streams that will be sorted based on stream select or lack there of
    if stream_select is not None:
        print(f"Sorting stream {stream_select} for each recording.")
    else:
        print(f"Sorting all streams for each recording.")    
    
    sorting_objects = {}
    for recording_path in recording_paths:
        sorting_objects[recording_path] = {}
        stream_nums = [0, 1, 2, 3, 4, 5]
        if stream_select is not None: stream_nums = [stream_select]
        for stream_num in stream_nums:
            os.makedirs(sorting_output_path, exist_ok=True)
            print(f"Sorting output path: {sorting_output_path}")
            sorting_object = sort_data(recording_path, stream_num, sorting_output_path)
            print(f"Sorted data for stream {stream_num} at {sorting_output_path}")
            sorting_objects[recording_path][stream_num] = sorting_object      
    return sorting_objects

def sort_data(recording_path, stream_num, sorting_output_path):
    #alloc to enable running in perlmutter:
    #salloc --nodes 1 --qos interactive --time 04:00:00 --constraint gpu --gpus 1 --account m2043_g --image=docker:adammwea/axonkilo_docker:v7
    #shifter --image=docker:adammwea/axonkilo_docker:v7 python /pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/_1_derive_features_from_experimental_data/derive_CDKL5-E6D_T2_C1_DIV21_features.py

    #spike sort using kilosort2
    #from modules.mea_processing_library import load_recordings, kilosort2_wrapper
    #print module location
    print(f"Using kilosort2_wrapper from {mea_lib.__file__}")
    
    #load recording
    h5_file_path = recording_path
    stream_select = stream_num
    recording_object = mea_lib.load_recordings(h5_file_path, stream_select=stream_select)
    stream_objs = recording_object[1]
    
    #sort data
    for wellid, stream_obj in stream_objs.items():
        # don't set sorting params, just do defualt.
        maxwell_recording_extractor = stream_obj['recording_segments'][0] #note: network scans should only have one recording segment
        recording = maxwell_recording_extractor
        output_path = sorting_output_path
        sorting_params = ss.Kilosort2Sorter.default_params()
        sorting_object=mea_lib.kilosort2_wrapper(recording, output_path, sorting_params)
        print(f"Sorting complete for {wellid}")
    
    return sorting_object

def get_default_kilosort2_params():
    return ss.Kilosort2Sorter.default_params()

'''network metrics'''
def add_repo_root_to_sys_path():
    #path_for_external_repos = "/pscratch/sd/a/adammwea/workspace/RBS_neuronal_network_models/"
    repo_path = get_git_root()
    sys.path.insert(0, repo_path)
    #sys.path.append(path_for_external_repos)    

def get_git_root():
    try:
        current_dir = os.path.dirname(os.path.abspath(__file__))
        while current_dir != os.path.dirname(current_dir):  # Root directory check
            if os.path.isdir(os.path.join(current_dir, '.git')):
                return current_dir
            current_dir = os.path.dirname(current_dir)
        return None  # No .git directory found
    except Exception as e:
        print(f"An error occurred while determining the Git root: {e}")
        return None

def get_base_path(recording_path):
    return os.path.dirname(recording_path).replace('xRBS_input_data', 'yRBS_spike_sorted_data')

def check_for_sorting_objects_associated_to_recordings(recording_paths, sorting_output_path):
    #if recording_paths is one path, convert to list
    if isinstance(recording_paths, str):
        recording_paths = [recording_paths]
    
    # look for sorting objects in this directory
    sorting_output_paths = []
    sorting_output_paths = []
    for recording_path in recording_paths:
        #for stream in streams:
        #mea_lib.load_recordings(recording_path, stream_select=0)
        recording_details = mea_lib.extract_recording_details(recording_path)[0]
        project = recording_details['projectName']
        date = recording_details['date']
        chip_id = recording_details['chipID']
        scan_type = recording_details['scanType']
        run_id = recording_details['runID']
        
        for root, dirs, files in os.walk(sorting_output_path):
            if "spike_times.npy" in files:
                #return root
                if project in root and date in root and chip_id in root and scan_type in root and run_id in root:
                    sorting_output_paths.append(root)
            
    return sorting_output_paths

def import_module_from_path(module_path):
    #module_path = os.path.abspath(module_path)
    spec = importlib.util.spec_from_file_location("module.name", module_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module

def plot_raster_plot_experimental(ax, spiking_data_by_unit):
    """Plot a raster plot for spiking data."""
    
    # Calculate the average firing rate for each unit
    firing_rates = {}
    for gid in spiking_data_by_unit:
        spike_times = spiking_data_by_unit[gid]['spike_times']
        spike_times = [spike_times] if isinstance(spike_times, (int, float)) else spike_times
        firing_rate = len(spike_times) / (max(spike_times) - min(spike_times)) if len(spike_times) > 1 else 0
        firing_rates[gid] = firing_rate
    
    # Sort the units based on their average firing rates
    sorted_units = sorted(firing_rates, key=firing_rates.get)
    
    # Create a mapping from original gid to new y-axis position
    gid_to_ypos = {gid: pos for pos, gid in enumerate(sorted_units)}
    
    # Plot the units in the sorted order
    for gid in sorted_units:
        spike_times = spiking_data_by_unit[gid]['spike_times']
        spike_times = [spike_times] if isinstance(spike_times, (int, float)) else spike_times
        ax.plot(spike_times, [gid_to_ypos[gid]] * len(spike_times), 'b.', markersize=2)

    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Unit ID (sorted by firing rate)')
    ax.set_title('Raster Plot')
    plt.tight_layout()
    return ax

# def plot_network_bursting_experimental(ax, bursting_ax):
#     # """Plot network bursting activity."""
#     #ax_old = bursting_data['ax']
#     ax_old = bursting_ax
    
    
#     #copy ax features to new ax
#     ax.set_xlim(ax_old.get_xlim())
#     ax.set_ylim(ax_old.get_ylim())
#     #ax.set_ylabel(ax_old.get_ylabel())
#     ax.set_ylabel('Firing Rate (Hz)')
#     #ax.set_xlabel(ax_old.get_xlabel())
#     ax.set_xlabel('Time (s)')
#     ax.set_title(ax_old.get_title())
#     #ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
#     ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
#     ax.plot(ax_old.get_lines()[1].get_xdata(), ax_old.get_lines()[1].get_ydata(), 'or')
    
#     return ax
    
def generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path):
    """Generate and save the network bursting plot."""
    #fig = network_data['bursting_data']['bursting_summary_data']['fig']
    ax = network_data['bursting_data']['bursting_summary_data']['ax']
    ax_old = ax
    
    # Create a new figure with shared x-axis
    fig, ax = plt.subplots(figsize=(16, 4.5))
    
    #copy ax features to new ax
    ax.set_xlim(ax_old.get_xlim())
    ax.set_ylim(ax_old.get_ylim())
    #ax.set_ylabel(ax_old.get_ylabel())
    ax.set_ylabel('Firing Rate (Hz)')
    #ax.set_xlabel(ax_old.get_xlabel())
    ax.set_xlabel('Time (s)')
    ax.set_title(ax_old.get_title())
    #ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    ax.plot(ax_old.get_lines()[1].get_xdata(), ax_old.get_lines()[1].get_ydata(), 'or')
    

    fig.savefig(bursting_plot_path) #save as pdf
    print(f"Bursting plot saved to {bursting_plot_path}")

    # Save bursting plot as pickle
    with open(bursting_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Bursting plot data saved to {bursting_fig_path}")

def load_spike_sorted_data(path):
    # sorting_object_list = []
    # for path in paths:
    #sorter_output_folder = os.path.abspath(path)
    sorter_output_folder = path
    assert os.path.exists(sorter_output_folder), f"Error: path does not exist: {sorter_output_folder}"
    sorting_object = ss.Kilosort2Sorter._get_result_from_folder(sorter_output_folder)
    #sorting_object_list.append(sorting_object)
    return sorting_object

'''refine convolution params'''
def collect_fitness_data(simulation_run_paths, get_extra_data=False):
    """
    Collect fitness data from JSON files in specified simulation paths.

    Parameters:
        simulation_run_paths (list): List of paths to search for fitness JSON files.

    Returns:
        list: Sorted list of tuples containing (average_fitness, file_path).
    """
    if not get_extra_data:
        fitness_data = []
        for path in simulation_run_paths:
            for root, _, files in os.walk(path):
                for file in files:
                    if file.endswith('_fitness.json'):
                        fitness_path = os.path.join(root, file)
                        try:
                            with open(fitness_path, 'r') as f:
                                fitness_content = json.load(f)
                                average_fitness = fitness_content.get('average_fitness', float('inf'))
                                fitness_data.append((average_fitness, fitness_path))
                        except (json.JSONDecodeError, OSError) as e:
                            print(f"Error reading file {fitness_path}: {e}")
        
        # Sort by average fitness
        fitness_data.sort(key=lambda x: x[0])
        return fitness_data
    else:
        fitness_data_dict = {}
        for path in simulation_run_paths:
            for root, _, files in os.walk(path):
                for file in files:
                    if file.endswith('_fitness.json'):
                        fitness_path = os.path.join(root, file)
                        try:
                            with open(fitness_path, 'r') as f:
                                fitness_content = json.load(f)
                                #average_fitness = fitness_content.get('average_fitness', float('inf'))
                                #extra_data = fitness_content.get('extra_data', {})
                                #fitness_data.append((average_fitness, fitness_path, extra_data))
                                fitness_data_dict[fitness_path] = fitness_content
                        except (json.JSONDecodeError, OSError) as e:
                            print(f"Error reading file {fitness_path}: {e}")
        
        # Sort by average fitness
        return fitness_data_dict

def analyze_simulation(simulation_run_paths, reference=False, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Analyze simulations by generating network summary plots for all `_data.json` files.

    Parameters:
        simulation_run_paths (list): List of paths to search for simulation data files.
        reference (bool): Whether to include reference plots in the analysis.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    for path in simulation_run_paths:
        for root, _, files in os.walk(path):
            for file in files:
                if file.endswith('_data.json'):
                    data_path = os.path.join(root, file)
                    try:
                        generate_network_summary_slide_content(
                            data_path,
                            reference=reference,
                            reference_data=reference_data,
                            #reference_raster=reference_raster,
                            #reference_bursting=reference_bursting,
                        )
                    except Exception as e:
                        print(f"Error analyzing {file} at {data_path}: {e}")

def generate_network_summary_slide_content(data_path=None, 
                                            #reference=False, 
                                            reference_data=None, 
                                            #reference_raster=None, 
                                            #reference_bursting=None, 
                                            convolution_params=None,
                                            **kwargs                                           
                                            ):
    #from _external.RBS_network_simulation_optimization_tools.modules.analysis.analyze_network_activity import get_simulated_network_activity_metrics
    from RBS_network_models.developing.utils.analysis_helper import get_detailed_simulation_data, get_simulated_network_activity_metrics
    
    """
    Analyze simulation data and generate network activity plots.

    Parameters:
        data_path (str): Path to the simulation data file.
        reference (bool): Whether to include reference plots in the summary.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    assert data_path is not None, f"Error: data_path is None"
    
    print('') # Add a newline for clarity
    print(f"Analyzing data from: {data_path}...")

    # Define output paths
    raster_plot_path = data_path.replace('_data.json', '_raster_plot.pdf')
    bursting_plot_path = data_path.replace('_data.json', '_bursting_plot.pdf')
    raster_fig_path = data_path.replace('_data.json', '_raster_fig.pkl')
    bursting_fig_path = data_path.replace('_data.json', '_bursting_fig.pkl')
    summary_plot_path = data_path.replace('_data.json', '_summary_plot.pdf')

    # # Update reference paths if provided
    # if reference_raster and reference_bursting:
    #     reference = True
    #     raster_fig_path, bursting_fig_path = reference_raster, reference_bursting

    try:
        # Extract simulation data
        extracted_data = get_detailed_simulation_data(data_path)
    except Exception as e:
        print(f"Error extracting simulation data from {data_path}: {e}")

    try:
        conv_params = import_module_from_path(convolution_params)
        conv_params = conv_params.conv_params
        extracted_data['conv_params'] = conv_params
        network_data = get_simulated_network_activity_metrics(**extracted_data)
    except Exception as e:
        print(f"Error computing network activity metrics for {data_path}: {e}")
    
    # try:
    #     # Generate raster plot
    #     generate_raster_plot(network_data, raster_plot_path, raster_fig_path)
    # except Exception as e:
    #     print(f"Error generating raster plot for {data_path}: {e}")

    # try:
    #     # Generate network bursting plot
    #     generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path)
    # except Exception as e:
    #     print(f"Error generating network bursting plot for {data_path}: {e}")

    try:
        conv_param_tunning = kwargs.get('conv_param_tunning', False)
        # Generate network activity summary plot
        plot_and_save_network_summary_slide(
            network_data,
            summary_plot_path,
            #reference=reference,
            reference_data=reference_data,
            #reference_raster=reference_raster,
            #reference_bursting=reference_bursting,
            conv_params_path=convolution_params,
            conv_param_tunning=conv_param_tunning,
        )
    except Exception as e:
        print(f"Error generating network activity summary plot for {data_path}: {e}")

    print("Plots generated successfully.")

def plot_and_save_network_summary_slide(network_data, 
                                        summary_plot_path, 
                                        #reference=True, 
                                        reference_data=None, 
                                        #reference_raster=None, reference_bursting=None,
                                        conv_params_path=None,
                                        conv_param_tunning=False):
    """
    Generate a network activity summary plot with raster and bursting subplots.

    Parameters:
        network_data (dict): Contains simulation and bursting data for the network.
        summary_plot_path (str): Path to save the summary plot as a PDF.
        reference (bool): Whether to include reference raster and bursting plots.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids, I_gids = simulated_data['E_Gids'], simulated_data['I_Gids']
    bursting_ax = network_data['bursting_data']['bursting_summary_data']['ax']

    # Determine the number of figure rows based on reference data
    #fig_rows = 4 if reference else 2
    fig_rows = 4
    fig, axs = plt.subplots(fig_rows, 1, figsize=(16, 9), sharex=True)

    # Plot simulated raster
    axs[0].set_title('Raster Plot')
    #from RBS_network_models.developing.utils.analysis_helper import plot_raster_plot_simulated
    plot_raster_plot_simulated(axs[0], spiking_data_by_unit, E_gids, I_gids)
    axs[0].set_ylabel('Neuron GID')

    # Plot simulated bursting
    axs[1].set_title('Bursting Plot')
    axs[1].set_xlabel('Time (s)')
    axs[1].set_ylabel('Firing Rate (Hz)')
    axs[1].plot(bursting_ax.get_lines()[0].get_xdata(), bursting_ax.get_lines()[0].get_ydata(), color='royalblue')
    #axs[1].scatter(bursting_ax.get_lines()[1].get_xdata(), bursting_ax.get_lines()[1].get_ydata(), 'or')
    axs[1].scatter(bursting_ax.get_lines()[1].get_xdata(), bursting_ax.get_lines()[1].get_ydata(), marker='o', color='r')

    # Plot reference data        
    #load reference data
    ref_data = np.load(reference_data, allow_pickle=True).item()
    ref_spike_data = ref_data['spiking_data']['spiking_data_by_unit']
    ref_bursting_ax = ref_data['bursting_data']['bursting_summary_data']['ax']
    
    # Sort reference spike data by FR
    sorted_ref_spike_data = {k: v for k, v in sorted(ref_spike_data.items(), key=lambda item: item[1]['FireRate'])}
    
    # trim reference data to match the length of simulated data
    max_sim_time = max(axs[1].get_lines()[0].get_xdata())
    trimmed_ref_spike_data = trim_spike_data(sorted_ref_spike_data, max_sim_time)
    
    # Plot reference raster
    axs[2] = plot_raster_plot_simulated(axs[2], trimmed_ref_spike_data, [], [], mode='experimental')
    axs[2].set_ylabel('Unit ID')
    axs[2].set_title('Reference Raster Plot')
    axs[2].set_xlabel('Time (s)')
            
    # Plot reference bursting
    axs[3].set_title('Reference Bursting Plot')
    axs[3].set_xlabel('Time (s)')
    axs[3].set_ylabel('Firing Rate (Hz)')
    
    #copy, trim, and plot reference bursting data ax
    max_sim_time = max(axs[1].get_lines()[0].get_xdata())
    axs[3] = trim_and_plot_bursting_data(axs[3], ref_bursting_ax, max_sim_time=max_sim_time)
    
    # Match y-limits across bursting plots
    match_ylims(axs[1], axs[3])
    
    #match x across all plots
    for ax in axs:
        ax.set_xlim([0, max_sim_time])

    # Finalize and save the figure
    plt.tight_layout()
    if not conv_param_tunning:
        #typical case for running in parallel I think
        fig.savefig(summary_plot_path, format='pdf')
        print(f"Summary plot saved to: {summary_plot_path}")
        summary_plot_path = summary_plot_path.replace('.pdf', '.png')
        fig.savefig(summary_plot_path, format='png', dpi=600)
        print(f"Summary plot saved to: {summary_plot_path}")
        print("Summary plot generated successfully.")
    else:
        assert conv_params_path is not None, f"Error: conv_params_path is None"
        #special case when running for conv param tunning
        #make paths for saving in conv_params_path
        conv_params_dir = os.path.dirname(conv_params_path)
        save_path = os.path.join(conv_params_dir, 'summary_plot.pdf')
        fig.savefig(save_path, format='pdf')
        print(f"Summary plot saved to: {save_path}")
        save_path = os.path.join(conv_params_dir, 'summary_plot.png')
        fig.savefig(save_path, format='png', dpi=600)
        print(f"Summary plot saved to: {save_path}")

# def plot_raster_plot_simulated(ax, spiking_data_by_unit, unit_ids=None, E_gids=None, I_gids=None, data_type='simulated'):
#     """
#     Plot a raster plot of spiking activity.

#     Parameters:
#         ax (matplotlib.axes.Axes): The matplotlib axis to draw the plot on.
#         spiking_data_by_unit (dict): A dictionary where keys are neuron IDs (GIDs) 
#                                      and values are spike time data for each unit.
#         E_gids (list): List of excitatory neuron IDs.
#         I_gids (list): List of inhibitory neuron IDs.

#     Returns:
#         matplotlib.axes.Axes: The axis with the raster plot.
#     """
#     # Helper function to plot spikes for a group of neurons
#     def plot_group(gids, color, label, mode='simulated'):
#         if mode == 'simulated':
#             for gid in gids:
#                 if gid in spiking_data_by_unit:
#                     spike_times = spiking_data_by_unit[gid]['spike_times']
#                     if isinstance(spike_times, (int, float)):
#                         spike_times = [spike_times]
#                     ax.plot(spike_times, [gid] * len(spike_times), f'{color}', markersize=1, label=label if gid == gids[0] else None)
#         elif mode == 'experimental':
#             gids = None
#             # Sort spike data by firing rate
#             sorted_spike_data = {k: v for k, v in sorted(spiking_data_by_unit.items(), key=lambda item: item[1]['FireRate'])}
#             sorted_key = 0
#             for gid, data in sorted_spike_data.items():
#                 spike_times = data['spike_times']
#                 ax.plot(spike_times, [sorted_key] * len(spike_times), f'{color}', markersize=1, label=label if gids is None else None)
#                 sorted_key += 1                
                 
#             # for gid, data in spiking_data_by_unit.items():
#             #     spike_times = data['spike_times']
#             #     ax.plot(spike_times, [gid] * len(spike_times), f'{color}.', markersize=2, label=label if gids is None else None)
#             #     gids = gids or [gid]
        

#     # Plot excitatory (yellow) and inhibitory (blue) neurons
#     if mode == 'simulated':
#         plot_group(E_gids, 'y.', 'Excitatory')
#         plot_group(I_gids, 'b.', 'Inhibitory')
#     elif mode == 'experimental':
#         plot_group([], 
#                'k.', #orange
#                'Experimental', mode='experimental')
#         plt.savefig('raster_plot.png')  
#     else:
#         raise ValueError(f"Invalid mode: {mode}")

#     # Customize axis labels and title
#     ax.set_xlabel('Time (s)')
#     ax.set_ylabel('Neuron GID')
#     ax.set_title('Raster Plot')
#     plt.tight_layout()

#     return ax

## updated version of this TODO: update other functions using the function above later
def plot_raster(ax, spiking_data_by_unit, unit_ids=None, E_gids=None, I_gids=None, data_type='simulated'):
        """
        Plot a raster plot of spiking activity.

        Parameters:
            ax (matplotlib.axes.Axes): The matplotlib axis to draw the plot on.
            spiking_data_by_unit (dict): A dictionary where keys are neuron IDs (GIDs) 
                                        and values are spike time data for each unit.
            E_gids (list): List of excitatory neuron IDs.
            I_gids (list): List of inhibitory neuron IDs.

        Returns:
            matplotlib.axes.Axes: The axis with the raster plot.
        """
        
        #main logic
        if data_type == 'simulated':
            list_of_gids = [E_gids, I_gids]
            list_of_colors = ['y.', 'b.']
            list_of_labels = ['Excitatory', 'Inhibitory']
            def plot_mixed(list_of_gids, list_of_colors, list_of_labels):
                assert E_gids is not None, "E_gids must be provided"
                assert I_gids is not None, "I_gids must be provided"

                # Create a dict of all gids and their respective colors and labels
                gid_color_label_dict = {gid: (color, label) for gids, color, label in zip(list_of_gids, list_of_colors, list_of_labels) for gid in gids}
                
                #sort all gids by fr
                sorted_spike_data = {k: v for k, v in sorted(spiking_data_by_unit.items(), key=lambda item: item[1]['FireRate'])}
                sorted_key = 0
                for gid, data in sorted_spike_data.items():
                    spike_times = data['spike_times']
                    #sometimes spike_times is a single value, not a list - make it a list
                    if isinstance(spike_times, (int, float)):
                        spike_times = [spike_times]
                    color, label = gid_color_label_dict[gid]
                    ax.plot(spike_times, [sorted_key] * len(spike_times), f'{color}', markersize=2, label=label)
                    sorted_key += 1
            plot_mixed(list_of_gids, list_of_colors, list_of_labels)
            ax.set_title('Simulated Raster Plot')
            # ax.set_xlabel('Time (s)')
            ax.set_ylabel('Neuron GID')
        elif data_type == 'experimental':            
            color = 'k.'
            label = 'Experimental'
            def plot_exprimental(unit_ids, color, label):           
                #sort
                sorted_spike_data = {k: v for k, v in sorted(spiking_data_by_unit.items(), key=lambda item: item[1]['FireRate'])}
                sorted_key = 0
                for gid, data in sorted_spike_data.items():
                    spike_times = data['spike_times']
                    ax.plot(spike_times, [sorted_key] * len(spike_times), f'{color}', markersize=1, label=label)
                    sorted_key += 1
                
                # #plot
                # for gid in unit_ids:
                #     if gid in spiking_data_by_unit:
                #         spike_times = spiking_data_by_unit[gid]['spike_times']
                #         if isinstance(spike_times, (int, float)):
                #             spike_times = [spike_times]
                #         #ax.plot(spike_times, [gid] * len(spike_times), f'{color}', markersize=1, label=label if gid == gids[0] else None)
                #         ax.plot(spike_times, [gid] * len(spike_times), f'{color}', markersize=1, label=label)
            plot_exprimental(unit_ids, color, label)
            #ax.set_xlabel('Time (s)')
            ax.set_ylabel('Unit ID')
            ax.set_title('Reference Raster Plot')
        else:
            raise ValueError(f"Invalid data_type: {data_type}")
        
        ax.set_title('Raster Plot')
        ax.legend()
        
        # Set the marker size in the legend
        legend = ax.legend()
        for handle in legend.legendHandles:
            #handle._legmarker.set_markersize(5)
            handle._markersize = 5

        # only show unique legend items
        handles, labels = ax.get_legend_handles_labels()
        by_label = dict(zip(labels, handles))
        ax.legend(by_label.values(), by_label.keys())
           
        plt.tight_layout()
        return ax

def plot_experimental_raster(ax, spiking_data_by_unit):
    #TODO: implement this function
    print("Implement this function later")
    return

def trim_spike_data(spiking_data_by_unit, max_sim_time):
    """Trim spike data to fit within the simulation time."""
    trimmed_data = spiking_data_by_unit.copy()
    for gid, data in trimmed_data.items():
        data['spike_times'] = [t for t in data['spike_times'] if t <= max_sim_time]
        trimmed_data[gid] = data   
    return trimmed_data   
    # return {gid: {'spike_times': [t for t in data['spike_times'] if t <= max_sim_time]}
    #         for gid, data in spiking_data_by_unit.items()}

def trim_and_plot_bursting_data(ax, reference_ax, max_sim_time=None):
    ### snipet from MEA Pipeline function where ax is created
    # # Plot the smoothed network activity
    # ax.plot(timeVector, firingRate, color='royalblue')
    # # Restrict the plot to the first and last 100 ms
    # ax.set_xlim([min(relativeSpikeTimes), max(relativeSpikeTimes)])
    # ax.set_ylim([min(firingRate)*0.8, max(firingRate)*1.2])  # Set y-axis limits to min and max of firingRate
    # ax.set_ylabel('Firing Rate [Hz]')
    # ax.set_xlabel('Time [ms]')
    # ax.set_title('Network Activity', fontsize=11)
    #...
    # ax.plot(burstPeakTimes, burstPeakValues, 'or')  # Plot burst peaks as red circles
    
    #copy ax features to new ax
    #ax.set_ylabel(ax_old.get_ylabel())
    #ax.set_xlabel(ax_old.get_xlabel())
    #ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    
    # Copy limits, labels, and data from the original axis
    ax.set_xlim(reference_ax.get_xlim())
    ax.set_ylim(reference_ax.get_ylim())
    #ax.set_xlabel(reference_ax.get_xlabel())
    ax.set_xlabel('Time (s)')
    ax.set_ylabel(reference_ax.get_ylabel())
    ax.set_title(reference_ax.get_title())
    
    # Copy data from the original axis
    #ax.plot(reference_ax.get_lines()[0].get_xdata(), reference_ax.get_lines()[0].get_ydata(), color='royalblue')
    #ax.plot(reference_ax.get_lines()[1].get_xdata(), reference_ax.get_lines()[1].get_ydata(), 'or')   
    #ax.plot(trimmed_x, trimmed_y, color='purple', marker='o')
    
    #trim and plotdata to match the length of the simulated data
    assert max_sim_time is not None, "max_sim_time must be provided"
    x_data, y_data = reference_ax.get_lines()[0].get_xdata(), reference_ax.get_lines()[0].get_ydata()
    x_data, y_data = x_data.copy(), y_data.copy()
    trimmed_x = x_data[x_data <= max_sim_time]
    trimmed_y = y_data[:len(trimmed_x)]
    ax.plot(trimmed_x, trimmed_y, color='purple')
    
    x_data, y_data = reference_ax.get_lines()[1].get_xdata(), reference_ax.get_lines()[1].get_ydata()
    x_data, y_data = x_data.copy(), y_data.copy()
    trimmed_x = x_data[x_data <= max_sim_time]
    trimmed_y = y_data[:len(trimmed_x)]
    ax.scatter(trimmed_x, trimmed_y, color='red', marker='o')
    
    # TODO: maybe implement this later but the above code should work for now
    # # Re-plot data from the original axis
    # for line in reference_ax.get_lines():
    #     ax.plot(line.get_xdata(), line.get_ydata(), color=line.get_color(), marker=line.get_marker())

    return ax

def match_ylims(ax1, ax2):
    """Match the y-axis limits between two axes."""
    min_ylim = min(ax1.get_ylim()[0], ax2.get_ylim()[0])
    max_ylim = max(ax1.get_ylim()[1], ax2.get_ylim()[1])
    padding = 0.1 * (max_ylim - min_ylim)
    ax1.set_ylim(min_ylim - padding, max_ylim + padding)
    ax2.set_ylim(min_ylim - padding, max_ylim + padding)

def get_detailed_simulation_data(data_path):
    """
    Load and process simulation data from the specified path.

    Parameters:
        data_path (str): Path to the simulation data file.

    Returns:
        dict: Extracted simulation data, including simulation configuration,
              network parameters, and population and cell data.
    """
    from netpyne import sim

    # Load simulation components
    sim.loadSimCfg(data_path)
    sim.loadNet(data_path)
    sim.loadNetParams(data_path)
    sim.loadSimData(data_path)

    # Extract and copy data
    extracted_data = {
        'simData': sim.allSimData.copy(),
        'cellData': sim.net.allCells.copy(),
        'popData': sim.net.allPops.copy(),
        'simCfg': sim.cfg.__dict__.copy(),
        'netParams': sim.net.params.__dict__.copy(),
    }

    # Clear all data from the simulation to free memory
    sim.clearAll()

    return extracted_data
    
def analyze_simulations(simulation_run_paths, reference=False, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Analyze simulations by generating network summary plots for all `_data.json` files.

    Parameters:
        simulation_run_paths (list): List of paths to search for simulation data files.
        reference (bool): Whether to include reference plots in the analysis.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    for path in simulation_run_paths:
        for root, _, files in os.walk(path):
            for file in files:
                if file.endswith('_data.json'):
                    data_path = os.path.join(root, file)
                    try:
                        generate_network_summary_slide_content(
                            data_path,
                            reference=reference,
                            reference_data=reference_data,
                            #reference_raster=reference_raster,
                            #reference_bursting=reference_bursting,
                        )
                    except Exception as e:
                        print(f"Error analyzing {file} at {data_path}: {e}")

def collect_network_summary_plots(simulation_run_paths, fitness_data, max_workers=4):
    """
    Collect network summary plots from simulation paths and add them to a PowerPoint presentation.

    Parameters:
        simulation_run_paths (list): Paths to directories containing simulation results.
        fitness_data (list): List of tuples (average_fitness, file_path) from collect_fitness_data.
        file_extension (str): File extension of the summary plot files (default is 'png').
        max_workers (int): Number of parallel processes for processing PDFs.

    Returns:
        None
    """
    #prs = Presentation()
    fitness_dict = {path: fitness for fitness, path in fitness_data}
    #pdf_list = []

    # Process summary plots
    tasks = submit_pdf_tasks(
        #prs,
        fitness_dict,
        simulation_run_paths,
        max_workers,
        debug = False,
        )

    # Handle PDF tasks and update the presentation
    try:
        png_slides_paths, pdf_slides_paths = process_pdf_tasks(tasks)
    except Exception as e:
        print(f"Error processing PDF tasks: {e}")
        return

    #get slide output name based on run path name
    for path in simulation_run_paths:
        
        #divy png and pdf slides by run path
        run_path_basename = os.path.basename(path)
        curated_png_slides_paths = [slide for slide in png_slides_paths if run_path_basename in slide]
        curated_pdf_slides_paths = [slide for slide in pdf_slides_paths if run_path_basename in slide]
        
        # # build prs
        # prs = Presentation()
        # for png_slide_path in png_slides_paths:
        #     slide_layout = prs.slide_layouts[5]
        #     slide = prs.slides.add_slide(slide_layout)
        #     for curated_png_slide_path in curated_png_slides_paths:
        #         slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
        #         print(f"Adding slide: {curated_png_slide_path}")
        
        # import concurrent.futures
        # # from pptx import Presentation
        # # from pptx.util import Inches
        
        # def add_slide(prs, curated_png_slide_path):
        #     slide_layout = prs.slide_layouts[5]
        #     slide = prs.slides.add_slide(slide_layout)
        #     slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
        #     print(f"Adding slide: {curated_png_slide_path}")
        
        # # build prs
        # prs = Presentation()
        # with concurrent.futures.ThreadPoolExecutor() as executor:
        #     futures = [executor.submit(add_slide, prs, curated_png_slide_path) for curated_png_slide_path in curated_png_slides_paths]
        #     concurrent.futures.wait(futures)
            
        #     #wait for all slides to be added
        #     for future in concurrent.futures.as_completed(futures):
        #         print(future.result())
        
        from threading import Lock
        from pptx import Presentation
        from pptx.util import Inches
        import concurrent.futures

        # Define the add_slide function with a lock
        lock = Lock()

        def add_slide(prs, curated_png_slide_path):
            with lock:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
                print(f"Adding slide: {curated_png_slide_path}")

        #build prs
        prs = Presentation()
        # Use ThreadPoolExecutor
        with concurrent.futures.ThreadPoolExecutor() as executor:
            futures = [executor.submit(add_slide, prs, path) for path in curated_png_slides_paths]
            concurrent.futures.wait(futures)
    
        # Save outputs
        save_presentation(prs, f"{run_path_basename}.pptx")
        merge_and_save_pdfs(curated_pdf_slides_paths, f"{run_path_basename}_slides.pdf")

'''feature extraction'''
def build_network_metric_targets_dict(network_metrics):
    network_metric_targets = {
        #General Data
        'source': network_metrics['source'], # 'simulated' or 'experimental'
        #'timeVector': network_metrics['timeVector'],
        
        # Spiking Data
        'spiking_data': {
            #'spike_times': network_metrics['spiking_data']['spike_times'],
            #'spiking_times_by_unit': network_metrics['spiking_data']['spiking_times_by_unit'],
            #'spiking_data_by_unit': network_metrics['spiking_data']['spiking_data_by_unit'],
            'spiking_summary_data': {
                'MeanFireRate': {
                    'target': network_metrics['spiking_data']['spiking_summary_data']['MeanFireRate'],
                    'min': get_min(network_metrics['spiking_data']['spiking_data_by_unit'], 'FireRate'),
                    'max': get_max(network_metrics['spiking_data']['spiking_data_by_unit'], 'FireRate'),
                    'weight': 1, # TODO: update these with Nfactors
                },
                'CoVFireRate': {
                    'target': network_metrics['spiking_data']['spiking_summary_data']['CoVFireRate'],
                    'min': get_min(network_metrics['spiking_data']['spiking_data_by_unit'], 'fr_CoV'),
                    'max': get_max(network_metrics['spiking_data']['spiking_data_by_unit'], 'fr_CoV'),
                    'weight': 1, # TODO: update these with Nfactors
                },
                'MeanISI': {
                    'target': network_metrics['spiking_data']['spiking_summary_data']['MeanISI'],
                    'min': get_min(network_metrics['spiking_data']['spiking_data_by_unit'], 'meanISI'),
                    'max': get_max(network_metrics['spiking_data']['spiking_data_by_unit'], 'meanISI'),
                    'weight': 1, # TODO: update these with Nfactors
                },
                'CoV_ISI': {
                    'target': network_metrics['spiking_data']['spiking_summary_data']['CoV_ISI'],
                    'min': get_min(network_metrics['spiking_data']['spiking_data_by_unit'], 'isi_CoV'),
                    'max': get_max(network_metrics['spiking_data']['spiking_data_by_unit'], 'isi_CoV'),
                    'weight': 1, # TODO: update these with Nfactors
                },
            },
        },
        
        #Bursting Data
        'bursting_data': {
            'bursting_summary_data': {
                'MeanWithinBurstISI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('MeanWithinBurstISI'),
                    'min': get_min(network_metrics['bursting_data']['bursting_data_by_unit'], 'mean_isi_within'),
                    'max': get_max(network_metrics['bursting_data']['bursting_data_by_unit'], 'mean_isi_within'),
                    'weight': 1,
                },
                'CovWithinBurstISI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('CoVWithinBurstISI'),
                    'min': get_min(network_metrics['bursting_data']['bursting_data_by_unit'], 'cov_isi_within'),
                    'max': get_max(network_metrics['bursting_data']['bursting_data_by_unit'], 'cov_isi_within'),
                    'weight': 1,
                },
                'MeanOutsideBurstISI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('MeanOutsideBurstISI'),
                    'min': get_min(network_metrics['bursting_data']['bursting_data_by_unit'], 'mean_isi_outside'),
                    'max': get_max(network_metrics['bursting_data']['bursting_data_by_unit'], 'mean_isi_outside'),
                    'weight': 1,
                },
                'CoVOutsideBurstISI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('CoVOutsideBurstISI'),
                    'min': get_min(network_metrics['bursting_data']['bursting_data_by_unit'], 'cov_isi_outside'),
                    'max': get_max(network_metrics['bursting_data']['bursting_data_by_unit'], 'cov_isi_outside'),
                    'weight': 1,
                },
                'MeanNetworkISI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('MeanNetworkISI'),
                    'min': get_min(network_metrics['bursting_data']['bursting_data_by_unit'], 'mean_isi_all'),
                    'max': get_max(network_metrics['bursting_data']['bursting_data_by_unit'], 'mean_isi_all'),
                    'weight': 1,
                },
                'CoVNetworkISI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('CoVNetworkISI'),
                    'min': get_min(network_metrics['bursting_data']['bursting_data_by_unit'], 'cov_isi_all'),
                    'max': get_max(network_metrics['bursting_data']['bursting_data_by_unit'], 'cov_isi_all'),
                    'weight': 1,
                },
                'NumUnits': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('NumUnits'),
                    # 'min': 1,
                    # 'max': None,
                    # 'weight': 1,
                },
                'Number_Bursts': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('Number_Bursts'),
                    'min': 1,
                    'max': None,
                    'weight': 1, #TODO: update these with Nfactors
                },
                'mean_IBI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('mean_IBI'),
                    'min': None,
                    'max': None,
                    'weight': 1, #TODO: update these with Nfactors
                },
                'cov_IBI': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('cov_IBI'),
                    'min': None,
                    'max': None,
                    'weight': 1, #TODO: update these with Nfactors
                },
                'mean_Burst_Peak': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('mean_Burst_Peak'),
                    'min': None,
                    'max': None,
                    'weight': 1, #TODO: update these with Nfactors
                },
                'cov_Burst_Peak': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('cov_Burst_Peak'),
                    'min': None,
                    'max': None,
                    'weight': 1, #TODO: update these with Nfactors
                },
                'fano_factor': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('fano_factor'),
                    'min': None,
                    'max': None,
                    'weight': 1, #TODO: update these with Nfactors
                },
                'baseline': {
                    'target': network_metrics['bursting_data']['bursting_summary_data'].get('baseline'),
                    'min': None,
                    'max': None,
                    'weight': 1,
                },
            },
            #'bursting_data_by_unit': None,
        }
    }
    return network_metric_targets

def handle_numpy_float64(data):
    try: 
        if isinstance(data, np.ndarray):
            data = data.tolist()
    except: 
        pass
    return data

def get_min(data, key):
    data_min = min(unit[key] for unit in data.values() if unit[key] is not None and unit[key] > 0 and not (unit[key] != unit[key] or unit[key] == float('inf') or unit[key] == float('-inf')))
    #if data is an array with one value, return that value
    # if isinstance(data_min, list) and len(data_min) == 1:
    #     data_min = data_min[0]
    # try: 
    #     len(data_min)
    #     print('data_min is an array')
    # except TypeError: pass
    data_min = handle_numpy_float64(data_min)
    #print('data_min:', data_min)
    return data_min

def get_max(data, key):
    data_max = max(unit[key] for unit in data.values() if not (unit[key] is None or unit[key] != unit[key] or unit[key] == float('inf') or unit[key] == float('-inf')))
    #if data is an array with one value, return that value
    # if type(data_max) == list:
    #     data_max = data_max[0]
    data_max = handle_numpy_float64(data_max)
    return data_max

def save_network_metric_dict_with_timestamp(network_metrics, network_metrics_targets, output_dir):
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    output_path = output_dir + '/fitness_args_' + timestamp + '.py'
    output_path = output_path.replace('-', '_') #conver all '-' to '_' in the timestamp
    print('output_path:', output_path)

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    def format_dict(d, indent=0):
        formatted_str = ''
        for key, value in d.items():
            if isinstance(value, dict):
                formatted_str += ' ' * indent + f"'{key}': {{\n" + format_dict(value, indent + 4) + ' ' * indent + '},\n'
            elif isinstance(value, list):
                formatted_str += ' ' * indent + f"'{key}': [\n"
                for item in value:
                    formatted_str += ' ' * (indent + 4) + f"{item},\n"
                formatted_str += ' ' * indent + '],\n'
            else:
                formatted_str += ' ' * indent + f"'{key}': {value},\n"
        return formatted_str

    # Convert network_metric_targets to a formatted string
    formatted_fitness_args = 'fitness_args = {\n' + format_dict(network_metrics_targets, 4) + '}'

    with open(output_path, 'w') as f:
        f.write(formatted_fitness_args)

    print('Updated fitness args saved to:', output_path)

'''plot network summaries'''
def analyze_simulations_parallel(simulation_run_paths, 
                                 reference=False, 
                                 reference_data=None, 
                                 #reference_raster=None, 
                                 #reference_bursting=None, 
                                 max_workers=4, 
                                 debug=False,
                                 convolution_params=None, 
                                 progress_slides_path = None,
                                 seed_paths=None,
                                 ):
    
    # This is a shortcut for now. #TODO: update this to use the new function
    #global CONVOLUTION_PARAMS, PROGRESS_SLIDES_PATH
    CONVOLUTION_PARAMS = convolution_params
    PROGRESS_SLIDES_PATH = progress_slides_path
    
    def run_process_pool(debug=debug):
        with ProcessPoolExecutor(max_workers=max_workers) as executor:
            for simulation_path in simulation_run_paths:
                if simulation_path.endswith('_data.json'):
                    data_path = simulation_path
                    executor.submit(
                        generate_network_summary_slide_content,
                        data_path,
                        reference=reference,
                        reference_data=reference_data,
                        # reference_raster=reference_raster,
                        # reference_bursting=reference_bursting,
                        convolution_params=CONVOLUTION_PARAMS,
                    )
                else:
                    for root, _, files in os.walk(simulation_path):
                        for file in files:
                            if file.endswith('_data.json'):
                                data_path = os.path.join(root, file)
                                executor.submit(
                                    generate_network_summary_slide_content,
                                    data_path,
                                    reference=reference,
                                    reference_data=reference_data,
                                    # reference_raster=reference_raster,
                                    # reference_bursting=reference_bursting,
                                    convolution_params=CONVOLUTION_PARAMS,
                                )
                            if debug: #if debug mode, only analyze one simulation
                                return
                
    
    """
    Analyze simulations in parallel by generating network summary plots for `_data.json` files.

    Parameters:
        fitness_data (list): List of paths to simulation directories.
        reference (bool): Whether to include reference plots in the analysis.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
        max_workers (int): Maximum number of worker processes for parallel execution.
    """
    import time
    
    # Start the process pool
    start_time = time.time()
    run_process_pool()
    print(f"Time elapsed: {time.time() - start_time:.2f} seconds")




def plot_fitness_trend(data, gen_col='#gen', fitness_col='fitness'):
    """
    Plot fitness scores by generation with a 2nd-order trendline.

    Highlights:
    - Minimum and 50th fitness scores for each generation.
    - Mean and standard deviation for each generation.

    Parameters:
        data (pd.DataFrame): Dataset containing generation and fitness data.
        gen_col (str): Column name for generations.
        fitness_col (str): Column name for fitness scores.
    """
    # Ensure numeric data and remove invalid rows
    data[gen_col] = pd.to_numeric(data[gen_col], errors='coerce')
    data[fitness_col] = pd.to_numeric(data[fitness_col], errors='coerce')
    data.dropna(subset=[gen_col, fitness_col], inplace=True)

    # Group by generation and calculate statistics
    grouped_stats = data.groupby(gen_col)[fitness_col].agg(['mean', 'std']).reset_index()

    # Highlight points for min and 50th fitness in each generation
    highlighted_points = [
        (gen_data.iloc[0], gen_data.iloc[min(49, len(gen_data) - 1)])
        for gen in grouped_stats[gen_col]
        if (gen_data := data[data[gen_col] == gen].sort_values(by=fitness_col))
    ]

    # Fit a 2nd-order polynomial trendline
    x, y = data[gen_col], data[fitness_col]
    coefficients = np.polyfit(x, y, deg=2)
    trendline = np.poly1d(coefficients)

    # Create scatter plot
    plt.figure(figsize=(12, 6))
    plt.scatter(data[gen_col], data[fitness_col], alpha=0.6, s=10, label='Fitness Data')

    # Highlight min and 50th points
    for min_point, mid_point in highlighted_points:
        plt.scatter(min_point[gen_col], min_point[fitness_col], color='green', edgecolors='black', s=50, label='Min Fitness')
        plt.scatter(mid_point[gen_col], mid_point[fitness_col], color='blue', edgecolors='black', s=50, label='50th Fitness')

    # Plot trendline
    x_range = np.linspace(grouped_stats[gen_col].min(), grouped_stats[gen_col].max(), 500)
    plt.plot(x_range, trendline(x_range), color='orange', label='2nd Order Trendline', linewidth=2)

    # Add error bars for mean and standard deviation
    plt.errorbar(grouped_stats[gen_col], grouped_stats['mean'], yerr=grouped_stats['std'],
                 fmt='o', color='purple', label='Mean ± Std Dev', capsize=5)

    # Customize plot
    plt.xlabel('Generation')
    plt.ylabel('Fitness Score')
    plt.title('Fitness Trend with Mean, Std Dev, and Trendline')
    plt.legend(loc='upper right', fontsize='small')
    plt.grid(alpha=0.3)
    plt.tight_layout()
    plt.show()

    return f"Trendline: f(x) = {coefficients[0]:.3f}x^2 + {coefficients[1]:.3f}x + {coefficients[2]:.3f}"

def get_detailed_simulation_data(data_path):
    """
    Load and process simulation data from the specified path.

    Parameters:
        data_path (str): Path to the simulation data file.

    Returns:
        dict: Extracted simulation data, including simulation configuration,
              network parameters, and population and cell data.
    """
    from netpyne import sim

    # Load simulation components
    sim.loadSimCfg(data_path)
    sim.loadNet(data_path)
    sim.loadNetParams(data_path)
    sim.loadSimData(data_path)

    # Extract and copy data
    extracted_data = {
        'simData': sim.allSimData.copy(),
        'cellData': sim.net.allCells.copy(),
        'popData': sim.net.allPops.copy(),
        'simCfg': sim.cfg.__dict__.copy(),
        'netParams': sim.net.params.__dict__.copy(),
    }

    # Clear all data from the simulation to free memory
    sim.clearAll()

    return extracted_data

def plot_raster_plot(ax, spiking_data_by_unit, E_gids, I_gids, mode='simulated'):
    """
    Plot a raster plot of spiking activity.

    Parameters:
        ax (matplotlib.axes.Axes): The matplotlib axis to draw the plot on.
        spiking_data_by_unit (dict): A dictionary where keys are neuron IDs (GIDs) 
                                     and values are spike time data for each unit.
        E_gids (list): List of excitatory neuron IDs.
        I_gids (list): List of inhibitory neuron IDs.

    Returns:
        matplotlib.axes.Axes: The axis with the raster plot.
    """
    # Helper function to plot spikes for a group of neurons
    def plot_group(gids, color, label, mode='simulated'):
        if mode == 'simulated':
            for gid in gids:
                if gid in spiking_data_by_unit:
                    spike_times = spiking_data_by_unit[gid]['spike_times']
                    if isinstance(spike_times, (int, float)):
                        spike_times = [spike_times]
                    ax.plot(spike_times, [gid] * len(spike_times), f'{color}', markersize=1, label=label if gid == gids[0] else None)
        elif mode == 'experimental':
            gids = None
            # Sort spike data by firing rate
            sorted_spike_data = {k: v for k, v in sorted(spiking_data_by_unit.items(), key=lambda item: item[1]['FireRate'])}
            sorted_key = 0
            for gid, data in sorted_spike_data.items():
                spike_times = data['spike_times']
                ax.plot(spike_times, [sorted_key] * len(spike_times), f'{color}', markersize=1, label=label if gids is None else None)
                sorted_key += 1                
                 
            # for gid, data in spiking_data_by_unit.items():
            #     spike_times = data['spike_times']
            #     ax.plot(spike_times, [gid] * len(spike_times), f'{color}.', markersize=2, label=label if gids is None else None)
            #     gids = gids or [gid]
        

    # Plot excitatory (yellow) and inhibitory (blue) neurons
    if mode == 'simulated':
        plot_group(E_gids, 'y.', 'Excitatory')
        plot_group(I_gids, 'b.', 'Inhibitory')
    elif mode == 'experimental':
        plot_group([], 
               'k.', #orange
               'Experimental', mode='experimental')
        plt.savefig('raster_plot.png')  
    else:
        raise ValueError(f"Invalid mode: {mode}")

    # Customize axis labels and title
    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Neuron GID')
    ax.set_title('Raster Plot')
    plt.tight_layout()

    return ax

def generate_raster_plot(network_data, raster_plot_path, raster_fig_path, fig_size=(16, 4.5)):
    """
    Generate and save a raster plot of spiking activity.

    Parameters:
        network_data (dict): Contains simulation data, including spiking activity and neuron IDs.
        raster_plot_path (str): Path to save the raster plot as a PDF.
        raster_fig_path (str): Path to save the raster plot figure as a pickle file.
        fig_size (tuple): Dimensions of the plot (width, height) in inches.
    """
    # Extract relevant data
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids = simulated_data['E_Gids']
    I_gids = simulated_data['I_Gids']

    # Create the raster plot
    fig, ax = plt.subplots(figsize=fig_size)
    plot_raster_plot(ax, spiking_data_by_unit, E_gids, I_gids)

    # Save the raster plot as a PDF
    fig.savefig(raster_plot_path, format='pdf')
    print(f"Raster plot saved to: {raster_plot_path}")

    # Save the figure object as a pickle file
    with open(raster_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Raster plot figure saved to: {raster_fig_path}")

def generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path, fig_size=(16, 4.5)):
    """
    Generate and save the network bursting plot.

    Parameters:
        network_data (dict): Contains bursting summary data and plot axis from the simulation.
        bursting_plot_path (str): Path to save the bursting plot as a PDF.
        bursting_fig_path (str): Path to save the bursting plot figure as a pickle file.
        fig_size (tuple): Dimensions of the plot (width, height) in inches.
    """
    # Extract the original bursting axis
    original_ax = network_data['bursting_data']['bursting_summary_data']['ax']

    # Create a new figure with the specified size
    fig, ax = plt.subplots(figsize=fig_size)

    # Copy limits, labels, and data from the original axis
    ax.set_xlim(original_ax.get_xlim())
    ax.set_ylim(original_ax.get_ylim())
    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Firing Rate (spike count)')
    ax.set_title(original_ax.get_title())
    
    # Re-plot data from the original axis
    # for line in original_ax.get_lines():
    #     ax.plot(line.get_xdata(), line.get_ydata(), line.get_color(), marker=line.get_marker())
    ax.plot(original_ax.get_lines()[0].get_xdata(), original_ax.get_lines()[0].get_ydata(), color='royalblue')
    #ax.scatter(original_ax.get_lines()[1].get_xdata(), original_ax.get_lines()[1].get_ydata(), 'or')
    ax.scatter(original_ax.get_lines()[1].get_xdata(), original_ax.get_lines()[1].get_ydata(), marker='o', color='r')
    
    # Save the bursting plot as a PDF
    fig.savefig(bursting_plot_path, format='pdf')
    print(f"Bursting plot saved to: {bursting_plot_path}")

    # Save the figure object as a pickle file
    with open(bursting_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Bursting plot figure saved to: {bursting_fig_path}")

def generate_network_activity_summary_plot(network_data, summary_plot_path, reference=True, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Generate a network activity summary plot with raster and bursting subplots.

    Parameters:
        network_data (dict): Contains simulation and bursting data for the network.
        summary_plot_path (str): Path to save the summary plot as a PDF.
        reference (bool): Whether to include reference raster and bursting plots.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids, I_gids = simulated_data['E_Gids'], simulated_data['I_Gids']
    bursting_ax = network_data['bursting_data']['bursting_summary_data']['ax']

    # Determine the number of figure rows based on reference data
    fig_rows = 4 if reference else 2
    fig, axs = plt.subplots(fig_rows, 1, figsize=(16, 9), sharex=True)

    # Plot simulated raster
    axs[0].set_title('Raster Plot')
    plot_raster_plot(axs[0], spiking_data_by_unit, E_gids, I_gids)
    axs[0].set_ylabel('Neuron GID')

    # Plot simulated bursting
    axs[1].set_title('Bursting Plot')
    axs[1].set_xlabel('Time (s)')
    axs[1].set_ylabel('Firing Rate (Hz)')
    axs[1].plot(bursting_ax.get_lines()[0].get_xdata(), bursting_ax.get_lines()[0].get_ydata(), color='royalblue')
    #axs[1].scatter(bursting_ax.get_lines()[1].get_xdata(), bursting_ax.get_lines()[1].get_ydata(), 'or')
    axs[1].scatter(bursting_ax.get_lines()[1].get_xdata(), bursting_ax.get_lines()[1].get_ydata(), marker='o', color='r')

    # Plot reference data if provided
    if reference:
        
        #load reference data
        ref_data = np.load(reference_data, allow_pickle=True).item()
        ref_spike_data = ref_data['spiking_data']['spiking_data_by_unit']
        ref_bursting_ax = ref_data['bursting_data']['bursting_summary_data']['ax']
        
        # Sort reference spike data by FR
        sorted_ref_spike_data = {k: v for k, v in sorted(ref_spike_data.items(), key=lambda item: item[1]['FireRate'])}
        
        # trim reference data to match the length of simulated data
        max_sim_time = max(axs[1].get_lines()[0].get_xdata())
        trimmed_ref_spike_data = trim_spike_data(sorted_ref_spike_data, max_sim_time)
        
        # Plot reference raster
        axs[2] = plot_raster_plot(axs[2], trimmed_ref_spike_data, [], [], mode='experimental')
        axs[2].set_ylabel('Unit ID')
        axs[2].set_title('Reference Raster Plot')
        axs[2].set_xlabel('Time (s)')
                
        # Plot reference bursting
        axs[3].set_title('Reference Bursting Plot')
        axs[3].set_xlabel('Time (s)')
        axs[3].set_ylabel('Firing Rate (Hz)')
        
        #copy, trim, and plot reference bursting data ax
        max_sim_time = max(axs[1].get_lines()[0].get_xdata())
        axs[3] = trim_and_plot_bursting_data(axs[3], ref_bursting_ax, max_sim_time=max_sim_time)
        
        # Match y-limits across bursting plots
        match_ylims(axs[1], axs[3])
        
        #match x across all plots
        for ax in axs:
            ax.set_xlim([0, max_sim_time])

    # Finalize and save the figure
    plt.tight_layout()
    fig.savefig(summary_plot_path, format='pdf')
    #if summary plot path ends in .pdf, change to .png
    print(f"Summary plot saved to: {summary_plot_path}")
    summary_plot_path = summary_plot_path.replace('.pdf', '.png')
    fig.savefig(summary_plot_path, format='png', dpi=600)
    print(f"Summary plot saved to: {summary_plot_path}")
    print("Summary plot generated successfully.")

def trim_and_plot_bursting_data(ax, reference_ax, max_sim_time=None):
    ### snipet from MEA Pipeline function where ax is created
    # # Plot the smoothed network activity
    # ax.plot(timeVector, firingRate, color='royalblue')
    # # Restrict the plot to the first and last 100 ms
    # ax.set_xlim([min(relativeSpikeTimes), max(relativeSpikeTimes)])
    # ax.set_ylim([min(firingRate)*0.8, max(firingRate)*1.2])  # Set y-axis limits to min and max of firingRate
    # ax.set_ylabel('Firing Rate [Hz]')
    # ax.set_xlabel('Time [ms]')
    # ax.set_title('Network Activity', fontsize=11)
    #...
    # ax.plot(burstPeakTimes, burstPeakValues, 'or')  # Plot burst peaks as red circles
    
    #copy ax features to new ax
    #ax.set_ylabel(ax_old.get_ylabel())
    #ax.set_xlabel(ax_old.get_xlabel())
    #ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    
    # Copy limits, labels, and data from the original axis
    ax.set_xlim(reference_ax.get_xlim())
    ax.set_ylim(reference_ax.get_ylim())
    #ax.set_xlabel(reference_ax.get_xlabel())
    ax.set_xlabel('Time (s)')
    ax.set_ylabel(reference_ax.get_ylabel())
    ax.set_title(reference_ax.get_title())
    
    # Copy data from the original axis
    #ax.plot(reference_ax.get_lines()[0].get_xdata(), reference_ax.get_lines()[0].get_ydata(), color='royalblue')
    #ax.plot(reference_ax.get_lines()[1].get_xdata(), reference_ax.get_lines()[1].get_ydata(), 'or')   
    #ax.plot(trimmed_x, trimmed_y, color='purple', marker='o')
    
    #trim and plotdata to match the length of the simulated data
    assert max_sim_time is not None, "max_sim_time must be provided"
    x_data, y_data = reference_ax.get_lines()[0].get_xdata(), reference_ax.get_lines()[0].get_ydata()
    x_data, y_data = x_data.copy(), y_data.copy()
    trimmed_x = x_data[x_data <= max_sim_time]
    trimmed_y = y_data[:len(trimmed_x)]
    ax.plot(trimmed_x, trimmed_y, color='purple')
    
    x_data, y_data = reference_ax.get_lines()[1].get_xdata(), reference_ax.get_lines()[1].get_ydata()
    x_data, y_data = x_data.copy(), y_data.copy()
    trimmed_x = x_data[x_data <= max_sim_time]
    trimmed_y = y_data[:len(trimmed_x)]
    ax.scatter(trimmed_x, trimmed_y, color='red', marker='o')
    
    # TODO: maybe implement this later but the above code should work for now
    # # Re-plot data from the original axis
    # for line in reference_ax.get_lines():
    #     ax.plot(line.get_xdata(), line.get_ydata(), color=line.get_color(), marker=line.get_marker())

    return ax

def trim_spike_data(spiking_data_by_unit, max_sim_time):
    """Trim spike data to fit within the simulation time."""
    trimmed_data = spiking_data_by_unit.copy()
    for gid, data in trimmed_data.items():
        data['spike_times'] = [t for t in data['spike_times'] if t <= max_sim_time]
        trimmed_data[gid] = data   
    return trimmed_data   
    # return {gid: {'spike_times': [t for t in data['spike_times'] if t <= max_sim_time]}
    #         for gid, data in spiking_data_by_unit.items()}

def match_ylims(ax1, ax2):
    """Match the y-axis limits between two axes."""
    min_ylim = min(ax1.get_ylim()[0], ax2.get_ylim()[0])
    max_ylim = max(ax1.get_ylim()[1], ax2.get_ylim()[1])
    padding = 0.1 * (max_ylim - min_ylim)
    ax1.set_ylim(min_ylim - padding, max_ylim + padding)
    ax2.set_ylim(min_ylim - padding, max_ylim + padding)

def collect_fitness_data(simulation_run_paths, get_extra_data=False):
    """
    Collect fitness data from JSON files in specified simulation paths.

    Parameters:
        simulation_run_paths (list): List of paths to search for fitness JSON files.

    Returns:
        list: Sorted list of tuples containing (average_fitness, file_path).
    """
    if not get_extra_data:
        fitness_data = []
        for path in simulation_run_paths:
            for root, _, files in os.walk(path):
                for file in files:
                    if file.endswith('_fitness.json'):
                        fitness_path = os.path.join(root, file)
                        try:
                            with open(fitness_path, 'r') as f:
                                fitness_content = json.load(f)
                                average_fitness = fitness_content.get('average_fitness', float('inf'))
                                fitness_data.append((average_fitness, fitness_path))
                        except (json.JSONDecodeError, OSError) as e:
                            print(f"Error reading file {fitness_path}: {e}")
        
        # Sort by average fitness
        fitness_data.sort(key=lambda x: x[0])
        return fitness_data
    else:
        fitness_data_dict = {}
        for path in simulation_run_paths:
            for root, _, files in os.walk(path):
                for file in files:
                    if file.endswith('_fitness.json'):
                        fitness_path = os.path.join(root, file)
                        try:
                            with open(fitness_path, 'r') as f:
                                fitness_content = json.load(f)
                                #average_fitness = fitness_content.get('average_fitness', float('inf'))
                                #extra_data = fitness_content.get('extra_data', {})
                                #fitness_data.append((average_fitness, fitness_path, extra_data))
                                fitness_data_dict[fitness_path] = fitness_content
                        except (json.JSONDecodeError, OSError) as e:
                            print(f"Error reading file {fitness_path}: {e}")
        
        # Sort by average fitness
        return fitness_data_dict

def process_pdf_to_images(fitness_dict, fitness_path, pdf_path):
    """
    Convert PDF pages to high-resolution image objects.

    Parameters:
        pdf_path (str): Path to the PDF file.
        fitness_path (str): Path to the fitness file associated with the PDF.

    Returns:
        list: A list of tuples containing:
            - Candidate name
            - Fitness file path
            - PDF file path
            - Page number
            - Pillow image object
    """
    print(f"Processing PDF: {pdf_path}")
    
    try:
        import re
        # Extract candidate name from the file name
        file_name = os.path.basename(pdf_path)
        cand_name_match = re.search(r"(.+)_summary_plot\.pdf", file_name)
        cand_name = cand_name_match.group(1) if cand_name_match else "Unknown Candidate"
        
        #each pdf file can be treated as a single image
        # Open the PDF file
        #doc = fitz.open(pdf_path)
        #without using fitz
        image = convert_from_path(pdf_path, dpi=600, size=(8000, None), first_page=0, last_page=1)[0]
        #get output image path from fitness path
        output_image_path = os.path.join(os.path.dirname(fitness_path), f"{cand_name}_summary_plot.png")
        
        image.save(output_image_path)
        #save image to local disk
        #image.save(f"workspace/{cand_name}_summary_plot.png")
        #print('image size:', image.size) 
        print(f"Image saved to: {output_image_path}")
        #print(f"Image saved to: workspace/{cand_name}_summary_plot.png")
        # # Convert PDF pages to high-resolution images
        # pil_images = convert_from_path(pdf_path, dpi=600)  # High resolution

        # # Create a list of tuples for each page
        # images = []
        # for page_num, pil_image in enumerate(pil_images, start=1):
        #     # Store the Pillow image object directly
        #     images.append((cand_name, fitness_path, pdf_path, page_num, pil_image))
        
        average_fitness = fitness_dict.get(fitness_path, "N/A")
        output_dir = os.path.dirname(fitness_path)
        image_slide_path, pdf_slide_path = add_slide_with_image(output_image_path, cand_name, average_fitness, output_dir)
        
        images = (cand_name, fitness_path, pdf_path, image_slide_path, pdf_slide_path)
        return images

    except Exception as e:
        print(f"Error processing PDF {pdf_path}: {e}")
        return []    

def add_slide_with_image(image_path, title, average_fitness, output_dir):
    """
    Add a slide with an image and associated text to a PowerPoint presentation.

    Parameters:
        prs (Presentation): The PowerPoint presentation object.
        image_stream (BytesIO): The byte stream of the image to add.
        title (str): Title of the slide (candidate name).
        average_fitness (float): Average fitness score of the candidate.
        output_dir (str): Directory to save the slide as an image and PDF.

    Returns:
        str: Path to the saved PDF of the slide.
    """
    # # Add a slide with the selected layout
    # slide_layout = prs.slide_layouts[5]
    # slide = prs.slides.add_slide(slide_layout)

    # # Add the image to the slide
    # slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(10))

    # # Add title text
    # title_box = slide.shapes.add_textbox(Inches(0), Inches(6.25), Inches(2), Inches(0.5))
    # title_frame = title_box.text_frame
    # title_frame.text = f"Candidate: {title}"
    # title_frame.paragraphs[0].font.size = Pt(12)

    # # Add fitness score text
    # fitness_box = slide.shapes.add_textbox(Inches(0), Inches(6.5), Inches(2), Inches(0.5))
    # fitness_frame = fitness_box.text_frame
    # fitness_frame.text = f"Avg Fitness: {average_fitness:.2f}"
    # fitness_frame.paragraphs[0].font.size = Pt(12)
    
    # # Define paths for saving the slide as an image and PDF
    input_image_path = image_path

    # Define paths for saving the slide as an image and PDF
    image_path = os.path.join(output_dir, f"{title}_summary_slide.png")
    pdf_path = os.path.join(output_dir, f"{title}_summary_slide.pdf")

    # Save slide content as an image and a PDF
    generate_image_and_pdf(input_image_path, title, average_fitness, image_path, pdf_path)
    
    # #load image and add to slide
    # slide_layout = prs.slide_layouts[5]
    # slide = prs.slides.add_slide(slide_layout)
    # slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10))

    return image_path, pdf_path

def generate_image_and_pdf(input_image_path, title, average_fitness, image_path, pdf_path):
    """
    Generate and save an image and PDF of a candidate's summary slide.

    Parameters:
        image_stream (BytesIO): Byte stream of the image to display.
        title (str): Candidate name to display on the slide.
        average_fitness (float): Average fitness score to display on the slide.
        image_path (str): Path to save the image (PNG format).
        pdf_path (str): Path to save the PDF.
    """
    #from PIL import Image

    # Create a figure and axis
    fig, ax = plt.subplots(figsize=(10, 7.5))
    ax.axis("off")
    
    #load image from file
    img = Image.open(input_image_path)
    #print('image size:', img.size)
    ax.imshow(img)
    
    # Add title and fitness score as text
    #replace _summary_plot.ext with _data.json in input_image_path to get data_path
    import re
    data_path = re.sub(r"_summary_plot\..*", "_data.json", input_image_path)
    #data_path = input_image_path.replace("_summary_plot.png", "_data.json")
    plt.text(
        0.15,  # x position for 1-inch margin (10% of the figure width)
        0.15, # y position for 1-inch margin (10% of the figure height)
        f"Candidate: {title}\nAvg Fitness: {average_fitness:.2f}\ndata_path: {data_path}",
        fontsize=7.5,
        ha="left",
        transform=fig.transFigure,
        #image_path = input_image_path,
    )

    # Save the figure as an image
    plt.savefig(image_path, format="png", bbox_inches="tight", pad_inches=0, dpi=600)
    print(f"Image saved to: {image_path}")

    # Save the figure as a PDF
    plt.savefig(pdf_path, format="pdf", bbox_inches="tight", pad_inches=0, dpi=600)
    print(f"PDF saved to: {pdf_path}")

    # Close the figure to free memory
    plt.close(fig)
    
    return image_path

def merge_pdfs(pdf_list, output_path):
    #import PyPDF2
    from PyPDF2 import PdfMerger
    #from PyPF2 import PdfMerger.
    """
    Merge multiple PDF files into a single PDF.

    Parameters:
        pdf_list (list): List of paths to the PDF files to merge.
        output_path (str): Path to save the merged PDF.

    Returns:
        None
    """
    try:
        # Initialize the PDF merger
        merger = PdfMerger()

        # Append each PDF to the merger
        for pdf in pdf_list:
            merger.append(pdf)

        # Write the merged PDF to the output path
        merger.write(output_path)
        print(f"Merged PDF saved to: {output_path}")

    except Exception as e:
        print(f"Error merging PDFs: {e}")

    finally:
        # Ensure the merger is closed
        merger.close()
        
def submit_pdf_tasks(fitness_dict, simulation_run_paths, seed_paths, max_workers, debug=False):
    import time
    
    def submit_tasks_to_process_executor(debug=False):
        
        def return_tasks(tasks):
            task_states = [task._state for task in tasks]
            
            #if any state=PENDING, wait and check again until all tasks are not PENDING
            if debug:
                while 'PENDING' or 'RUNNING' in task_states:
                    task_states = [task._state for task in tasks]
                    #print(f"Task states: {task_states}")
                    num_pending = task_states.count('PENDING')
                    num_running = task_states.count('RUNNING')
                    print(f"Tasks running: {num_running}")
                    print(f"Tasks pending: {num_pending}")
                    if num_pending == 0 and num_running == 0:
                        break
                    time.sleep(5)
                    
                print(f"Task states: {task_states}")
            
            not_finished = [state for state in task_states if state != 'FINISHED']
            not_finished_count = len(not_finished)
            return tasks, tasks_submitted, not_finished, not_finished_count
        
        def submit(tasks_submitted=0):
            network_summary_path = os.path.join(root, file)
            cand_name = file.split("_summary_plot")[0]
            fitness_path = os.path.join(root, f"{cand_name}_fitness.json")
            pdf_path = network_summary_path

            if file_extension == "pdf":
                tasks.append(
                    executor.submit(
                        process_pdf_to_images, 
                        #prs,
                        fitness_dict,
                        fitness_path,
                        pdf_path, 
                    )
                )
                tasks_submitted += 1
            
            #If in debug mode, only submit one task
            if debug and tasks_submitted > 0:
                return return_tasks(tasks)
            
        task_states = [task._state for task in tasks]

        tasks_submitted = 0
        error_state=[]
        with ProcessPoolExecutor(max_workers=max_workers) as executor:
            for path in simulation_run_paths:
                #if path is not file:
                if not os.path.isfile(path):
                    for root, _, files in os.walk(path):
                        for file in files:
                            if file.endswith(f"_summary_plot.{file_extension}"):
                                
                                #if seed_paths is None, submit all tasks
                                if seed_paths is None:
                                    submit(tasks_submitted=tasks_submitted)
                                    
                                #if seed_paths is not None, only submit if seed_name is in seed_paths
                                else:
                                    cand_name = re.search(r"(.+)_summary_plot\..*", file).group(1)
                                    if any(cand_name in path for path in seed_paths):
                                        submit(tasks_submitted=tasks_submitted)
                                        
                                        # else:
                #     file = os.path.basename(path)
                #     root = os.path.dirname(path)
                #     #network_summary_path = path
                #     #get simulation run path from network_summary_path
                #     gen_path = os.path.dirname(path)
                #     #get candidate from path, assuming format _gen_mmm_cand_nnn_data.json
                #     cand_name = re.search(r"gen_\d_cand_\d+", path).group(0)

                #     if file.endswith(f"_summary_plot.{file_extension}"):
                #         submit(tasks_submitted=tasks_submitted)
            

            
            #if any state=PENDING, wait and check again until all tasks are not PENDING
            while 'PENDING' or 'RUNNING' in task_states:
                task_states = [task._state for task in tasks]
                #print(f"Task states: {task_states}")
                num_pending = task_states.count('PENDING')
                num_running = task_states.count('RUNNING')
                print(f"Tasks running: {num_running}")
                print(f"Tasks pending: {num_pending}")
                num_finished = task_states.count('FINISHED')
                print(f"Tasks finished: {num_finished}")
                #error_bool = any([task._exception is not None for task in tasks])
                error_count = sum([task._exception is not None for task in tasks])
                print(f"Tasks with errors: {error_count}")
                if num_pending == 0 and num_running == 0:
                    break
                time.sleep(5)
        
        # Return the list of tasks submitted to the executor
        return return_tasks(tasks)
    
    """
    Submit tasks for processing PDFs in parallel.

    Parameters:
        simulation_run_paths (list): Paths to directories containing simulation results.
        fitness_dict (dict): Dictionary mapping fitness file paths to their scores.
        file_extension (str): File extension of the summary plot files (default is 'png').
        max_workers (int): Number of parallel processes for processing PDFs.

    Returns:
        list: A list of tasks submitted to the executor.
    """
    start_time = time.time()
    tasks = []
    file_extension = "pdf"
    tasks, tasks_submitted, not_finished, not_finished_count = submit_tasks_to_process_executor(debug=debug) #debug just forces one task to be subnitted for easier debugging
    print(f"Tasks submitted: {tasks_submitted}")
    print(f"Tasks not finished: {not_finished_count}")
    print(f"Time elapsed: {time.time() - start_time:.2f} seconds")
    return tasks

def process_pdf_tasks(tasks):
    """
    Process PDF tasks and add slides to the presentation.

    Parameters:
        tasks (list): List of tasks for processing PDFs.
        prs (Presentation): PowerPoint presentation object to update.
        fitness_dict (dict): Dictionary mapping fitness file paths to their scores.
        pdf_list (list): List to append generated PDF slide paths.

    Returns:
        None
    """
    pdf_list = []
    png_list = []
    
    #prs = Presentation()
    
    error_state_count = 0
    for task in tasks:
        try:
            image = task.result() 
            #for cand_name, fitness_path, pdf_path, output_image_path in images:
            cand_name, fitness_path, pdf_path, png_slide_path, pdf_slide_path = image
            assert cand_name in fitness_path, f"Fitness path mismatch: {cand_name} vs. {fitness_path}"
            # average_fitness = fitness_dict.get(fitness_path, "N/A")
            # output_dir = os.path.dirname(fitness_path)
            #pdf_slide_path = add_slide_with_image(prs, image_stream, cand_name, average_fitness, output_dir)\
            #pdf_slide_path = add_slide_with_image(prs, output_image_path, cand_name, average_fitness, output_dir)
            
            png_list.append(png_slide_path)
            pdf_list.append(pdf_slide_path)
        except Exception as e:
            error_state_count += 1
            print(f"Error processing a task: {e}")
        
    #merge pdfs
    #merge_pdfs(pdf_list, os.path.join(PROGRESS_SLIDES_PATH, "merged_slides.pdf"))
    print(f"Error state count: {error_state_count}")
    # this error is expected: Error processing a task: not enough values to unpack (expected 5, got 0) - just a bad simulation run        
    return png_list, pdf_list

def save_presentation(prs, pptx_filename, progress_slides_path=None):
    """
    Save the PowerPoint presentation to disk.

    Parameters:
        prs (Presentation): PowerPoint presentation object.
        pptx_filename (str): Name of the PowerPoint file.
        progress_slides_path (str): Path to save the PowerPoint file.

    Returns:
        None
    """
    # look for PROGRESS_SLIDES_PATH in globals
    #PROGRESS_SLIDES_PATH = globals().get('PROGRESS_SLIDES_PATH', None)
    #assert PROGRESS_SLIDES_PATH is not None, "PROGRESS_SLIDES_PATH not found in globals"
    
    if progress_slides_path is None:
        progress_slides_path = PROGRESS_SLIDES_PATH
    pptx_path = os.path.join(progress_slides_path, pptx_filename)
    pptx_dir = os.path.dirname(pptx_path)
    if not os.path.exists(pptx_dir):
        os.makedirs(pptx_dir)
    prs.save(pptx_path)
    print(f"PowerPoint presentation saved to: {pptx_path}")

def merge_and_save_pdfs(pdf_list, pdf_filename, progress_slides_path=None):
    """
    Merge PDFs into a single file and save it.

    Parameters:
        pdf_list (list): List of PDF file paths to merge.
        pdf_filename (str): Name of the merged PDF file.

    Returns:
        None
    """
    if progress_slides_path is None:
        progress_slides_path = PROGRESS_SLIDES_PATH
        
    pdf_path = os.path.join(progress_slides_path, pdf_filename)
    merge_pdfs(pdf_list, pdf_path)
    #print(f"Merged PDF saved to: {pdf_path}")
    print(f"PDFs successfully merged and saved to: {pdf_path}")

def analyze_simulations_dep(simulation_run_paths, reference=False, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Analyze simulations by generating network summary plots for all `_data.json` files.

    Parameters:
        simulation_run_paths (list): List of paths to search for simulation data files.
        reference (bool): Whether to include reference plots in the analysis.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    for path in simulation_run_paths:
        for root, _, files in os.walk(path):
            for file in files:
                if file.endswith('_data.json'):
                    data_path = os.path.join(root, file)
                    try:
                        generate_network_summary_slide_content(
                            data_path,
                            reference=reference,
                            reference_data=reference_data,
                            #reference_raster=reference_raster,
                            #reference_bursting=reference_bursting,
                        )
                    except Exception as e:
                        print(f"Error analyzing {file} at {data_path}: {e}")

def collect_network_summary_plots_parallel(simulation_run_paths, 
                                           fitness_data, 
                                           seed_paths = None,
                                           progress_slides_path = None,
                                           max_workers=4):
    """
    Collect network summary plots from simulation paths and add them to a PowerPoint presentation.

    Parameters:
        simulation_run_paths (list): Paths to directories containing simulation results.
        fitness_data (list): List of tuples (average_fitness, file_path) from collect_fitness_data.
        file_extension (str): File extension of the summary plot files (default is 'png').
        max_workers (int): Number of parallel processes for processing PDFs.

    Returns:
        None
    """
    #prs = Presentation()
    fitness_dict = {path: fitness for fitness, path in fitness_data}
    #pdf_list = []

    # Process summary plots
    tasks = submit_pdf_tasks(
        #prs,
        fitness_dict,
        simulation_run_paths,
        seed_paths,
        max_workers,
        debug = False,
        )

    # Handle PDF tasks and update the presentation
    try:
        png_slides_paths, pdf_slides_paths = process_pdf_tasks(tasks)
    except Exception as e:
        print(f"Error processing PDF tasks: {e}")
        return

    #get slide output name based on run path name
    for path in simulation_run_paths:
        
        #divy png and pdf slides by run path
        run_path_basename = os.path.basename(path)
        curated_png_slides_paths = [slide for slide in png_slides_paths if run_path_basename in slide]
        curated_pdf_slides_paths = [slide for slide in pdf_slides_paths if run_path_basename in slide]
        
        # # build prs
        # prs = Presentation()
        # for png_slide_path in png_slides_paths:
        #     slide_layout = prs.slide_layouts[5]
        #     slide = prs.slides.add_slide(slide_layout)
        #     for curated_png_slide_path in curated_png_slides_paths:
        #         slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
        #         print(f"Adding slide: {curated_png_slide_path}")
        
        # import concurrent.futures
        # # from pptx import Presentation
        # # from pptx.util import Inches
        
        # def add_slide(prs, curated_png_slide_path):
        #     slide_layout = prs.slide_layouts[5]
        #     slide = prs.slides.add_slide(slide_layout)
        #     slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
        #     print(f"Adding slide: {curated_png_slide_path}")
        
        # # build prs
        # prs = Presentation()
        # with concurrent.futures.ThreadPoolExecutor() as executor:
        #     futures = [executor.submit(add_slide, prs, curated_png_slide_path) for curated_png_slide_path in curated_png_slides_paths]
        #     concurrent.futures.wait(futures)
            
        #     #wait for all slides to be added
        #     for future in concurrent.futures.as_completed(futures):
        #         print(future.result())
        
        from threading import Lock
        from pptx import Presentation
        from pptx.util import Inches
        import concurrent.futures

        # Define the add_slide function with a lock
        lock = Lock()

        def add_slide(prs, curated_png_slide_path):
            with lock:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
                print(f"Adding slide: {curated_png_slide_path}")

        #build prs
        prs = Presentation()
        # Use ThreadPoolExecutor
        with concurrent.futures.ThreadPoolExecutor() as executor:
            futures = [executor.submit(add_slide, prs, path) for path in curated_png_slides_paths]
            concurrent.futures.wait(futures)
    
        # Save outputs
        slide_deck_name = f"{run_path_basename}_slides"
        if seed_paths is not None:
            slide_deck_name = f"{run_path_basename}_slides_seeds"
        save_presentation(
            prs,
            f"{slide_deck_name}.pptx", 
            progress_slides_path=progress_slides_path)
        merge_and_save_pdfs(
            curated_pdf_slides_paths, 
            f"{slide_deck_name}.pdf",
            progress_slides_path=progress_slides_path)

'''rerun simulation of interest'''

#control printing
import builtins
original_print = builtins.print # Save the original print function

def suppress_print():
    """Suppress all print statements."""
    builtins.print = lambda *args, **kwargs: None

def activate_print():
    """Activate print statements."""
    builtins.print = original_print

def privileged_print(*args, **kwargs):
    """Print messages with higher privileges."""
    original_print(*args, **kwargs)

'''deprecated'''
def get_presumed_sorting_output_path(recording_path, stream, **kwargs):
    base_path = get_base_path(recording_path)

    #generate initial path
    path = os.path.join(base_path, f"well{str(stream).zfill(3)}/sortings/sorter_output")

    #replace xInputs with yThroughput
    path = path.replace('xInputs', 'yThroughput')
    path = path.replace('xRBS_input_data', 'yRBS_spike_sorted_data')

    #assert replacements have been made
    assert 'xInputs' not in path, f"Error: 'xInputs' not replaced in path: {path}"
    assert 'xRBS_input_data' not in path, f"Error: 'xRBS_input_data' not replaced in path: {path}"    

    return path