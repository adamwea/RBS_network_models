import os
from PyPDF2 import PdfMerger
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
import re
from concurrent.futures import ThreadPoolExecutor

# ===============================================================================
outputs_dir = "/pscratch/sd/a/adammwea/workspace/RBS_network_models/optimizing/CDKL5/CDKL5_DIV21/outputs"
pdf_dir = "/pscratch/sd/a/adammwea/workspace/RBS_network_models/optimizing/CDKL5/CDKL5_DIV21/slides"
# ===============================================================================
keyword_exclude = []
keyword_include = [
    'cand_',
    'gen_',
    '_slide',
]

def process_pdf_to_ppt(pair):
    run_name, gen_name, pdf_files, pdf_dir = pair

    # Combine the PDFs into one
    combined_pdf_path = os.path.join(pdf_dir, f"{run_name}_{gen_name}.pdf")
    merger = PdfMerger()
    try:
        for pdf in pdf_files:
            merger.append(pdf)
        merger.write(combined_pdf_path)
        merger.close()
        print(f"Combined PDF saved to {combined_pdf_path}")
    except Exception as e:
        print(f"Error combining PDFs for {run_name}/{gen_name}: {e}")
        return

    # Convert the combined PDF to a PowerPoint presentation
    presentation = Presentation()
    try:
        pages = convert_from_path(combined_pdf_path)
        for page in pages:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            image_path = os.path.join(pdf_dir, "temp_page.png")
            page.save(image_path, "PNG")

            # Add the image to the slide
            left = top = Inches(0)
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            slide.shapes.add_picture(image_path, left, top, width=slide_width, height=slide_height)

            # Clean up the temporary image
            os.remove(image_path)

        pptx_name = f"{run_name}_{gen_name}.pptx"
        pptx_path = os.path.join(pdf_dir, pptx_name)
        presentation.save(pptx_path)
        print(f"PowerPoint presentation saved to {pptx_path}")
    except Exception as e:
        print(f"Error converting PDF to PPT for {run_name}/{gen_name}: {e}")

def combine_pdfs_to_pptx(
    outputs_dir, 
    pdf_dir, 
    combined_pdf_name="all_pdfs.pdf", 
    pptx_name="all_pdfs.pptx",
    keyword_exclude=None,
    keyword_include=None,
    num_workers=None
):
    # Ensure the PDF directory exists
    os.makedirs(pdf_dir, exist_ok=True)

    # Find all PDF files in the outputs directory
    pdf_files = []
    for root, _, files in os.walk(outputs_dir):
        for file in files:
            if file.endswith(".pdf"):
                if keyword_exclude is not None and any([keyword in file for keyword in keyword_exclude]):
                    continue
                if keyword_include is not None and not all([keyword in file for keyword in keyword_include]):
                    continue
                pdf_files.append(os.path.join(root, file))

    if not pdf_files:
        print("No PDF files found in the outputs directory.")
        return

    # Organize PDFs into a structure for processing
    pdf_groups = {}
    for pdf in pdf_files:
        try:
            run_name = os.path.basename(os.path.dirname(os.path.dirname(pdf)))
            gen_name = os.path.basename(os.path.dirname(pdf))

            if run_name not in pdf_groups:
                pdf_groups[run_name] = {}
            if gen_name not in pdf_groups[run_name]:
                pdf_groups[run_name][gen_name] = []

            pdf_groups[run_name][gen_name].append(pdf)
        except Exception as e:
            print(f"Error organizing file {pdf}: {e}")
            continue

    # Prepare tasks for parallel processing
    tasks = [
        (run_name, gen_name, pdf_files, pdf_dir)
        for run_name, runs in pdf_groups.items()
        for gen_name, pdf_files in runs.items()
    ]

    # Set environment variables for thread and worker management
    available_cpus = os.cpu_count()
    if num_workers is None:
        num_workers = min(len(tasks), available_cpus)
    else:
        num_workers = min(num_workers, available_cpus)
    print(f'Using {num_workers} workers out of {available_cpus} available CPUs.')

    threads_per_worker = max(1, available_cpus // num_workers)
    os.environ["OMP_NUM_THREADS"] = str(threads_per_worker)
    os.environ["MKL_NUM_THREADS"] = str(threads_per_worker)
    os.environ["OPENBLAS_NUM_THREADS"] = str(threads_per_worker)
    os.environ["NUMEXPR_NUM_THREADS"] = str(threads_per_worker)

    # Process each group in parallel
    with ThreadPoolExecutor(max_workers=num_workers) as executor:
        executor.map(process_pdf_to_ppt, tasks)

combine_pdfs_to_pptx(
    outputs_dir, 
    pdf_dir,
    keyword_exclude=keyword_exclude,
    keyword_include=keyword_include,
    num_workers=None
)
