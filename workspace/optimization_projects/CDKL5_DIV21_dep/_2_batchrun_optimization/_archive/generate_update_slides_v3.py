from setup_environment import set_pythonpath
set_pythonpath()
import os
import json
import re
import analysis_functions as af
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import pptx.enum.text as PP_ALIGN
from pptx.dml.color import RGBColor

'''helper functions'''

'''main script'''

progress_slides_path = '/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/progress_slides'

simulation_run_paths = [
    '/pscratch/sd/a/adammwea/workspace/yThroughput/zRBS_network_simulation_outputs/CDKL5_DIV21/241126_Run2_improved_netparams',
]

fitness_data = []

# Collect fitness data
for simulation_run_path in simulation_run_paths:
    for root, dirs, files in os.walk(simulation_run_path):
        for file in files:
            if file.endswith('_fitness.json'):
                fitness_path = os.path.join(root, file)
                with open(fitness_path, 'r') as f:
                    fitness_content = json.load(f)
                    average_fitness = fitness_content.get('average_fitness', float('inf'))
                    fitness_data.append((average_fitness, root))

# Sort by average fitness (least to greatest)
fitness_data.sort()

# Create a PowerPoint presentation
prs = Presentation()

# Analyze simulations in order of fitness
for average_fitness, simulation_path in fitness_data:
    for root, dirs, files in os.walk(simulation_path):
        for file in files:
            if file.endswith('_data.json'):
                data_path = os.path.join(root, file)
                try:
                    af.analyze_simulation_data(data_path)
                except Exception as e:
                    print(f'Error analyzing data for file: {file}, Error: {e}')
                    
# collect network_summary plots in analyzed simulation folders
for simulation_run_path in simulation_run_paths:
    for root, dirs, files in os.walk(simulation_run_path):
        for file in files:
            if file.endswith('_network_summary.png'):
                network_summary_path = os.path.join(root, file)
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.add_picture(network_summary_path, Inches(0), Inches(0), width=Inches(10))
                
    #prs.save(os.path.join(progress_slides_path, 'CDKL5_DIV21_progress_slides.pptx'))
print('Done analyzing simulations')
                    
#how to run on interacive node in perlmutter command line:
#salloc --nodes=1 --ntasks-per-node=256 -C cpu -q interactive -t 04:00:00
'''
module load conda
conda activate netsims_env
srun --tasks-per-node=256 --cpu-bind=cores python /pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/generate_update_slides_v3.py
'''
#module load conda
#conda activate netsims_env
#srun -n 256 --cpu-bind=cores python /pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/generate_update_slides_v3.py
#srun --tasks-per-node=256 --cpu-bind=cores python /pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/generate_update_slides_v3.py
