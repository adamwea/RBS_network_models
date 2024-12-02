from setup_environment import set_pythonpath
set_pythonpath()

import os
import json
import re
import io
import pickle
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches, Pt
from PyPDF2 import PdfMerger
from PIL import Image
from concurrent.futures import ProcessPoolExecutor

import fitz

from pdf2image import convert_from_path
from PIL import Image
import os
import io
import re

import os

from netpyne import sim
from modules.analysis.analyze_network_activity import get_simulated_network_activity_metrics
from workspace.RBS_network_simulations.workspace.optimization_projects.CDKL5_DIV21._2_batchrun_optimization.calculate_fitness import fitnessFunc

'''helper functions'''

def plot_fitness_trend(data, gen_col='#gen', fitness_col='fitness'):
    """
    Plot fitness scores by generation with a 2nd-order trendline.

    Highlights:
    - Minimum and 50th fitness scores for each generation.
    - Mean and standard deviation for each generation.

    Parameters:
        data (pd.DataFrame): Dataset containing generation and fitness data.
        gen_col (str): Column name for generations.
        fitness_col (str): Column name for fitness scores.
    """
    # Ensure numeric data and remove invalid rows
    data[gen_col] = pd.to_numeric(data[gen_col], errors='coerce')
    data[fitness_col] = pd.to_numeric(data[fitness_col], errors='coerce')
    data.dropna(subset=[gen_col, fitness_col], inplace=True)

    # Group by generation and calculate statistics
    grouped_stats = data.groupby(gen_col)[fitness_col].agg(['mean', 'std']).reset_index()

    # Highlight points for min and 50th fitness in each generation
    highlighted_points = [
        (gen_data.iloc[0], gen_data.iloc[min(49, len(gen_data) - 1)])
        for gen in grouped_stats[gen_col]
        if (gen_data := data[data[gen_col] == gen].sort_values(by=fitness_col))
    ]

    # Fit a 2nd-order polynomial trendline
    x, y = data[gen_col], data[fitness_col]
    coefficients = np.polyfit(x, y, deg=2)
    trendline = np.poly1d(coefficients)

    # Create scatter plot
    plt.figure(figsize=(12, 6))
    plt.scatter(data[gen_col], data[fitness_col], alpha=0.6, s=10, label='Fitness Data')

    # Highlight min and 50th points
    for min_point, mid_point in highlighted_points:
        plt.scatter(min_point[gen_col], min_point[fitness_col], color='green', edgecolors='black', s=50, label='Min Fitness')
        plt.scatter(mid_point[gen_col], mid_point[fitness_col], color='blue', edgecolors='black', s=50, label='50th Fitness')

    # Plot trendline
    x_range = np.linspace(grouped_stats[gen_col].min(), grouped_stats[gen_col].max(), 500)
    plt.plot(x_range, trendline(x_range), color='orange', label='2nd Order Trendline', linewidth=2)

    # Add error bars for mean and standard deviation
    plt.errorbar(grouped_stats[gen_col], grouped_stats['mean'], yerr=grouped_stats['std'],
                 fmt='o', color='purple', label='Mean ± Std Dev', capsize=5)

    # Customize plot
    plt.xlabel('Generation')
    plt.ylabel('Fitness Score')
    plt.title('Fitness Trend with Mean, Std Dev, and Trendline')
    plt.legend(loc='upper right', fontsize='small')
    plt.grid(alpha=0.3)
    plt.tight_layout()
    plt.show()

    return f"Trendline: f(x) = {coefficients[0]:.3f}x^2 + {coefficients[1]:.3f}x + {coefficients[2]:.3f}"

def get_detailed_simulation_data(data_path):
    """
    Load and process simulation data from the specified path.

    Parameters:
        data_path (str): Path to the simulation data file.

    Returns:
        dict: Extracted simulation data, including simulation configuration,
              network parameters, and population and cell data.
    """
    from netpyne import sim

    # Load simulation components
    sim.loadSimCfg(data_path)
    sim.loadNet(data_path)
    sim.loadNetParams(data_path)
    sim.loadSimData(data_path)

    # Extract and copy data
    extracted_data = {
        'simData': sim.allSimData.copy(),
        'cellData': sim.net.allCells.copy(),
        'popData': sim.net.allPops.copy(),
        'simCfg': sim.cfg.__dict__.copy(),
        'netParams': sim.net.params.__dict__.copy(),
    }

    # Clear all data from the simulation to free memory
    sim.clearAll()

    return extracted_data

def plot_raster_plot(ax, spiking_data_by_unit, E_gids, I_gids, mode='simulated'):
    """
    Plot a raster plot of spiking activity.

    Parameters:
        ax (matplotlib.axes.Axes): The matplotlib axis to draw the plot on.
        spiking_data_by_unit (dict): A dictionary where keys are neuron IDs (GIDs) 
                                     and values are spike time data for each unit.
        E_gids (list): List of excitatory neuron IDs.
        I_gids (list): List of inhibitory neuron IDs.

    Returns:
        matplotlib.axes.Axes: The axis with the raster plot.
    """
    # Helper function to plot spikes for a group of neurons
    def plot_group(gids, color, label, mode='simulated'):
        if mode == 'simulated':
            for gid in gids:
                if gid in spiking_data_by_unit:
                    spike_times = spiking_data_by_unit[gid]['spike_times']
                    if isinstance(spike_times, (int, float)):
                        spike_times = [spike_times]
                    ax.plot(spike_times, [gid] * len(spike_times), f'{color}.', markersize=2, label=label if gid == gids[0] else None)
        elif mode == 'experimental':
            gids = None
            # Sort spike data by firing rate
            sorted_spike_data = {k: v for k, v in sorted(spiking_data_by_unit.items(), key=lambda item: item[1]['FireRate'])}
            sorted_key = 0
            for gid, data in sorted_spike_data.items():
                spike_times = data['spike_times']
                ax.plot(spike_times, [sorted_key] * len(spike_times), f'{color}.', markersize=1, label=label if gids is None else None)
                sorted_key += 1                
                 
            # for gid, data in spiking_data_by_unit.items():
            #     spike_times = data['spike_times']
            #     ax.plot(spike_times, [gid] * len(spike_times), f'{color}.', markersize=2, label=label if gids is None else None)
            #     gids = gids or [gid]
        

    # Plot excitatory (yellow) and inhibitory (blue) neurons
    if mode == 'simulated':
        plot_group(E_gids, 'y', 'Excitatory')
        plot_group(I_gids, 'b', 'Inhibitory')
    elif mode == 'experimental':
        plot_group([], 'g', 'Experimental', mode='experimental')
    else:
        raise ValueError(f"Invalid mode: {mode}")

    # Customize axis labels and title
    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Neuron GID')
    ax.set_title('Raster Plot')
    plt.tight_layout()

    return ax

def generate_raster_plot(network_data, raster_plot_path, raster_fig_path, fig_size=(16, 4.5)):
    """
    Generate and save a raster plot of spiking activity.

    Parameters:
        network_data (dict): Contains simulation data, including spiking activity and neuron IDs.
        raster_plot_path (str): Path to save the raster plot as a PDF.
        raster_fig_path (str): Path to save the raster plot figure as a pickle file.
        fig_size (tuple): Dimensions of the plot (width, height) in inches.
    """
    # Extract relevant data
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids = simulated_data['E_Gids']
    I_gids = simulated_data['I_Gids']

    # Create the raster plot
    fig, ax = plt.subplots(figsize=fig_size)
    plot_raster_plot(ax, spiking_data_by_unit, E_gids, I_gids)

    # Save the raster plot as a PDF
    fig.savefig(raster_plot_path, format='pdf')
    print(f"Raster plot saved to: {raster_plot_path}")

    # Save the figure object as a pickle file
    with open(raster_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Raster plot figure saved to: {raster_fig_path}")

def generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path, fig_size=(16, 4.5)):
    """
    Generate and save the network bursting plot.

    Parameters:
        network_data (dict): Contains bursting summary data and plot axis from the simulation.
        bursting_plot_path (str): Path to save the bursting plot as a PDF.
        bursting_fig_path (str): Path to save the bursting plot figure as a pickle file.
        fig_size (tuple): Dimensions of the plot (width, height) in inches.
    """
    # Extract the original bursting axis
    original_ax = network_data['bursting_data']['bursting_summary_data']['ax']

    # Create a new figure with the specified size
    fig, ax = plt.subplots(figsize=fig_size)

    # Copy limits, labels, and data from the original axis
    ax.set_xlim(original_ax.get_xlim())
    ax.set_ylim(original_ax.get_ylim())
    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Firing Rate (spike count)')
    ax.set_title(original_ax.get_title())
    
    # Re-plot data from the original axis
    # for line in original_ax.get_lines():
    #     ax.plot(line.get_xdata(), line.get_ydata(), line.get_color(), marker=line.get_marker())
    ax.plot(original_ax.get_lines()[0].get_xdata(), original_ax.get_lines()[0].get_ydata(), color='royalblue')
    #ax.scatter(original_ax.get_lines()[1].get_xdata(), original_ax.get_lines()[1].get_ydata(), 'or')
    ax.scatter(original_ax.get_lines()[1].get_xdata(), original_ax.get_lines()[1].get_ydata(), marker='o', color='r')
    
    # Save the bursting plot as a PDF
    fig.savefig(bursting_plot_path, format='pdf')
    print(f"Bursting plot saved to: {bursting_plot_path}")

    # Save the figure object as a pickle file
    with open(bursting_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Bursting plot figure saved to: {bursting_fig_path}")

def generate_network_activity_summary_plot(network_data, summary_plot_path, reference=True, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Generate a network activity summary plot with raster and bursting subplots.

    Parameters:
        network_data (dict): Contains simulation and bursting data for the network.
        summary_plot_path (str): Path to save the summary plot as a PDF.
        reference (bool): Whether to include reference raster and bursting plots.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids, I_gids = simulated_data['E_Gids'], simulated_data['I_Gids']
    bursting_ax = network_data['bursting_data']['bursting_summary_data']['ax']

    # Determine the number of figure rows based on reference data
    fig_rows = 4 if reference else 2
    fig, axs = plt.subplots(fig_rows, 1, figsize=(16, 9), sharex=True)

    # Plot simulated raster
    axs[0].set_title('Raster Plot')
    plot_raster_plot(axs[0], spiking_data_by_unit, E_gids, I_gids)
    axs[0].set_ylabel('Neuron GID')

    # Plot simulated bursting
    axs[1].set_title('Bursting Plot')
    axs[1].set_xlabel('Time (s)')
    axs[1].set_ylabel('Firing Rate (Hz)')
    axs[1].plot(bursting_ax.get_lines()[0].get_xdata(), bursting_ax.get_lines()[0].get_ydata(), color='royalblue')
    #axs[1].scatter(bursting_ax.get_lines()[1].get_xdata(), bursting_ax.get_lines()[1].get_ydata(), 'or')
    axs[1].scatter(bursting_ax.get_lines()[1].get_xdata(), bursting_ax.get_lines()[1].get_ydata(), marker='o', color='r')

    # Plot reference data if provided
    if reference:
        
        #load reference data
        ref_data = np.load(reference_data, allow_pickle=True).item()
        ref_spike_data = ref_data['spiking_data']['spiking_data_by_unit']
        ref_bursting_ax = ref_data['bursting_data']['bursting_summary_data']['ax']
        
        # Sort reference spike data by FR
        sorted_ref_spike_data = {k: v for k, v in sorted(ref_spike_data.items(), key=lambda item: item[1]['FireRate'])}
        
        # trim reference data to match the length of simulated data
        max_sim_time = max(axs[1].get_lines()[0].get_xdata())
        trimmed_ref_spike_data = trim_spike_data(sorted_ref_spike_data, max_sim_time)
        
        # Plot reference raster
        axs[2] = plot_raster_plot(axs[2], trimmed_ref_spike_data, [], [], mode='experimental')
        axs[2].set_ylabel('Unit ID')
        axs[2].set_title('Reference Raster Plot')
        axs[2].set_xlabel('Time (s)')
                
        # Plot reference bursting
        axs[3].set_title('Reference Bursting Plot')
        axs[3].set_xlabel('Time (s)')
        axs[3].set_ylabel('Firing Rate (Hz)')
        
        #copy, trim, and plot reference bursting data ax
        max_sim_time = max(axs[1].get_lines()[0].get_xdata())
        axs[3] = trim_and_plot_bursting_data(axs[3], ref_bursting_ax, max_sim_time=max_sim_time)
        
        # Match y-limits across bursting plots
        match_ylims(axs[1], axs[3])
        
        #match x across all plots
        for ax in axs:
            ax.set_xlim([0, max_sim_time])

    # Finalize and save the figure
    plt.tight_layout()
    fig.savefig(summary_plot_path, format='pdf')
    print(f"Summary plot saved to: {summary_plot_path}")
    print("Summary plot generated successfully.")

def trim_and_plot_bursting_data(ax, reference_ax, max_sim_time=None):
    ### snipet from MEA Pipeline function where ax is created
    # # Plot the smoothed network activity
    # ax.plot(timeVector, firingRate, color='royalblue')
    # # Restrict the plot to the first and last 100 ms
    # ax.set_xlim([min(relativeSpikeTimes), max(relativeSpikeTimes)])
    # ax.set_ylim([min(firingRate)*0.8, max(firingRate)*1.2])  # Set y-axis limits to min and max of firingRate
    # ax.set_ylabel('Firing Rate [Hz]')
    # ax.set_xlabel('Time [ms]')
    # ax.set_title('Network Activity', fontsize=11)
    #...
    # ax.plot(burstPeakTimes, burstPeakValues, 'or')  # Plot burst peaks as red circles
    
    #copy ax features to new ax
    #ax.set_ylabel(ax_old.get_ylabel())
    #ax.set_xlabel(ax_old.get_xlabel())
    #ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    
    # Copy limits, labels, and data from the original axis
    ax.set_xlim(reference_ax.get_xlim())
    ax.set_ylim(reference_ax.get_ylim())
    #ax.set_xlabel(reference_ax.get_xlabel())
    ax.set_xlabel('Time (s)')
    ax.set_ylabel(reference_ax.get_ylabel())
    ax.set_title(reference_ax.get_title())
    
    # Copy data from the original axis
    #ax.plot(reference_ax.get_lines()[0].get_xdata(), reference_ax.get_lines()[0].get_ydata(), color='royalblue')
    #ax.plot(reference_ax.get_lines()[1].get_xdata(), reference_ax.get_lines()[1].get_ydata(), 'or')   
    #ax.plot(trimmed_x, trimmed_y, color='purple', marker='o')
    
    #trim and plotdata to match the length of the simulated data
    assert max_sim_time is not None, "max_sim_time must be provided"
    x_data, y_data = reference_ax.get_lines()[0].get_xdata(), reference_ax.get_lines()[0].get_ydata()
    x_data, y_data = x_data.copy(), y_data.copy()
    trimmed_x = x_data[x_data <= max_sim_time]
    trimmed_y = y_data[:len(trimmed_x)]
    ax.plot(trimmed_x, trimmed_y, color='purple')
    
    x_data, y_data = reference_ax.get_lines()[1].get_xdata(), reference_ax.get_lines()[1].get_ydata()
    x_data, y_data = x_data.copy(), y_data.copy()
    trimmed_x = x_data[x_data <= max_sim_time]
    trimmed_y = y_data[:len(trimmed_x)]
    ax.scatter(trimmed_x, trimmed_y, color='red', marker='o')
    
    # TODO: maybe implement this later but the above code should work for now
    # # Re-plot data from the original axis
    # for line in reference_ax.get_lines():
    #     ax.plot(line.get_xdata(), line.get_ydata(), color=line.get_color(), marker=line.get_marker())

    return ax

def trim_spike_data(spiking_data_by_unit, max_sim_time):
    """Trim spike data to fit within the simulation time."""
    trimmed_data = spiking_data_by_unit.copy()
    for gid, data in trimmed_data.items():
        data['spike_times'] = [t for t in data['spike_times'] if t <= max_sim_time]
        trimmed_data[gid] = data   
    return trimmed_data   
    # return {gid: {'spike_times': [t for t in data['spike_times'] if t <= max_sim_time]}
    #         for gid, data in spiking_data_by_unit.items()}

def match_ylims(ax1, ax2):
    """Match the y-axis limits between two axes."""
    min_ylim = min(ax1.get_ylim()[0], ax2.get_ylim()[0])
    max_ylim = max(ax1.get_ylim()[1], ax2.get_ylim()[1])
    padding = 0.1 * (max_ylim - min_ylim)
    ax1.set_ylim(min_ylim - padding, max_ylim + padding)
    ax2.set_ylim(min_ylim - padding, max_ylim + padding)

def collect_fitness_data(simulation_run_paths):
    """
    Collect fitness data from JSON files in specified simulation paths.

    Parameters:
        simulation_run_paths (list): List of paths to search for fitness JSON files.

    Returns:
        list: Sorted list of tuples containing (average_fitness, file_path).
    """
    fitness_data = []

    for path in simulation_run_paths:
        for root, _, files in os.walk(path):
            for file in files:
                if file.endswith('_fitness.json'):
                    fitness_path = os.path.join(root, file)
                    try:
                        with open(fitness_path, 'r') as f:
                            fitness_content = json.load(f)
                            average_fitness = fitness_content.get('average_fitness', float('inf'))
                            fitness_data.append((average_fitness, fitness_path))
                    except (json.JSONDecodeError, OSError) as e:
                        print(f"Error reading file {fitness_path}: {e}")

    # Sort by average fitness
    fitness_data.sort(key=lambda x: x[0])
    return fitness_data

def process_pdf_to_images(fitness_dict, fitness_path, pdf_path):
    """
    Convert PDF pages to high-resolution image objects.

    Parameters:
        pdf_path (str): Path to the PDF file.
        fitness_path (str): Path to the fitness file associated with the PDF.

    Returns:
        list: A list of tuples containing:
            - Candidate name
            - Fitness file path
            - PDF file path
            - Page number
            - Pillow image object
    """
    print(f"Processing PDF: {pdf_path}")
    
    try:
        # Extract candidate name from the file name
        file_name = os.path.basename(pdf_path)
        cand_name_match = re.search(r"(.+)_summary_plot\.pdf", file_name)
        cand_name = cand_name_match.group(1) if cand_name_match else "Unknown Candidate"
        
        #each pdf file can be treated as a single image
        # Open the PDF file
        #doc = fitz.open(pdf_path)
        #without using fitz
        image = convert_from_path(pdf_path, dpi=600, size=(8000, None), first_page=0, last_page=1)[0]
        #get output image path from fitness path
        output_image_path = os.path.join(os.path.dirname(fitness_path), f"{cand_name}_summary_plot.png")
        
        image.save(output_image_path)
        #save image to local disk
        #image.save(f"workspace/{cand_name}_summary_plot.png")
        #print('image size:', image.size) 
        print(f"Image saved to: {output_image_path}")
        #print(f"Image saved to: workspace/{cand_name}_summary_plot.png")
        # # Convert PDF pages to high-resolution images
        # pil_images = convert_from_path(pdf_path, dpi=600)  # High resolution

        # # Create a list of tuples for each page
        # images = []
        # for page_num, pil_image in enumerate(pil_images, start=1):
        #     # Store the Pillow image object directly
        #     images.append((cand_name, fitness_path, pdf_path, page_num, pil_image))
        
        average_fitness = fitness_dict.get(fitness_path, "N/A")
        output_dir = os.path.dirname(fitness_path)
        image_slide_path, pdf_slide_path = add_slide_with_image(output_image_path, cand_name, average_fitness, output_dir)
        
        images = (cand_name, fitness_path, pdf_path, image_slide_path, pdf_slide_path)
        return images

    except Exception as e:
        print(f"Error processing PDF {pdf_path}: {e}")
        return []    

def add_slide_with_image(image_path, title, average_fitness, output_dir):
    """
    Add a slide with an image and associated text to a PowerPoint presentation.

    Parameters:
        prs (Presentation): The PowerPoint presentation object.
        image_stream (BytesIO): The byte stream of the image to add.
        title (str): Title of the slide (candidate name).
        average_fitness (float): Average fitness score of the candidate.
        output_dir (str): Directory to save the slide as an image and PDF.

    Returns:
        str: Path to the saved PDF of the slide.
    """
    # # Add a slide with the selected layout
    # slide_layout = prs.slide_layouts[5]
    # slide = prs.slides.add_slide(slide_layout)

    # # Add the image to the slide
    # slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(10))

    # # Add title text
    # title_box = slide.shapes.add_textbox(Inches(0), Inches(6.25), Inches(2), Inches(0.5))
    # title_frame = title_box.text_frame
    # title_frame.text = f"Candidate: {title}"
    # title_frame.paragraphs[0].font.size = Pt(12)

    # # Add fitness score text
    # fitness_box = slide.shapes.add_textbox(Inches(0), Inches(6.5), Inches(2), Inches(0.5))
    # fitness_frame = fitness_box.text_frame
    # fitness_frame.text = f"Avg Fitness: {average_fitness:.2f}"
    # fitness_frame.paragraphs[0].font.size = Pt(12)
    
    # # Define paths for saving the slide as an image and PDF
    input_image_path = image_path

    # Define paths for saving the slide as an image and PDF
    image_path = os.path.join(output_dir, f"{title}_summary_slide.png")
    pdf_path = os.path.join(output_dir, f"{title}_summary_slide.pdf")

    # Save slide content as an image and a PDF
    generate_image_and_pdf(input_image_path, title, average_fitness, image_path, pdf_path)
    
    # #load image and add to slide
    # slide_layout = prs.slide_layouts[5]
    # slide = prs.slides.add_slide(slide_layout)
    # slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10))

    return image_path, pdf_path

def generate_image_and_pdf(input_image_path, title, average_fitness, image_path, pdf_path):
    """
    Generate and save an image and PDF of a candidate's summary slide.

    Parameters:
        image_stream (BytesIO): Byte stream of the image to display.
        title (str): Candidate name to display on the slide.
        average_fitness (float): Average fitness score to display on the slide.
        image_path (str): Path to save the image (PNG format).
        pdf_path (str): Path to save the PDF.
    """
    #from PIL import Image

    # Create a figure and axis
    fig, ax = plt.subplots(figsize=(10, 7.5))
    ax.axis("off")
    
    #load image from file
    img = Image.open(input_image_path)
    #print('image size:', img.size)
    ax.imshow(img)
    
    # Add title and fitness score as text
    plt.text(
        0.15,  # x position for 1-inch margin (10% of the figure width)
        0.15, # y position for 1-inch margin (10% of the figure height)
        f"Candidate: {title}\nAvg Fitness: {average_fitness:.2f}",
        fontsize=10,
        ha="left",
        transform=fig.transFigure,
        #image_path = input_image_path,
    )

    # Save the figure as an image
    plt.savefig(image_path, format="png", bbox_inches="tight", pad_inches=0, dpi=600)
    print(f"Image saved to: {image_path}")

    # Save the figure as a PDF
    plt.savefig(pdf_path, format="pdf", bbox_inches="tight", pad_inches=0, dpi=600)
    print(f"PDF saved to: {pdf_path}")

    # Close the figure to free memory
    plt.close(fig)
    
    return image_path

def merge_pdfs(pdf_list, output_path):
    """
    Merge multiple PDF files into a single PDF.

    Parameters:
        pdf_list (list): List of paths to the PDF files to merge.
        output_path (str): Path to save the merged PDF.

    Returns:
        None
    """
    try:
        # Initialize the PDF merger
        merger = PdfMerger()

        # Append each PDF to the merger
        for pdf in pdf_list:
            merger.append(pdf)

        # Write the merged PDF to the output path
        merger.write(output_path)
        print(f"Merged PDF saved to: {output_path}")

    except Exception as e:
        print(f"Error merging PDFs: {e}")

    finally:
        # Ensure the merger is closed
        merger.close()
        
def submit_pdf_tasks(fitness_dict, simulation_run_paths, max_workers, debug=False):
    import time
    
    def submit_tasks_to_process_executor(debug=False):
        
        def return_tasks(tasks):
            task_states = [task._state for task in tasks]
            
            #if any state=PENDING, wait and check again until all tasks are not PENDING
            if debug:
                while 'PENDING' or 'RUNNING' in task_states:
                    task_states = [task._state for task in tasks]
                    #print(f"Task states: {task_states}")
                    num_pending = task_states.count('PENDING')
                    num_running = task_states.count('RUNNING')
                    print(f"Tasks running: {num_running}")
                    print(f"Tasks pending: {num_pending}")
                    if num_pending == 0 and num_running == 0:
                        break
                    time.sleep(5)
                    
                print(f"Task states: {task_states}")
            
            not_finished = [state for state in task_states if state != 'FINISHED']
            not_finished_count = len(not_finished)
            return tasks, tasks_submitted, not_finished, not_finished_count
        
        tasks_submitted = 0
        error_state=[]
        with ProcessPoolExecutor(max_workers=max_workers) as executor:
            for path in simulation_run_paths:
                for root, _, files in os.walk(path):
                    for file in files:
                        if file.endswith(f"_summary_plot.{file_extension}"):
                            network_summary_path = os.path.join(root, file)
                            cand_name = file.split("_summary_plot")[0]
                            fitness_path = os.path.join(root, f"{cand_name}_fitness.json")
                            pdf_path = network_summary_path

                            if file_extension == "pdf":
                                tasks.append(
                                    executor.submit(
                                        process_pdf_to_images, 
                                        #prs,
                                        fitness_dict,
                                        fitness_path,
                                        pdf_path, 
                                    )
                                )
                                tasks_submitted += 1
                            
                            #If in debug mode, only submit one task
                            if debug and tasks_submitted > 0:
                                return return_tasks(tasks)
                            
                        task_states = [task._state for task in tasks]
            
            #if any state=PENDING, wait and check again until all tasks are not PENDING
            while 'PENDING' or 'RUNNING' in task_states:
                task_states = [task._state for task in tasks]
                #print(f"Task states: {task_states}")
                num_pending = task_states.count('PENDING')
                num_running = task_states.count('RUNNING')
                print(f"Tasks running: {num_running}")
                print(f"Tasks pending: {num_pending}")
                num_finished = task_states.count('FINISHED')
                print(f"Tasks finished: {num_finished}")
                #error_bool = any([task._exception is not None for task in tasks])
                error_count = sum([task._exception is not None for task in tasks])
                print(f"Tasks with errors: {error_count}")
                if num_pending == 0 and num_running == 0:
                    break
                time.sleep(5)
        
        # Return the list of tasks submitted to the executor
        return return_tasks(tasks)
    
    """
    Submit tasks for processing PDFs in parallel.

    Parameters:
        simulation_run_paths (list): Paths to directories containing simulation results.
        fitness_dict (dict): Dictionary mapping fitness file paths to their scores.
        file_extension (str): File extension of the summary plot files (default is 'png').
        max_workers (int): Number of parallel processes for processing PDFs.

    Returns:
        list: A list of tasks submitted to the executor.
    """
    start_time = time.time()
    tasks = []
    file_extension = "pdf"
    tasks, tasks_submitted, not_finished, not_finished_count = submit_tasks_to_process_executor(debug=debug) #debug just forces one task to be subnitted for easier debugging
    print(f"Tasks submitted: {tasks_submitted}")
    print(f"Tasks not finished: {not_finished_count}")
    print(f"Time elapsed: {time.time() - start_time:.2f} seconds")
    return tasks

def process_pdf_tasks(tasks):
    """
    Process PDF tasks and add slides to the presentation.

    Parameters:
        tasks (list): List of tasks for processing PDFs.
        prs (Presentation): PowerPoint presentation object to update.
        fitness_dict (dict): Dictionary mapping fitness file paths to their scores.
        pdf_list (list): List to append generated PDF slide paths.

    Returns:
        None
    """
    pdf_list = []
    png_list = []
    
    #prs = Presentation()
    
    error_state_count = 0
    for task in tasks:
        try:
            image = task.result() 
            #for cand_name, fitness_path, pdf_path, output_image_path in images:
            cand_name, fitness_path, pdf_path, png_slide_path, pdf_slide_path = image
            assert cand_name in fitness_path, f"Fitness path mismatch: {cand_name} vs. {fitness_path}"
            # average_fitness = fitness_dict.get(fitness_path, "N/A")
            # output_dir = os.path.dirname(fitness_path)
            #pdf_slide_path = add_slide_with_image(prs, image_stream, cand_name, average_fitness, output_dir)\
            #pdf_slide_path = add_slide_with_image(prs, output_image_path, cand_name, average_fitness, output_dir)
            
            png_list.append(png_slide_path)
            pdf_list.append(pdf_slide_path)
        except Exception as e:
            error_state_count += 1
            print(f"Error processing a task: {e}")
        
    #merge pdfs
    #merge_pdfs(pdf_list, os.path.join(PROGRESS_SLIDES_PATH, "merged_slides.pdf"))
    print(f"Error state count: {error_state_count}")
    # this error is expected: Error processing a task: not enough values to unpack (expected 5, got 0) - just a bad simulation run        
    return png_list, pdf_list

def save_presentation(prs, pptx_filename, progress_slides_path=None):
    """
    Save the PowerPoint presentation to disk.

    Parameters:
        prs (Presentation): PowerPoint presentation object.
        pptx_filename (str): Name of the PowerPoint file.
        progress_slides_path (str): Path to save the PowerPoint file.

    Returns:
        None
    """
    if progress_slides_path is None:
        progress_slides_path = PROGRESS_SLIDES_PATH
    pptx_path = os.path.join(progress_slides_path, pptx_filename)
    prs.save(pptx_path)
    print(f"PowerPoint presentation saved to: {pptx_path}")

def merge_and_save_pdfs(pdf_list, pdf_filename, progress_slides_path=None):
    """
    Merge PDFs into a single file and save it.

    Parameters:
        pdf_list (list): List of PDF file paths to merge.
        pdf_filename (str): Name of the merged PDF file.

    Returns:
        None
    """
    if progress_slides_path is None:
        progress_slides_path = PROGRESS_SLIDES_PATH
        
    pdf_path = os.path.join(progress_slides_path, pdf_filename)
    merge_pdfs(pdf_list, pdf_path)
    #print(f"Merged PDF saved to: {pdf_path}")
    print(f"PDFs successfully merged and saved to: {pdf_path}")

'''main functions'''

def analyze_simulations(simulation_run_paths, reference=False, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Analyze simulations by generating network summary plots for all `_data.json` files.

    Parameters:
        simulation_run_paths (list): List of paths to search for simulation data files.
        reference (bool): Whether to include reference plots in the analysis.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    for path in simulation_run_paths:
        for root, _, files in os.walk(path):
            for file in files:
                if file.endswith('_data.json'):
                    data_path = os.path.join(root, file)
                    try:
                        generate_network_summary_slide_content(
                            data_path,
                            reference=reference,
                            reference_data=reference_data,
                            #reference_raster=reference_raster,
                            #reference_bursting=reference_bursting,
                        )
                    except Exception as e:
                        print(f"Error analyzing {file} at {data_path}: {e}")

def analyze_simulations_parallel(simulation_run_paths, reference=False, reference_data=None, reference_raster=None, reference_bursting=None, max_workers=4):
    """
    Analyze simulations in parallel by generating network summary plots for `_data.json` files.

    Parameters:
        fitness_data (list): List of paths to simulation directories.
        reference (bool): Whether to include reference plots in the analysis.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
        max_workers (int): Maximum number of worker processes for parallel execution.
    """
    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        for simulation_path in simulation_run_paths:
            for root, _, files in os.walk(simulation_path):
                for file in files:
                    if file.endswith('_data.json'):
                        data_path = os.path.join(root, file)
                        executor.submit(
                            generate_network_summary_slide_content,
                            data_path,
                            reference=reference,
                            reference_data=reference_data,
                            # reference_raster=reference_raster,
                            # reference_bursting=reference_bursting,
                        )

def generate_network_summary_slide_content(data_path, reference=False, reference_data=None, reference_raster=None, reference_bursting=None):
    """
    Analyze simulation data and generate network activity plots.

    Parameters:
        data_path (str): Path to the simulation data file.
        reference (bool): Whether to include reference plots in the summary.
        reference_raster (str): Path to reference raster plot data (if any).
        reference_bursting (str): Path to reference bursting plot data (if any).
    """
    print('') # Add a newline for clarity
    print(f"Analyzing data from: {data_path}...")

    # Define output paths
    raster_plot_path = data_path.replace('_data.json', '_raster_plot.pdf')
    bursting_plot_path = data_path.replace('_data.json', '_bursting_plot.pdf')
    raster_fig_path = data_path.replace('_data.json', '_raster_fig.pkl')
    bursting_fig_path = data_path.replace('_data.json', '_bursting_fig.pkl')
    summary_plot_path = data_path.replace('_data.json', '_summary_plot.pdf')

    # Update reference paths if provided
    if reference_raster and reference_bursting:
        reference = True
        raster_fig_path, bursting_fig_path = reference_raster, reference_bursting

    try:
        # Extract simulation data
        extracted_data = get_detailed_simulation_data(data_path)
    except Exception as e:
        print(f"Error extracting simulation data from {data_path}: {e}")
        return

    try:
        # Compute network activity metrics
        network_data = get_simulated_network_activity_metrics(**extracted_data)
    except Exception as e:
        print(f"Error computing network activity metrics for {data_path}: {e}")
        return

    try:
        # Generate raster plot
        generate_raster_plot(network_data, raster_plot_path, raster_fig_path)
    except Exception as e:
        print(f"Error generating raster plot for {data_path}: {e}")

    try:
        # Generate network bursting plot
        generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path)
    except Exception as e:
        print(f"Error generating network bursting plot for {data_path}: {e}")

    try:
        # Generate network activity summary plot
        generate_network_activity_summary_plot(
            network_data,
            summary_plot_path,
            reference=reference,
            reference_data=reference_data,
            #reference_raster=reference_raster,
            #reference_bursting=reference_bursting,
        )
    except Exception as e:
        print(f"Error generating network activity summary plot for {data_path}: {e}")

    print("Plots generated successfully.")
    
def collect_network_summary_plots_parallel(simulation_run_paths, fitness_data, max_workers=4):
    """
    Collect network summary plots from simulation paths and add them to a PowerPoint presentation.

    Parameters:
        simulation_run_paths (list): Paths to directories containing simulation results.
        fitness_data (list): List of tuples (average_fitness, file_path) from collect_fitness_data.
        file_extension (str): File extension of the summary plot files (default is 'png').
        max_workers (int): Number of parallel processes for processing PDFs.

    Returns:
        None
    """
    #prs = Presentation()
    fitness_dict = {path: fitness for fitness, path in fitness_data}
    #pdf_list = []

    # Process summary plots
    tasks = submit_pdf_tasks(
        #prs,
        fitness_dict,
        simulation_run_paths,
        max_workers,
        debug = False,
        )

    # Handle PDF tasks and update the presentation
    try:
        png_slides_paths, pdf_slides_paths = process_pdf_tasks(tasks)
    except Exception as e:
        print(f"Error processing PDF tasks: {e}")
        return

    #get slide output name based on run path name
    for path in simulation_run_paths:
        
        #divy png and pdf slides by run path
        run_path_basename = os.path.basename(path)
        curated_png_slides_paths = [slide for slide in png_slides_paths if run_path_basename in slide]
        curated_pdf_slides_paths = [slide for slide in pdf_slides_paths if run_path_basename in slide]
        
        # # build prs
        # prs = Presentation()
        # for png_slide_path in png_slides_paths:
        #     slide_layout = prs.slide_layouts[5]
        #     slide = prs.slides.add_slide(slide_layout)
        #     for curated_png_slide_path in curated_png_slides_paths:
        #         slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
        #         print(f"Adding slide: {curated_png_slide_path}")
        
        # import concurrent.futures
        # # from pptx import Presentation
        # # from pptx.util import Inches
        
        # def add_slide(prs, curated_png_slide_path):
        #     slide_layout = prs.slide_layouts[5]
        #     slide = prs.slides.add_slide(slide_layout)
        #     slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
        #     print(f"Adding slide: {curated_png_slide_path}")
        
        # # build prs
        # prs = Presentation()
        # with concurrent.futures.ThreadPoolExecutor() as executor:
        #     futures = [executor.submit(add_slide, prs, curated_png_slide_path) for curated_png_slide_path in curated_png_slides_paths]
        #     concurrent.futures.wait(futures)
            
        #     #wait for all slides to be added
        #     for future in concurrent.futures.as_completed(futures):
        #         print(future.result())
        
        from threading import Lock
        from pptx import Presentation
        from pptx.util import Inches
        import concurrent.futures

        # Define the add_slide function with a lock
        lock = Lock()

        def add_slide(prs, curated_png_slide_path):
            with lock:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(curated_png_slide_path, Inches(0), Inches(0), width=Inches(10))
                print(f"Adding slide: {curated_png_slide_path}")

        #build prs
        prs = Presentation()
        # Use ThreadPoolExecutor
        with concurrent.futures.ThreadPoolExecutor() as executor:
            futures = [executor.submit(add_slide, prs, path) for path in curated_png_slides_paths]
            concurrent.futures.wait(futures)
    
        # Save outputs
        save_presentation(prs, f"{run_path_basename}.pptx")
        merge_and_save_pdfs(curated_pdf_slides_paths, f"{run_path_basename}_slides.pdf")

#===================================================================================================
'''main script'''
# Configuration parameters
global PROGRESS_SLIDES_PATH, SIMULATION_RUN_PATHS, REFERENCE_DATA_NPY
# PROGRESS_SLIDES_PATH = (
#     # "/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/"
#     # "workspace/optimization_projects/CDKL5_DIV21/"
#     # "_3_analyze_plot_review/progress_slides"
    
#     "/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/"
#     "_3_plot_all_runs_and_select_seeds/network_summary_slides"
# )
# SIMULATION_RUN_PATHS = [
#     "/pscratch/sd/a/adammwea/workspace/yThroughput/zRBS_network_simulation_outputs/"
#     "CDKL5_DIV21/241126_Run2_improved_netparams"
# ]
# REFERENCE_DATA_NPY = (
#     "/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace"
#     "/optimization_projects/CDKL5_DIV21/_1_derive_features_from_experimental_data/"
#     "network_metrics/network_metrics_well000.npy"
# )

# Main entry point
def main(analyze_in_parallel=False, analyze_in_sequence=False, generate_presentation_in_parallel=False, generate_presentation_in_sequence=False):
    print("Starting simulation analysis...")
    print('')
    print(f"analysis in parallel: {analyze_in_parallel}")
    print(f"analysis in sequence: {analyze_in_sequence}")
    #print(f"generate presentation: {generate_presentation}")
    print(f"generate presentation in parallel: {generate_presentation_in_parallel}")
    print(f"generate presentation in sequence: {generate_presentation_in_sequence}")
    print('') 
    
    # Determine CPU usage for parallel tasks
    number_physical_cpus = int(os.cpu_count() / 2)
    
    # Collect fitness data
    fitness_data = collect_fitness_data(SIMULATION_RUN_PATHS)
    print(f"Collected fitness data for {len(fitness_data)} simulations.")

    # Analyze simulations - parallel or sequential
    #assert analyze_in_parallel or analyze_in_sequence, "At least one analysis method must be selected."
    assert not (analyze_in_parallel and analyze_in_sequence), "Select only one analysis method."
    assert not (generate_presentation_in_parallel and generate_presentation_in_sequence), "Select only one presentation method."
    
    #sequential
    if analyze_in_sequence:
        print("Analyzing simulations sequentially...")
        analyze_simulations(
            SIMULATION_RUN_PATHS,
            reference=True,
            reference_data = REFERENCE_DATA_NPY,
        )
        
    #parallel
    #analyze_in_parallel = False
    if analyze_in_parallel:
        max_workers = max(1, number_physical_cpus // 4)
        print(f"Analyzing simulations in parallel with {max_workers} workers...")
        analyze_simulations_parallel(
            SIMULATION_RUN_PATHS,
            #fitness_data,
            reference=True,
            reference_data = REFERENCE_DATA_NPY,
            max_workers=max_workers,
        )
    else:
        print("Skipping simulation analysis step.")

    # Generate presentation and summary plots
    #generate_presentation = True
    
    #sequential
    if generate_presentation_in_sequence:
        # #TODO: I need to reimplement this function to work with the new data structure
        # implemented = False
        # assert implemented, "This function is not yet implemented."
        print("Generating network summary slides and plots sequentially...")
        collect_network_summary_plots_parallel(
            SIMULATION_RUN_PATHS, 
            fitness_data, 
            #file_extension="pdf",
            max_workers=1
        )
    
    #parallel
    if generate_presentation_in_parallel:
        print("Generating network summary slides and plots...")
        collect_network_summary_plots_parallel(
            SIMULATION_RUN_PATHS, 
            fitness_data, 
            #file_extension="pdf", 
            max_workers=number_physical_cpus // 2
        )
    
    print("Simulation analysis complete.")

'''notes and main script'''
# ===============================================================================================================================
# NOTE: This script is also meant to act as a very basic analysis log. Look for corresponding notes in aw obsidian notes.
# ===============================================================================================================================

SIMULATION_RUN_PATHS = []
# ok, this is working now - 2024-12-1

#run: 241126_Run2_improved_netparams
SIMULATION_RUN_PATHS.append("/pscratch/sd/a/adammwea/workspace/yThroughput/zRBS_network_simulation_outputs/CDKL5_DIV21/241126_Run2_improved_netparams")
REFERENCE_DATA_NPY = ("/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/_1_derive_features_from_experimental_data/network_metrics/network_metrics_well000.npy")
#SLIDES_OUTPUT_PATH = "/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/_3_analyze_plot_review/network_summary_slides"
SLIDES_OUTPUT_PATH = "/pscratch/sd/a/adammwea/workspace/yThroughput/zRBS_network_simulation_outputs/CDKL5_DIV21/241126_Run2_improved_netparams/_network_summary_slides"   

# Entry point check
if __name__ == "__main__":
    main(
        #comment out the options you don't want to run, as needed
        #analyze_in_sequence=False,
        analyze_in_parallel=True,
        #analyze_in_sequence=True,
        #analyze_in_parallel=False,
        generate_presentation_in_parallel=True,
        #generate_presentation_in_sequence=True,
    )
