from setup_environment import set_pythonpath
set_pythonpath()

import os
import json
import re
import io
import fitz  # PyMuPDF
from concurrent.futures import ProcessPoolExecutor
from pptx import Presentation
from pptx.util import Inches
import analysis_functions as af

from concurrent.futures import ProcessPoolExecutor
import os
from pptx import Presentation
from pptx.util import Inches

'''Setup Python environment for running the script'''
from pprint import pprint
import setup_environment
setup_environment.set_pythonpath()

import numpy as np
import matplotlib.pyplot as plt

import matplotlib.pyplot as plt

import os
import re
import json
import pickle
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from modules.analysis.analyze_network_activity import get_simulated_network_activity_metrics

# from workspace.RBS_network_simulations.workspace.optimization_projects.CDKL5_DIV21.calculate_fitness import fitnessFunc
from workspace.RBS_network_simulations.workspace.optimization_projects.CDKL5_DIV21._2_batchrun_optimization.calculate_fitness import fitnessFunc
from netpyne import sim
import dill
import os

def plot_fitness_trend(data, gen_col='#gen', fitness_col='fitness'):
    """
    Plots a scatter plot of fitness scores by generation with a 2nd-order trendline.
    Highlights minimum and 50th fitness values for each generation.
    Displays mean and standard deviation for each generation.
    
    Parameters:
    - data (pd.DataFrame): The dataset containing generation and fitness data.
    - gen_col (str): The name of the generation column.
    - fitness_col (str): The name of the fitness column.
    """
    # Ensure numeric data types
    data[gen_col] = pd.to_numeric(data[gen_col], errors='coerce')
    data[fitness_col] = pd.to_numeric(data[fitness_col], errors='coerce')
    data = data.dropna(subset=[gen_col, fitness_col])

    # Group by generation and calculate statistics
    grouped_stats = data.groupby(gen_col)[fitness_col].agg(['mean', 'std']).reset_index()
    grouped_stats[gen_col] = grouped_stats[gen_col].astype(int)  # Ensure integer generations

    # Highlight points for min and 50th fitness in each generation
    highlighted_points = []
    generations = data[gen_col].unique()
    for gen in generations:
        gen_data = data[data[gen_col] == gen].sort_values(by=fitness_col, ascending=True)
        min_point = gen_data.iloc[0]  # Lowest fitness
        if len(gen_data) >= 50:
            mid_point = gen_data.iloc[49]  # 50th fitness
        else:
            mid_point = gen_data.iloc[-1]  # Last available if fewer than 50
        highlighted_points.append((min_point, mid_point))

    # Calculate 2nd-order trendline
    x = data[gen_col]
    y = data[fitness_col]
    coefficients_2nd_order = np.polyfit(x, y, deg=2)  # 2nd-order regression
    trendline_2nd_order = np.poly1d(coefficients_2nd_order)

    # Create the scatter plot
    plt.figure(figsize=(12, 6))
    for gen in grouped_stats[gen_col]:
        gen_data = data[data[gen_col] == gen]
        plt.scatter(gen_data[gen_col], gen_data[fitness_col], label=f'Gen {gen}', alpha=0.6, s=10)

    # Highlight points for min and 50th fitness in each generation
    for min_point, mid_point in highlighted_points:
        plt.scatter(min_point[gen_col], min_point[fitness_col], color='green', label='Min Fitness' if 'Min Fitness' not in plt.gca().get_legend_handles_labels()[1] else "", edgecolors='black', s=50)
        plt.scatter(mid_point[gen_col], mid_point[fitness_col], color='blue', label='50th Fitness' if '50th Fitness' not in plt.gca().get_legend_handles_labels()[1] else "", edgecolors='black', s=50)

    # Plot the 2nd-order trendline
    x_range = np.linspace(grouped_stats[gen_col].min(), grouped_stats[gen_col].max(), 500)
    plt.plot(x_range, trendline_2nd_order(x_range), color='orange', label='2nd Order Trendline', linewidth=2)

    # Plot mean and standard deviation for each generation
    plt.errorbar(
        grouped_stats[gen_col], grouped_stats['mean'],
        yerr=grouped_stats['std'], fmt='o', color='purple', label='Mean ± Std Dev', capsize=5
    )

    # Customize plot
    plt.xticks(grouped_stats[gen_col])  # Ensure x-axis only shows whole-number generations
    plt.xlabel('Generation')
    plt.ylabel('Fitness Score')
    plt.title('Scatter Plot with Mean, Std Dev, and 2nd Order Trendline')
    plt.legend(loc='upper right', fontsize='small', ncol=2)
    plt.grid(alpha=0.3)
    plt.tight_layout()

    # Show the plot
    plt.show()

    # Return the trendline equation
    return f"Trendline Function: f(x) = {coefficients_2nd_order[0]:.3f}x^2 + {coefficients_2nd_order[1]:.3f}x + {coefficients_2nd_order[2]:.3f}"

def get_detailed_simulation_data(data_path):
    from netpyne import sim
    """
    Processes simulation data to extract relevant information for analysis.
    
    Parameters:
    - data_path (str): The path to the simulation data file.
    
    Returns:
    - extracted_data (dict): A dictionary containing the extracted data.
    """
    # Load the simulation data
    sim.loadSimCfg(data_path)
    sim.loadNet(data_path)
    sim.loadNetParams(data_path)
    sim.loadSimData(data_path)
    
    simData = sim.allSimData.copy()
    cellData = sim.net.allCells.copy()
    popData = sim.net.allPops.copy()
    simCfg = sim.cfg.__dict__.copy()
    netParams = sim.net.params.__dict__.copy()
    
    sim.clearAll()
    
    extracted_data = {
        'simData': simData,
        'cellData': cellData,
        'popData': popData,
        'simCfg': simCfg,
        'netParams': netParams,
    }
    
    return extracted_data

def plot_raster_plot(ax, spiking_data_by_unit, E_gids, I_gids):
    """Plot a raster plot for spiking data."""
    #fig, ax = plt.subplots(figsize=(10, 6))

    # Plot excitatory neurons in yellow
    for gid in E_gids:
        if gid in spiking_data_by_unit:
            spike_times = spiking_data_by_unit[gid]['spike_times']
            spike_times = [spike_times] if isinstance(spike_times, (int, float)) else spike_times
            ax.plot(spike_times, [gid] * len(spike_times), 'y.', markersize=2)

    # Plot inhibitory neurons in blue
    for gid in I_gids:
        if gid in spiking_data_by_unit:
            spike_times = spiking_data_by_unit[gid]['spike_times']
            spike_times = [spike_times] if isinstance(spike_times, (int, float)) else spike_times
            ax.plot(spike_times, [gid] * len(spike_times), 'b.', markersize=2)

    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Neuron GID')
    ax.set_title('Raster Plot')
    plt.tight_layout()
    return ax

def generate_raster_plot(network_data, raster_plot_path, raster_fig_path):
    """Generate and save the raster plot."""
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids = simulated_data['E_Gids']
    I_gids = simulated_data['I_Gids']
    
    #init figure and axis - planning for a network summary slide 16:9 aspect ratio where this plot
    #takes up the top half of the slide
    fig, ax = plt.subplots(figsize=(16, 4.5))

    ax = plot_raster_plot(ax, spiking_data_by_unit, E_gids, I_gids)
    fig.savefig(raster_plot_path) #save as pdf
    print(f"Raster plot saved to {raster_plot_path}")

    # Save raster plot as pickle
    with open(raster_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Raster plot data saved to {raster_fig_path}")

def generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path):
    """Generate and save the network bursting plot."""
    #fig = network_data['bursting_data']['bursting_summary_data']['fig']
    ax = network_data['bursting_data']['bursting_summary_data']['ax']
    ax_old = ax
    
    # Create a new figure with shared x-axis
    # planning for a network summary slide 16:9 aspect ratio where this plot
    # takes up the bottom half of the slide
    fig, ax = plt.subplots(figsize=(16, 4.5))
    
    ### snipet from MEA Pipeline function where ax is created
    # # Plot the smoothed network activity
    # ax.plot(timeVector, firingRate, color='royalblue')
    # # Restrict the plot to the first and last 100 ms
    # ax.set_xlim([min(relativeSpikeTimes), max(relativeSpikeTimes)])
    # ax.set_ylim([min(firingRate)*0.8, max(firingRate)*1.2])  # Set y-axis limits to min and max of firingRate
    # ax.set_ylabel('Firing Rate [Hz]')
    # ax.set_xlabel('Time [ms]')
    # ax.set_title('Network Activity', fontsize=11)
    #...
    # ax.plot(burstPeakTimes, burstPeakValues, 'or')  # Plot burst peaks as red circles
    
    #copy ax features to new ax
    ax.set_xlim(ax_old.get_xlim())
    ax.set_ylim(ax_old.get_ylim())
    #ax.set_ylabel(ax_old.get_ylabel())
    ax.set_ylabel('Firing Rate (spike count)')
    #ax.set_xlabel(ax_old.get_xlabel())
    ax.set_xlabel('Time (s)')
    ax.set_title(ax_old.get_title())
    #ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    ax.plot(ax_old.get_lines()[0].get_xdata(), ax_old.get_lines()[0].get_ydata(), color='royalblue')
    ax.plot(ax_old.get_lines()[1].get_xdata(), ax_old.get_lines()[1].get_ydata(), 'or')
    

    fig.savefig(bursting_plot_path) #save as pdf
    print(f"Bursting plot saved to {bursting_plot_path}")

    # Save bursting plot as pickle
    with open(bursting_fig_path, 'wb') as f:
        pickle.dump((fig, ax), f)
    print(f"Bursting plot data saved to {bursting_fig_path}")

def generate_network_activity_summary_plot(network_data, summary_plot_path):
    """Generate a network activity summary plot with raster and bursting subplots."""
    simulated_data = network_data['simulated_data']
    spiking_data_by_unit = simulated_data['spiking_data_by_unit']
    E_gids = simulated_data['E_Gids']
    I_gids = simulated_data['I_Gids']
    ax = network_data['bursting_data']['bursting_summary_data']['ax']

    bursting_data = network_data['bursting_data']['bursting_summary_data']

    # Create a new figure with shared x-axis
    fig, axs = plt.subplots(2, 1, figsize=(16, 9), sharex=True)

    # Raster plot (top subplot)
    axs[0].set_title('Raster Plot')
    axs[0] = plot_raster_plot(axs[0], spiking_data_by_unit, E_gids, I_gids)
    axs[0].set_ylabel('Neuron GID')

    # Bursting plot (bottom subplot)
    axs[1].set_title('Bursting Plot')
    axs[1].set_xlabel('Time (s)')
    axs[1].set_ylabel('Firing Rate (spike count)')
    axs[1].plot(ax.get_lines()[0].get_xdata(), ax.get_lines()[0].get_ydata(), color='royalblue')
    axs[1].plot(ax.get_lines()[1].get_xdata(), ax.get_lines()[1].get_ydata(), 'or')   

    # Finalize layout and save
    plt.tight_layout()
    fig.savefig(summary_plot_path)
    print(f"Summary plot saved to {summary_plot_path}")

def analyze_simulation_data(data_path):
    """Main function to analyze simulation data and generate plots."""
    # Define paths for output
    raster_plot_path = re.sub(r'_data.json', '_raster_plot.pdf', data_path)
    bursting_plot_path = re.sub(r'_data.json', '_bursting_plot.pdf', data_path)
    raster_fig_path = re.sub(r'_data.json', '_raster_fig.pkl', data_path)
    bursting_fig_path = re.sub(r'_data.json', '_bursting_fig.pkl', data_path)
    summary_plot_path = re.sub(r'_data.json', '_summary_plot.pdf', data_path)

    # Load existing plots if they exist
    if os.path.exists(raster_fig_path) and os.path.exists(bursting_fig_path) and os.path.exists(summary_plot_path):
        print("Existing plots found. Loading...")
        return

    # Placeholder for actual extraction function
    extracted_data = get_detailed_simulation_data(data_path)

    # Get network activity metrics
    network_data = get_simulated_network_activity_metrics(**extracted_data)    

    # Generate plots
    generate_raster_plot(network_data, raster_plot_path, raster_fig_path)
    generate_network_bursting_plot(network_data, bursting_plot_path, bursting_fig_path)
    generate_network_activity_summary_plot(network_data, summary_plot_path)
    
    print("Plots generated successfully.")

def collect_fitness_data(simulation_run_paths):
    fitness_data = []
    for simulation_run_path in simulation_run_paths:
        for root, dirs, files in os.walk(simulation_run_path):
            for file in files:
                if file.endswith('_fitness.json'):
                    fitness_path = os.path.join(root, file)
                    with open(fitness_path, 'r') as f:
                        fitness_content = json.load(f)
                        average_fitness = fitness_content.get('average_fitness', float('inf'))
                        fitness_data.append((average_fitness, fitness_path))
    fitness_data.sort()
    return fitness_data

def analyze_simulations(fitness_data):
    for average_fitness, simulation_path in fitness_data:
        for root, dirs, files in os.walk(simulation_path):
            for file in files:
                if file.endswith('_data.json'):
                    data_path = os.path.join(root, file)
                    try:
                        af.analyze_simulation_data(data_path)
                    except Exception as e:
                        print(f'Error analyzing data for file: {file}, Error: {e}')

def analyze_simulations_parallel(fitness_data, max_workers=4):
    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        for average_fitness, simulation_path in fitness_data:
            for root, dirs, files in os.walk(simulation_path):
                for file in files:
                    if file.endswith('_data.json'):
                        data_path = os.path.join(root, file)
                        executor.submit(af.analyze_simulation_data, data_path)

def process_pdf_to_images(pdf_path, fitness_path):
    """
    Processes a PDF and returns a list of images as byte streams.
    """
    print(f'Processing {pdf_path}...')
    doc = fitz.open(pdf_path)
    file_name = os.path.basename(pdf_path)
    cand_name_match = re.search(r'(.+)_summary_plot\.pdf', file_name)
    cand_name = cand_name_match.group(1) if cand_name_match else "Unknown Candidate"

    images = []
    for page_num in range(len(doc)):
        pix = doc[page_num].get_pixmap()
        image_stream = io.BytesIO(pix.tobytes("png"))
        images.append((cand_name, fitness_path, page_num + 1, image_stream))
    doc.close()
    return images

def add_slide_with_image(prs, image_path, title, average_fitness):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10))
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(Inches(0), Inches(5), Inches(2), Inches(0.5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = f"Avg Fitness: {average_fitness}"
    print(f'Added slide for candidate: {title} with Avg Fitness: {average_fitness}')

def collect_network_summary_plots(simulation_run_paths, prs, fitness_data, file_extension='png', max_workers=4):
    """
    Collects network summary plots from simulation run paths and adds them to a PowerPoint presentation.
    """
    tasks = []
    fitness_dict = {path: fitness for fitness, path in fitness_data}

    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        for simulation_run_path in simulation_run_paths:
            for root, _, files in os.walk(simulation_run_path):
                for file in files:
                    if file.endswith(f'_summary_plot.{file_extension}'):
                        network_summary_path = os.path.join(root, file)
                        cand_name = file.split('_summary_plot')[0]
                        fitness_path = os.path.join(root, cand_name + '_fitness.json')
                        average_fitness = fitness_dict.get(fitness_path, "N/A")
                        if file_extension == 'pdf':
                            tasks.append(executor.submit(process_pdf_to_images, network_summary_path, fitness_path))
                        else:
                            add_slide_with_image(prs, network_summary_path, cand_name, average_fitness)
    
    # Collect results from tasks and add slides sequentially
    for task in tasks:
        try:
            images = task.result()
            for cand_name, fitness_path, page_num, image_stream in images:                
                #fitness_path = os.path.join(root, cand_name + '_fitness.json')
                average_fitness = fitness_dict.get(fitness_path, "N/A")
                add_slide_with_image(prs, image_stream, f"{cand_name}, Page {page_num}", average_fitness)
        except Exception as e:
            print(f"Error processing a task: {e}")

# Example usage
prs = Presentation()
simulation_run_paths = ['/path/to/simulation/run']
fitness_data = [(0.5, '/path/to/fitness/file')]
collect_network_summary_plots(simulation_run_paths, prs, fitness_data, file_extension='png', max_workers=4)

if __name__ == "__main__":
    progress_slides_path = '/pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/_2_batchrun_optimization/progress_slides'
    simulation_run_paths = [
        '/pscratch/sd/a/adammwea/workspace/yThroughput/zRBS_network_simulation_outputs/CDKL5_DIV21/241126_Run2_improved_netparams',
    ]
    number_physical_cpus = int(os.cpu_count()/2)
    #number_available_threads = os.thread_count()
    #get_experimental_data_network_summary_plot()

    fitness_data = collect_fitness_data(simulation_run_paths)
    #analyze_simulations(fitness_data)
    analyze_simulations_parallel(fitness_data, max_workers=number_physical_cpus)

    prs = Presentation()
    #collect_network_summary_plots(simulation_run_paths, prs, file_extension='pdf')

    collect_network_summary_plots(simulation_run_paths, prs, fitness_data, file_extension='pdf', max_workers=number_physical_cpus)
    
    # Uncomment to save the presentation
    prs.save(os.path.join(progress_slides_path, 'CDKL5_DIV21_campaign_1.pptx'))
    print('Done analyzing simulations')

# Notes and analysis log:

# 2024-11-29 

# How to run on interactive node in Perlmutter command line:
# salloc --nodes=1 --ntasks-per-node=256 -C cpu -q interactive -t 04:00:00
# Copy paste the rest:
# '''
# module load conda
# conda activate netsims_env
# srun --tasks-per-node=256 --cpu-bind=cores python /pscratch/sd/a/adammwea/workspace/RBS_network_simulations/workspace/optimization_projects/CDKL5_DIV21/generate_update_slides_v3.py
# '''

# Still developing at this point. Soon I will track analysis notes here - at least as initial scratch before putting in obsidian.



